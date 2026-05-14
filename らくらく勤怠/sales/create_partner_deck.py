# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

GREEN  = RGBColor(0x06,0xC7,0x55)
DGREEN = RGBColor(0x04,0x9A,0x40)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x1A,0x1A,0x2E)
GRAY   = RGBColor(0x55,0x55,0x55)
LGRAY  = RGBColor(0xF4,0xF4,0xF4)
RED    = RGBColor(0xDC,0x26,0x26)
ORANGE = RGBColor(0xF5,0x9E,0x0B)
BLUE   = RGBColor(0x1D,0x4E,0xD8)
LBLUE  = RGBColor(0xEF,0xF6,0xFF)
LGREEN = RGBColor(0xDC,0xFB,0xE5)
GOLD   = RGBColor(0xD9,0x7A,0x06)
NAVY   = RGBColor(0x0F,0x17,0x2A)
PURPLE = RGBColor(0x6D,0x28,0xD9)
blank  = prs.slide_layouts[6]
TOTAL  = 10

def rect(s,l,t,w,h,fill=None,line=None,lw=None):
    sh=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb=fill
    sh.line.fill.background()
    if line:
        sh.line.color.rgb=line
        if lw: sh.line.width=lw
    else: sh.line.fill.background()
    return sh

def txt(s,text,l,t,w,h,size=14,bold=False,color=DARK,align=PP_ALIGN.LEFT,wrap=True,italic=False):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.word_wrap=wrap
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color
    return tb

def hdr(s,title,sub=None):
    rect(s,0,0,13.33,1.25,fill=NAVY)
    rect(s,0,1.25,0.06,6.25,fill=GREEN)
    txt(s,title,0.45,0.12,11,0.65,size=26,bold=True,color=WHITE)
    if sub: txt(s,sub,0.45,0.78,11,0.4,size=13,color=RGBColor(0xAA,0xCC,0xFF))

def pnum(s,n):
    txt(s,f"{n} / {TOTAL}",11.8,7.1,1.4,0.35,size=11,color=RGBColor(0x99,0x99,0x99),align=PP_ALIGN.RIGHT)

# ===== S1 表紙 =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=NAVY)
rect(s,0,0,13.33,0.06,fill=GREEN)
rect(s,0,7.44,13.33,0.06,fill=GREEN)
for i in range(6): rect(s,10.5+i*0.45,0,0.35,7.5,fill=RGBColor(0x1F,0x2A,0x3E))
txt(s,"PARTNERSHIP PROGRAM",1.0,0.9,11,0.45,size=13,bold=True,color=GREEN,italic=True)
txt(s,"社労士・顧問パートナー",1.0,1.45,11,1.0,size=42,bold=True,color=WHITE)
txt(s,"紹介プログラムのご案内",1.0,2.45,11,0.9,size=36,bold=True,color=WHITE)
txt(s,"ラクラク勤怠 x 社労士事務所",1.0,3.7,11,0.5,size=18,color=GREEN)
txt(s,"派遣会社の「勤怠DX」と「離職防止」を一緒に解決しませんか",1.0,4.3,11,0.45,size=14,color=RGBColor(0xCC,0xDD,0xFF))
rect(s,1.0,5.3,4.5,0.06,fill=GREEN)
txt(s,"2026年 ラクラク勤怠 株式会社（準備中）",1.0,5.45,8,0.4,size=12,color=RGBColor(0x88,0xAA,0xCC))
pnum(s,1)

# ===== S2 ご挨拶・サマリー =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"はじめに","本提案のサマリー")
pnum(s,2)
txt(s,"いつもお客様のご支援、誠にお疲れさまでございます。",0.4,1.45,12.5,0.5,size=15,color=DARK)
txt(s,"このたび、派遣業界に特化した勤怠管理サービス「ラクラク勤怠」の紹介パートナーを募集しております。",0.4,2.0,12.5,0.5,size=13,color=GRAY)

points=[
    (GREEN, "1社あたり月3,000〜6,000円", "の継続紹介手数料を 12ヶ月間 お支払い"),
    (BLUE,  "初回成約ボーナス ¥10,000", "顧客が30日無料トライアルを完了後に支給"),
    (ORANGE,"紹介資料・営業同行サポート", "貴所の負担なく顧問先にご紹介可能"),
    (PURPLE,"顧問先の労務管理コストを削減", "結果として顧問契約の維持・拡大にも貢献"),
]
for i,(c,ttl,desc) in enumerate(points):
    y=2.85+i*0.95
    rect(s,0.4,y,12.5,0.85,fill=LGRAY)
    rect(s,0.4,y,0.1,0.85,fill=c)
    txt(s,ttl,0.7,y+0.1,7.5,0.45,size=15,bold=True,color=c)
    txt(s,desc,0.7,y+0.5,11.5,0.35,size=12,color=GRAY)

# ===== S3 派遣会社の課題 =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"01  顧問先派遣会社の課題","先生方が普段から相談を受けている問題と一致しているはずです")
pnum(s,3)

problems=[
    (RED,   "離職率の高さ","派遣社員の年間離職率は正社員の2〜3倍\n採用・教育コストが経営を圧迫"),
    (ORANGE,"紙タイムカード・手集計","電子帳簿保存法対応が急務\n月末の集計工数が膨大"),
    (BLUE,  "労務トラブル","残業時間の記録不備による未払い残業金請求\n打刻修正依頼の頻発"),
    (PURPLE,"既存ツール導入の壁","アプリDL拒否・操作教育コストでスタッフが離脱\nIT不慣れな中小派遣会社は導入できず"),
]
for i,(c,ttl,desc) in enumerate(problems):
    x=0.35+(i%2)*6.35; y=1.5+(i//2)*2.85
    rect(s,x,y,6.1,2.6,fill=LGRAY); rect(s,x,y,0.08,2.6,fill=c)
    txt(s,ttl,x+0.2,y+0.2,5.7,0.55,size=15,bold=True,color=DARK)
    txt(s,desc,x+0.2,y+0.85,5.7,1.6,size=12,color=GRAY)

# ===== S4 ラクラク勤怠とは =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"02  ラクラク勤怠 とは","LINEで完結する派遣会社特化の勤怠管理SaaS")
pnum(s,4)

txt(s,"「LINEを開いてボタンを1回押すだけ」 アプリDL不要・当日から使える勤怠管理サービス",0.4,1.35,12.5,0.5,size=14,bold=True,color=DARK)

features=[
    (GREEN, "1タップ打刻","LINEで出退勤"),
    (ORANGE,"コンディション報告","5秒の絵文字選択"),
    (RED,   "離職アラート","管理者に即時通知"),
    (BLUE,  "管理ダッシュボード","リアルタイム可視化"),
]
for i,(c,ttl,desc) in enumerate(features):
    x=0.6+i*3.05
    rect(s,x,2.0,2.85,2.5,fill=LGRAY); rect(s,x,2.0,2.85,0.08,fill=c)
    txt(s,ttl,x,2.3,2.85,0.5,size=14,bold=True,color=DARK,align=PP_ALIGN.CENTER)
    txt(s,desc,x,2.85,2.85,0.5,size=12,color=GRAY,align=PP_ALIGN.CENTER)

rect(s,0.4,4.85,12.55,2.0,fill=LGREEN)
txt(s,"料金",0.6,4.95,3,0.4,size=14,bold=True,color=DGREEN)
txt(s,"スタッフ1人あたり 150〜200円/月　　初期費用：0円　　最低契約期間：なし",0.6,5.4,12.3,0.45,size=14,color=DARK)
txt(s,"例）登録スタッフ100名の派遣会社：月額 ¥20,000",0.6,5.9,12.3,0.4,size=13,color=GRAY)
txt(s,"30日間 全機能 無料トライアル（クレジットカード不要）",0.6,6.35,12.3,0.4,size=13,bold=True,color=DGREEN)

# ===== S5 なぜ社労士パートナーか =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"03  なぜ社労士の先生にご紹介をお願いするのか","三方良しの関係を構築できるためです")
pnum(s,5)

triangle=[
    (GREEN, "派遣会社（顧問先）",
     "離職防止＋勤怠DXを同時実現\n紙タイムカード廃止でコスト削減\n打刻トラブル激減で先生への問合せ減少"),
    (BLUE,  "社労士の先生",
     "顧問先の労務リスク低減\n継続紹介手数料による安定収益\n他事務所との差別化"),
    (ORANGE,"ラクラク勤怠",
     "信頼ある先生経由でリード獲得\n中小派遣会社への確度高い接触\n業界知識を持つパートナーとの連携"),
]
for i,(c,ttl,desc) in enumerate(triangle):
    x=0.35+i*4.3
    rect(s,x,1.5,4.1,5.4,fill=LGRAY); rect(s,x,1.5,4.1,0.6,fill=c)
    txt(s,ttl,x,1.55,4.1,0.5,size=15,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s,desc,x+0.2,2.3,3.75,4.5,size=12,color=DARK)

# ===== S6 紹介手数料スキーム =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"04  紹介手数料スキーム","業界水準を上回る還元率")
pnum(s,6)

txt(s,"成果報酬型・在籍中はずっと収益",0.4,1.35,12.5,0.45,size=15,bold=True,color=DARK)

rect(s,0.35,1.95,6.0,2.4,fill=GREEN)
txt(s,"初回成約ボーナス",0.55,2.1,5.6,0.5,size=15,bold=True,color=WHITE)
txt(s,"¥10,000",0.55,2.65,5.6,0.9,size=42,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"30日無料トライアル完了時に支給",0.55,3.7,5.6,0.4,size=12,color=LGREEN,align=PP_ALIGN.CENTER)

rect(s,6.55,1.95,6.4,2.4,fill=BLUE)
txt(s,"継続紹介手数料",6.75,2.1,6.0,0.5,size=15,bold=True,color=WHITE)
txt(s,"15〜20% / 月",6.75,2.65,6.0,0.9,size=36,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"契約継続中 12ヶ月間 お支払い",6.75,3.7,6.0,0.4,size=12,color=LBLUE,align=PP_ALIGN.CENTER)

txt(s,"シミュレーション例",0.4,4.55,5,0.4,size=14,bold=True,color=DARK)
sim_data=[
    ("100名規模 1社紹介→成約", "月額売上 ¥20,000",  "紹介料 ¥3,000/月", "¥10,000 + ¥36,000 = ¥46,000"),
    ("100名規模 5社紹介→成約", "月額売上 ¥100,000", "紹介料 ¥15,000/月","¥50,000 + ¥180,000 = ¥230,000"),
    ("300名規模 3社紹介→成約", "月額売上 ¥180,000", "紹介料 ¥27,000/月","¥30,000 + ¥324,000 = ¥354,000"),
]
rect(s,0.35,5.0,12.6,0.42,fill=NAVY)
headers=["紹介ケース","顧問先の月額","月次紹介料","12ヶ月合計"]
col_w=[4.3,2.6,2.6,3.1]
col_x=[0.35,4.65,7.25,9.85]
for j,(cx,cw,h) in enumerate(zip(col_x,col_w,headers)):
    txt(s,h,cx+0.1,5.03,cw-0.2,0.36,size=11,bold=True,color=WHITE)
for i,row in enumerate(sim_data):
    y=5.42+i*0.55
    bg=WHITE if i%2==0 else LGRAY
    rect(s,0.35,y,12.6,0.5,fill=bg)
    for j,(cx,cw,cell) in enumerate(zip(col_x,col_w,row)):
        fc=DGREEN if j==3 else DARK
        bld=j==3
        txt(s,cell,cx+0.1,y+0.12,cw-0.2,0.36,size=11,bold=bld,color=fc)

# ===== S7 紹介の流れ =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"05  紹介の流れ","先生の手間は最小限に設計しています")
pnum(s,7)

steps=[
    (GREEN, "STEP 1","顧問先にご紹介",
     "ご訪問時の雑談で\n「こんなサービスありますよ」と\nお伝えいただくだけ"),
    (BLUE,  "STEP 2","当社が商談を担当",
     "ご紹介いただいた後は\n当社が直接デモ・導入支援を実施\n先生の同席は不要"),
    (ORANGE,"STEP 3","30日無料トライアル",
     "顧問先に体験していただき\n社内で評価\n離脱リスクなし"),
    (PURPLE,"STEP 4","契約・手数料お支払い",
     "正式契約後、月次で紹介料\n振込（毎月末締め・翌月15日払い）"),
]
for i,(c,step,ttl,desc) in enumerate(steps):
    x=0.35+i*3.15
    rect(s,x,1.5,2.95,4.6,fill=LGRAY); rect(s,x,1.5,2.95,0.85,fill=c)
    txt(s,step,x,1.55,2.95,0.4,size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s,ttl,x,1.97,2.95,0.4,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s,desc,x+0.15,2.55,2.65,3.5,size=12,color=DARK,align=PP_ALIGN.CENTER)
    if i<3: txt(s,"→",3.4+i*3.15,3.5,0.25,0.5,size=20,bold=True,color=GRAY,align=PP_ALIGN.CENTER)

rect(s,0.35,6.3,12.6,0.85,fill=LGREEN)
txt(s,"先生にお願いするのは「紹介の一言」だけ。商談・導入・サポートはすべて当社が実施します。",0.5,6.4,12.3,0.4,size=13,bold=True,color=DGREEN)
txt(s,"営業同行ご希望の場合は、無料で当社担当が同席いたします。",0.5,6.8,12.3,0.32,size=11,color=DGREEN)

# ===== S8 マーケティング・営業ツール提供 =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"06  パートナー特典・サポート","紹介活動を支援するツールをご提供")
pnum(s,8)

tools=[
    (GREEN, "📄","紹介用パンフレット","A4両面の印刷可能PDF\n貴所オリジナルロゴ入りも可",),
    (BLUE,  "🎬","商品紹介動画","3分の説明動画\nメールやLINEで送信可能"),
    (ORANGE,"👥","共同セミナー開催","派遣会社経営者向けセミナーを\n貴所と共催（参加費は折半）"),
    (PURPLE,"📞","商談同行・電話サポート","当社の営業担当が無料で同行\n顧問先への説明をサポート"),
    (RED,   "📊","紹介状況ダッシュボード","紹介した顧問先の利用状況・\n手数料の発生状況を確認可能"),
    (NAVY,  "🎓","30分Zoom研修","商品理解・紹介トーク習得\n初回登録時に実施"),
]
for i,(c,icon,ttl,desc) in enumerate(tools):
    col,row=i%3,i//3
    x,y=0.35+col*4.3,1.5+row*2.85
    rect(s,x,y,4.1,2.6,fill=LGRAY); rect(s,x,y,4.1,0.08,fill=c)
    txt(s,icon,x+0.15,y+0.2,0.8,0.7,size=28,color=c)
    txt(s,ttl,x+1.0,y+0.25,2.95,0.55,size=14,bold=True,color=DARK)
    txt(s,desc,x+0.2,y+0.95,3.75,1.6,size=11,color=GRAY)

# ===== S9 よくあるご質問 =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"07  よくあるご質問","")
pnum(s,9)

faqs=[
    ("Q. 紹介した顧問先が解約した場合、手数料はどうなりますか？",
     "A. 解約された月以降の手数料は発生しません。継続利用中の月のみお支払いです。"),
    ("Q. 顧問先の規模に上限はありますか？",
     "A. ありません。スタッフ数1,000名超の派遣会社でも対応可能です。規模に応じてエンタープライズプランで個別見積もりとなります。"),
    ("Q. 私の事務所のロゴを入れたパンフレットを作れますか？",
     "A. はい。月3社以上の紹介実績がある先生にはオリジナルロゴ入りパンフレットを無料で作成いたします。"),
    ("Q. 派遣業以外の顧問先にも紹介できますか？",
     "A. 派遣業以外の活用については個別ご相談ください。多店舗展開する小売・飲食店などでも一部活用可能です。"),
    ("Q. 提携契約に最低期間はありますか？",
     "A. ありません。いつでも提携を終了できますが、終了後も既存紹介分の手数料は契約期間まで継続して支払われます。"),
]
for i,(q,a) in enumerate(faqs):
    y=1.4+i*1.05
    rect(s,0.35,y,12.6,0.95,fill=LGRAY)
    txt(s,q,0.5,y+0.05,12.3,0.4,size=13,bold=True,color=DARK)
    txt(s,a,0.5,y+0.5,12.3,0.4,size=11,color=GRAY)

# ===== S10 次のステップ =====
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=NAVY)
rect(s,0,0,13.33,0.06,fill=GREEN); rect(s,0,7.44,13.33,0.06,fill=GREEN)
for i in range(6): rect(s,10.5+i*0.45,0,0.35,7.5,fill=RGBColor(0x1F,0x2A,0x3E))

txt(s,"NEXT STEP",0.6,0.5,12,0.45,size=14,bold=True,color=GREEN,italic=True)
txt(s,"ご検討いただける場合は",0.6,1.05,12,0.6,size=24,bold=True,color=WHITE)

actions=[
    (GREEN, "1","30分のオンライン面談",
     "サービスデモ＋紹介スキームの詳細をご説明します。\nZoom/Google Meet/LINE通話いずれでも対応可能"),
    (BLUE,  "2","パートナー契約書のご確認",
     "業務委託契約書のドラフトをお送りします。\nご検討後、押印いただいて締結"),
    (ORANGE,"3","研修・紹介開始",
     "30分のオンライン研修後、即座に紹介開始可能。\nツール一式をお渡しします"),
]
for i,(c,num,ttl,desc) in enumerate(actions):
    y=1.95+i*1.4
    rect(s,0.6,y,9.5,1.25,fill=RGBColor(0x1A,0x2A,0x3E))
    rect(s,0.6,y,0.08,1.25,fill=c)
    txt(s,num,0.75,y+0.1,0.8,0.6,size=28,bold=True,color=c)
    txt(s,ttl,1.55,y+0.1,8.4,0.5,size=15,bold=True,color=WHITE)
    txt(s,desc,1.55,y+0.65,8.4,0.6,size=11,color=RGBColor(0xCC,0xDD,0xFF))

rect(s,0.6,6.35,9.5,0.85,fill=GREEN)
txt(s,"お問い合わせ",0.8,6.4,7,0.35,size=12,bold=True,color=WHITE)
txt(s,"contact@rakuraku-kintai.jp（準備中）",0.8,6.75,7,0.35,size=14,bold=True,color=WHITE)
txt(s,"担当：ラクラク勤怠 株式会社 パートナー営業窓口",0.8,7.1,7,0.3,size=10,color=LGREEN)
pnum(s,10)

out = r"C:\Users\PC_User\Desktop\AI動画\rakuraku-kintai\らくらく勤怠\sales\社労士パートナー提携プログラム.pptx"
prs.save(out)
print("saved:", out)