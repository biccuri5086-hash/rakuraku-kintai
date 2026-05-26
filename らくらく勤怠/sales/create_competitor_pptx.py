# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

LINE_GREEN = RGBColor(0x06, 0xC7, 0x55)
DARK_GREEN = RGBColor(0x04, 0x9A, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x7B, 0x8D)
LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xF8)
LIGHT_GREEN = RGBColor(0xE8, 0xF8, 0xEE)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
RED = RGBColor(0xE5, 0x3E, 0x3E)
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
BLUE_BG = RGBColor(0xE7, 0xF0, 0xFA)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def box(slide, x, y, w, h, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def rbox(slide, x, y, w, h, color, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def multiline(slide, lines, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, spacing=1.3):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(size * (spacing - 1) * 0.6)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold


def header(slide, num, title, subtitle=""):
    box(slide, 0, 0, 13.33, 0.85, NAVY)
    txt(slide, f"  {num}", 0.3, 0.1, 1.0, 0.65, 14, GOLD, bold=True)
    txt(slide, title, 0.9, 0.13, 10.5, 0.6, 22, WHITE, bold=True)
    if subtitle:
        txt(slide, subtitle, 0.9, 0.5, 10.5, 0.35, 12, RGBColor(0xCC, 0xDD, 0xEE))
    box(slide, 0, 0.85, 13.33, 0.05, GOLD)


def footer(slide):
    box(slide, 0, 7.15, 13.33, 0.35, NAVY)
    txt(slide, "ラクラク勤怠  |  競合比較・差別化資料  |  小原 健太  080-9895-7770", 0.3, 7.2, 13.0, 0.3, 9, WHITE)


# =====================================================
# Slide 1 : Title
# =====================================================
s = blank_slide()
bg(s, NAVY)
box(s, 0, 0, 13.33, 0.15, GOLD)
box(s, 0, 7.35, 13.33, 0.15, GOLD)

# Left panel - white
box(s, 0.6, 1.0, 12.13, 5.5, WHITE)

txt(s, "Competitive Analysis", 1.0, 1.4, 11, 0.5, 14, GOLD, bold=True)
txt(s, "ラクラク勤怠の", 1.0, 2.0, 11, 1.0, 36, NAVY, bold=True)
txt(s, "強みと差別化ポイント", 1.0, 2.85, 11, 1.0, 42, LINE_GREEN, bold=True)

# divider
box(s, 1.0, 4.1, 1.5, 0.05, GOLD)

multiline(s, [
    "派遣会社向け勤怠管理SaaS市場における、",
    "ラクラク勤怠の独自ポジションと、競合に対する明確な優位性。",
], 1.0, 4.35, 11, 1.0, 16, GRAY, spacing=1.5)

# bottom info
box(s, 1.0, 5.6, 11.3, 0.7, LIGHT_GRAY)
txt(s, "比較対象：KING OF TIME / ジョブカン / マネーフォワード勤怠 / freee人事労務 / タッチオンタイム / 紙タイムカード",
    1.15, 5.78, 11.0, 0.4, 11, NAVY, bold=True)

txt(s, "2026年版  |  ラクラク勤怠", 0.6, 6.85, 12.13, 0.35, 11, GRAY, align=PP_ALIGN.RIGHT)


# =====================================================
# Slide 2 : 結論 (Executive Summary)
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "01", "結論：3行で言うと", "本資料の要約")

txt(s, "ラクラク勤怠の独自ポジション", 0.6, 1.2, 12.0, 0.5, 22, NAVY, bold=True)
box(s, 0.6, 1.75, 1.0, 0.04, GOLD)

# 3 big cards
y = 2.1
for i, (num, title, body) in enumerate([
    ("01", "唯一の機能",
     "「コンディション報告」を標準搭載しているのは、市場で当社だけ。\n離職予兆を毎日5秒で可視化する、業界唯一のSaaSです。"),
    ("02", "圧倒的に安い",
     "月150〜200円/人。KING OF TIMEの約半額、紙タイムカードと同等水準。\n中小派遣会社が無理なく導入できる価格設計です。"),
    ("03", "今日から使える",
     "LINEで完結、アプリDL不要、初期費用ゼロ、30分でセットアップ完了。\n機器購入も不要、当日から運用開始できます。"),
]):
    x = 0.6 + i * 4.15
    rbox(s, x, y, 3.95, 4.4, LIGHT_GREEN)
    rbox(s, x + 0.15, y + 0.15, 0.8, 0.8, LINE_GREEN)
    txt(s, num, x + 0.15, y + 0.27, 0.8, 0.6, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.2, y + 1.1, 3.6, 0.6, 20, NAVY, bold=True)
    box(s, x + 0.2, y + 1.7, 0.6, 0.04, GOLD)
    multiline(s, body.split("\n"), x + 0.2, y + 1.9, 3.6, 2.4, 13, DARK, spacing=1.5)

footer(s)


# =====================================================
# Slide 3 : 市場マップ
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "02", "派遣業界の勤怠管理 市場マップ", "ポジショニングで見る競合関係")

txt(s, "縦軸：価格  /  横軸：派遣業界への最適化度", 0.6, 1.15, 12, 0.4, 12, GRAY)

# Matrix area
mx, my, mw, mh = 1.5, 1.7, 10.5, 5.0
box(s, mx, my, mw, mh, LIGHT_GRAY)

# Axes labels
txt(s, "高い", mx - 0.85, my + 0.1, 0.8, 0.3, 11, GRAY, align=PP_ALIGN.RIGHT)
txt(s, "安い", mx - 0.85, my + mh - 0.4, 0.8, 0.3, 11, GRAY, align=PP_ALIGN.RIGHT)
txt(s, "価格", mx - 0.95, my + mh / 2 - 0.15, 0.6, 0.3, 13, NAVY, bold=True, align=PP_ALIGN.RIGHT)

txt(s, "汎用", mx + 0.1, my + mh + 0.08, 1.0, 0.3, 11, GRAY)
txt(s, "派遣特化", mx + mw - 1.1, my + mh + 0.08, 1.0, 0.3, 11, GRAY, align=PP_ALIGN.RIGHT)
txt(s, "派遣業界への最適化度 →", mx + mw/2 - 1.3, my + mh + 0.35, 2.6, 0.3, 13, NAVY, bold=True)

# Cross lines
box(s, mx + mw/2 - 0.005, my, 0.01, mh, GRAY)
box(s, mx, my + mh/2 - 0.005, mw, 0.01, GRAY)

# Plot competitors
def plot(slide, x, y, label, color, size=1.0):
    rbox(slide, x - size/2, y - size/2, size, size, color)
    txt(slide, label, x - 1.5, y - 0.15, 3.0, 0.3, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)

# KING OF TIME: 高価格・汎用 (左上)
rbox(s, mx + 1.5, my + 0.5, 1.8, 0.8, NAVY)
txt(s, "KING OF TIME", mx + 1.5, my + 0.55, 1.8, 0.3, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "330円/人", mx + 1.5, my + 0.85, 1.8, 0.3, 9, RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

# ジョブカン: 高価格・少し汎用寄り
rbox(s, mx + 3.5, my + 1.0, 1.6, 0.8, GRAY)
txt(s, "ジョブカン", mx + 3.5, my + 1.05, 1.6, 0.3, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "200-400円/人", mx + 3.5, my + 1.35, 1.6, 0.3, 9, WHITE, align=PP_ALIGN.CENTER)

# マネーフォワード勤怠
rbox(s, mx + 2.3, my + 1.8, 1.9, 0.7, GRAY)
txt(s, "マネーフォワード", mx + 2.3, my + 1.85, 1.9, 0.3, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "300円/人", mx + 2.3, my + 2.15, 1.9, 0.3, 9, WHITE, align=PP_ALIGN.CENTER)

# freee人事労務
rbox(s, mx + 5.0, my + 1.5, 1.5, 0.7, GRAY)
txt(s, "freee人事労務", mx + 5.0, my + 1.55, 1.5, 0.3, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "300円〜", mx + 5.0, my + 1.85, 1.5, 0.3, 9, WHITE, align=PP_ALIGN.CENTER)

# タッチオンタイム
rbox(s, mx + 4.0, my + 2.7, 1.7, 0.7, GRAY)
txt(s, "タッチオンタイム", mx + 4.0, my + 2.75, 1.7, 0.3, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "300円/人+機器", mx + 4.0, my + 3.05, 1.7, 0.3, 9, WHITE, align=PP_ALIGN.CENTER)

# 紙タイムカード
rbox(s, mx + 0.6, my + 3.5, 1.6, 0.7, GRAY)
txt(s, "紙タイムカード", mx + 0.6, my + 3.55, 1.6, 0.3, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "0円〜小コスト", mx + 0.6, my + 3.85, 1.6, 0.3, 9, WHITE, align=PP_ALIGN.CENTER)

# ラクラク勤怠 (右下・大きく強調)
rbox(s, mx + 7.5, my + 3.2, 2.4, 1.3, LINE_GREEN)
txt(s, "★ ラクラク勤怠", mx + 7.5, my + 3.35, 2.4, 0.4, 14, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "150〜200円/人", mx + 7.5, my + 3.78, 2.4, 0.3, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "派遣特化×低価格×LINE完結", mx + 7.5, my + 4.12, 2.4, 0.3, 9, WHITE, align=PP_ALIGN.CENTER)

# Sweet spot label
txt(s, "← 空白地帯（誰もいない）", mx + 5.5, my + 3.6, 2.0, 0.3, 10, RED, bold=True, align=PP_ALIGN.RIGHT)

footer(s)


# =====================================================
# Slide 4 : 一目で分かる比較表
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "03", "一目で分かる 機能・価格 比較表", "派遣会社が重視する10項目で評価")

# Table layout
tx, ty = 0.5, 1.2
col_widths = [2.5, 1.7, 1.6, 1.65, 1.65, 1.65, 1.55]
row_h = 0.49
headers = ["項目", "ラクラク勤怠", "KING OF TIME", "ジョブカン", "マネーフォワード", "freee人事労務", "紙タイムカード"]

# header row
x_acc = tx
for i, h_text in enumerate(headers):
    col_color = LINE_GREEN if i == 1 else NAVY
    box(s, x_acc, ty, col_widths[i], row_h, col_color)
    txt(s, h_text, x_acc, ty + 0.08, col_widths[i], row_h - 0.1, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
    x_acc += col_widths[i]

rows = [
    ("月額（1人）",            "150〜200円", "330円",   "200-400円", "300円",    "300円〜",   "0円"),
    ("初期費用",              "0円",       "0円",     "0円",       "0円",      "0円",       "0円"),
    ("LINE完結",              "○",         "×",       "×",         "×",        "×",         "×"),
    ("アプリDL不要",          "○",         "×",       "×",         "×",        "×",         "○"),
    ("コンディション報告",    "○",         "×",       "×",         "×",        "×",         "×"),
    ("離職予兆アラート",      "○",         "×",       "×",         "×",        "×",         "×"),
    ("即日導入",              "○",         "△(数日)", "△(数日)",   "△(数日)",  "△(数日)",   "○"),
    ("GPS打刻",               "○",         "○",       "○",         "○",        "○",         "×"),
    ("CSV出力",               "○",         "○",       "○",         "○",        "○",         "×"),
    ("派遣業界特化",          "○",         "×",       "×",         "×",        "×",         "×"),
]

for ri, row in enumerate(rows):
    y = ty + (ri + 1) * row_h
    x_acc = tx
    bg_color = WHITE if ri % 2 == 0 else LIGHT_GRAY
    for ci, cell in enumerate(row):
        if ci == 1:
            cell_bg = LIGHT_GREEN
        else:
            cell_bg = bg_color
        box(s, x_acc, y, col_widths[ci], row_h, cell_bg, line_color=RGBColor(0xDD, 0xDD, 0xDD))
        color = NAVY if ci == 0 else (DARK_GREEN if ci == 1 else DARK)
        bold = ci == 0 or ci == 1
        if cell == "○" and ci == 1:
            txt(s, cell, x_acc, y + 0.1, col_widths[ci], row_h - 0.1, 16, DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
        elif cell == "×":
            txt(s, cell, x_acc, y + 0.1, col_widths[ci], row_h - 0.1, 14, RED, bold=True, align=PP_ALIGN.CENTER)
        elif cell == "○":
            txt(s, cell, x_acc, y + 0.1, col_widths[ci], row_h - 0.1, 14, DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
        elif cell.startswith("△"):
            txt(s, cell, x_acc, y + 0.1, col_widths[ci], row_h - 0.1, 11, ORANGE, bold=True, align=PP_ALIGN.CENTER)
        else:
            txt(s, cell, x_acc, y + 0.1, col_widths[ci], row_h - 0.1, 11, color, bold=bold, align=PP_ALIGN.CENTER)
        x_acc += col_widths[ci]

# bottom note
txt(s, "※ 価格は2026年5月時点の公式情報より。プランにより変動。",
    0.5, 7.0, 12.0, 0.3, 9, GRAY)


# =====================================================
# Slide 5 : 差別化① コンディション報告
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "04", "差別化ポイント①  コンディション報告", "業界唯一の離職予兆検知機能")

# Left
txt(s, "なぜ重要か", 0.6, 1.15, 6.0, 0.5, 18, NAVY, bold=True)
box(s, 0.6, 1.6, 0.6, 0.04, GOLD)

multiline(s, [
    "■ 派遣スタッフの離職は「突然」起きる",
    "  体調不良・モチベ低下のサインに、",
    "  管理者が気づけないまま辞めていく。",
    "",
    "■ 1人辞めると、採用コスト10〜30万円が消える",
    "  再募集・面談・研修すべてやり直し。",
    "",
    "■ 既存ツールでは「打刻」しか取れない",
    "  KING OF TIMEもジョブカンも、",
    "  時刻データの集計しかできない。",
], 0.6, 1.85, 6.0, 5.0, 13, DARK, spacing=1.5)

# Right - feature card
rbox(s, 7.0, 1.15, 5.8, 5.7, LIGHT_GREEN)
txt(s, "ラクラク勤怠の解決策", 7.2, 1.4, 5.4, 0.5, 16, NAVY, bold=True)
box(s, 7.2, 1.85, 0.6, 0.04, LINE_GREEN)

# 5 emoji cards
emojis = [
    ("😄", "絶好調", DARK_GREEN),
    ("😊", "良い",   DARK_GREEN),
    ("😐", "普通",   GOLD),
    ("😔", "疲れ",   ORANGE),
    ("😢", "しんどい", RED),
]
for i, (em, lbl, col) in enumerate(emojis):
    y = 2.15 + i * 0.45
    box(s, 7.2, y, 0.5, 0.4, WHITE)
    txt(s, em, 7.2, y + 0.02, 0.5, 0.4, 16, DARK, align=PP_ALIGN.CENTER)
    txt(s, lbl, 7.8, y + 0.07, 1.5, 0.3, 12, col, bold=True)
    if i >= 3:
        rbox(s, 9.5, y + 0.05, 3.1, 0.3, col)
        txt(s, "→ 管理者アラート", 9.5, y + 0.08, 3.1, 0.3, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s, "毎日5秒・絵文字1タップで完結", 7.2, 4.6, 5.4, 0.4, 13, NAVY, bold=True, align=PP_ALIGN.CENTER)

box(s, 7.2, 5.15, 5.4, 0.04, GOLD)

multiline(s, [
    "■ 「疲れ」「しんどい」が出たら即フォロー",
    "■ 体調・モチベの推移をグラフで可視化",
    "■ スタッフ「気にかけてもらえる」実感UP",
], 7.2, 5.35, 5.4, 1.5, 11, DARK, spacing=1.5)

footer(s)


# =====================================================
# Slide 6 : 差別化② LINE完結
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "05", "差別化ポイント②  LINEで完結", "アプリDL不要、ID/パスワード不要")

txt(s, "他社の常識：", 0.6, 1.2, 5.5, 0.5, 18, GRAY, bold=True)
box(s, 0.6, 2.0, 5.8, 1.0, LIGHT_GRAY)
multiline(s, [
    "❌ 専用アプリのダウンロードが必要",
    "❌ ID・パスワードの登録・管理",
    "❌ 機種変更時の再ログイン",
], 0.8, 2.15, 5.5, 0.9, 12, RED)

box(s, 0.6, 3.2, 5.8, 1.0, LIGHT_GRAY)
multiline(s, [
    "❌ スタッフ「使い方分からない」連絡が頻発",
    "❌ 60代スタッフが脱落",
    "❌ 結局、紙のタイムカードが残る",
], 0.8, 3.35, 5.5, 0.9, 12, RED)

# Right - LINE side
txt(s, "ラクラク勤怠：", 7.0, 1.2, 5.5, 0.5, 18, DARK_GREEN, bold=True)
box(s, 7.0, 2.0, 5.8, 1.0, LIGHT_GREEN)
multiline(s, [
    "○ LINEを開いて「出勤」をタップするだけ",
    "○ 友だち追加だけで利用開始",
    "○ パスワード忘れる心配ゼロ",
], 7.2, 2.15, 5.5, 0.9, 12, DARK_GREEN)

box(s, 7.0, 3.2, 5.8, 1.0, LIGHT_GREEN)
multiline(s, [
    "○ 60代以上でも当日から使える",
    "○ 機種変更してもLINE引き継ぎでそのまま",
    "○ スタッフ「これなら使える」",
], 7.2, 3.35, 5.5, 0.9, 12, DARK_GREEN)

# Bottom data
rbox(s, 0.6, 4.5, 12.13, 2.3, NAVY)
txt(s, "数字で見る差", 0.85, 4.7, 12.0, 0.5, 16, GOLD, bold=True)

for i, (label, num, sub) in enumerate([
    ("LINE利用率", "95%", "日本人の95%が日常利用"),
    ("アプリDL離脱率", "70%", "1ステップで7割が脱落"),
    ("シニア層対応", "○", "60-70代でも使える"),
    ("導入時間", "5分", "QR読むだけで完了"),
]):
    x = 1.0 + i * 2.85
    rbox(s, x, 5.3, 2.7, 1.35, WHITE)
    txt(s, label, x, 5.4, 2.7, 0.3, 11, GRAY, bold=True, align=PP_ALIGN.CENTER)
    txt(s, num, x, 5.65, 2.7, 0.5, 26, LINE_GREEN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, sub, x, 6.25, 2.7, 0.3, 9, NAVY, align=PP_ALIGN.CENTER)

footer(s)


# =====================================================
# Slide 7 : 差別化③ 価格
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "06", "差別化ポイント③  圧倒的に安い", "中小派遣会社が無理なく払える価格設計")

txt(s, "1人あたり月額の比較（スタッフ100名で試算）", 0.6, 1.15, 12, 0.5, 16, NAVY, bold=True)
box(s, 0.6, 1.65, 1.0, 0.04, GOLD)

# Horizontal bar chart
companies = [
    ("ラクラク勤怠 ★",       200, 20000, LINE_GREEN, True),
    ("マネーフォワード勤怠", 300, 30000, GRAY, False),
    ("freee人事労務",        300, 30000, GRAY, False),
    ("タッチオンタイム",     300, 30000, GRAY, False),
    ("KING OF TIME",         330, 33000, NAVY, False),
    ("ジョブカン",           400, 40000, RED, False),
]

base_x = 3.5
max_w = 7.5
max_val = 400

for i, (name, price, monthly, color, hl) in enumerate(companies):
    y = 2.0 + i * 0.7
    txt(s, name, 0.6, y + 0.1, 2.7, 0.4, 13, NAVY, bold=hl, align=PP_ALIGN.RIGHT)
    bar_w = max_w * (price / max_val)
    rbox(s, base_x, y, bar_w, 0.5, color, radius=0.2)
    txt(s, f"{price}円/人", base_x + bar_w + 0.15, y + 0.1, 1.3, 0.4, 12, color, bold=True)
    txt(s, f"月額 {monthly:,}円", base_x + bar_w + 1.5, y + 0.1, 1.7, 0.4, 11, GRAY)

# Bottom savings callout
rbox(s, 0.6, 6.4, 12.13, 0.7, LIGHT_GREEN)
txt(s, "→ KING OF TIME比 月13,000円・年156,000円の削減  /  ジョブカン比 月20,000円・年24万円の削減",
    0.8, 6.55, 11.8, 0.45, 13, DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)

footer(s)


# =====================================================
# Slide 8 : 差別化④ 即日導入
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "07", "差別化ポイント④  当日から使える", "30分セットアップで即運用開始")

txt(s, "他社の導入フロー（数日〜数週間）", 0.6, 1.15, 12, 0.4, 14, GRAY, bold=True)

# Competitor flow
steps_c = [
    ("①", "問合せ"),
    ("②", "ヒアリング"),
    ("③", "見積もり"),
    ("④", "契約書"),
    ("⑤", "初期設定"),
    ("⑥", "研修"),
    ("⑦", "運用開始"),
]
for i, (n, lb) in enumerate(steps_c):
    x = 0.6 + i * 1.75
    rbox(s, x, 1.65, 1.55, 0.7, LIGHT_GRAY)
    txt(s, n, x, 1.7, 1.55, 0.3, 12, GRAY, bold=True, align=PP_ALIGN.CENTER)
    txt(s, lb, x, 1.95, 1.55, 0.3, 11, GRAY, align=PP_ALIGN.CENTER)
    if i < len(steps_c) - 1:
        txt(s, "▸", x + 1.45, 1.75, 0.4, 0.4, 14, GRAY, align=PP_ALIGN.CENTER)

txt(s, "営業日 5〜30日", 10.0, 2.6, 3.0, 0.4, 12, RED, bold=True)

# Rakuraku flow
txt(s, "ラクラク勤怠の導入フロー（合計30分）", 0.6, 3.4, 12, 0.4, 14, DARK_GREEN, bold=True)

steps_r = [
    ("①", "管理者登録", "5分"),
    ("②", "QR共有", "5分"),
    ("③", "スタッフ登録", "5分/人"),
    ("④", "運用開始", "即時"),
]
for i, (n, lb, tm) in enumerate(steps_r):
    x = 0.6 + i * 3.1
    rbox(s, x, 3.95, 2.85, 1.0, LIGHT_GREEN)
    rbox(s, x + 0.15, 4.05, 0.7, 0.7, LINE_GREEN)
    txt(s, n, x + 0.15, 4.17, 0.7, 0.5, 14, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, lb, x + 1.0, 4.05, 1.7, 0.3, 13, NAVY, bold=True)
    txt(s, tm, x + 1.0, 4.4, 1.7, 0.3, 12, DARK_GREEN, bold=True)
    if i < len(steps_r) - 1:
        txt(s, "▸", x + 2.7, 4.2, 0.4, 0.4, 18, LINE_GREEN, align=PP_ALIGN.CENTER)

# Bottom comparison
y2 = 5.3
rbox(s, 0.6, y2, 12.13, 1.5, NAVY)
txt(s, "理由：余計な手続きをすべて省いた", 0.85, y2 + 0.15, 12, 0.4, 14, GOLD, bold=True)

multiline(s, [
    "■ 機器購入なし（既存スマホ・PCで完結）",
    "■ 専用アカウント不要（LINEを使うだけ）",
    "■ 訪問・研修なし（マニュアル1枚で済む）",
    "■ 紙の契約書なし（オンライン完結）",
], 0.85, y2 + 0.55, 11.5, 1.0, 11, WHITE, spacing=1.4)

footer(s)


# =====================================================
# Slide 9 : 競合別 vs KING OF TIME
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "08", "詳細比較①  vs KING OF TIME", "業界最大手との位置づけ")

# Left - KOT
rbox(s, 0.6, 1.15, 6.0, 5.85, LIGHT_GRAY)
txt(s, "KING OF TIME", 0.85, 1.4, 5.5, 0.5, 22, NAVY, bold=True)
txt(s, "ヒューマンテクノロジーズ（業界シェアNo.1）", 0.85, 1.85, 5.5, 0.3, 11, GRAY)
box(s, 0.85, 2.15, 0.6, 0.04, GOLD)

txt(s, "強み", 0.85, 2.3, 5.5, 0.3, 13, NAVY, bold=True)
multiline(s, [
    "○ 機能の豊富さ（シフト・有給・残業など）",
    "○ 大企業の導入実績",
    "○ 給与ソフト連携が豊富",
], 0.95, 2.6, 5.3, 1.2, 11, DARK, spacing=1.4)

txt(s, "弱み（派遣会社目線）", 0.85, 3.95, 5.5, 0.3, 13, RED, bold=True)
multiline(s, [
    "× 月330円/人と高め",
    "× 機能が多すぎて中小には過剰",
    "× スタッフが専用アプリor打刻機を使う",
    "× 派遣特化の機能はなし（離職予兆等）",
    "× 設定が複雑、研修に数時間必要",
], 0.95, 4.25, 5.3, 2.5, 11, DARK, spacing=1.4)

# Right - Rakuraku
rbox(s, 6.8, 1.15, 6.0, 5.85, LIGHT_GREEN)
txt(s, "ラクラク勤怠 ★", 7.05, 1.4, 5.5, 0.5, 22, DARK_GREEN, bold=True)
txt(s, "派遣業界に特化したLINE型勤怠SaaS", 7.05, 1.85, 5.5, 0.3, 11, DARK_GREEN)
box(s, 7.05, 2.15, 0.6, 0.04, GOLD)

txt(s, "強み", 7.05, 2.3, 5.5, 0.3, 13, DARK_GREEN, bold=True)
multiline(s, [
    "○ 月150〜200円（KOT比 半額以下）",
    "○ シンプル機能で迷わない",
    "○ LINEだけで完結、機器・アプリ不要",
    "○ コンディション報告で離職予兆を検知",
    "○ 30分で導入完了、当日から使える",
], 7.15, 2.6, 5.3, 2.6, 11, DARK, spacing=1.4)

txt(s, "競合への明確な答え", 7.05, 5.4, 5.5, 0.3, 13, NAVY, bold=True)
multiline(s, [
    "「機能が足りない」 → 中小派遣には十分",
    "「実績が浅い」     → 月単位解約OK・無料試用",
], 7.15, 5.7, 5.3, 1.0, 11, DARK, spacing=1.4)

footer(s)


# =====================================================
# Slide 10 : 詳細比較 vs ジョブカン
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "09", "詳細比較②  vs ジョブカン", "オールインワン勤怠との違い")

# Left
rbox(s, 0.6, 1.15, 6.0, 5.85, LIGHT_GRAY)
txt(s, "ジョブカン勤怠管理", 0.85, 1.4, 5.5, 0.5, 22, NAVY, bold=True)
txt(s, "DONUTS（シリーズ累計15万社）", 0.85, 1.85, 5.5, 0.3, 11, GRAY)
box(s, 0.85, 2.15, 0.6, 0.04, GOLD)

txt(s, "強み", 0.85, 2.3, 5.5, 0.3, 13, NAVY, bold=True)
multiline(s, [
    "○ 給与・採用・経費とシリーズ連携",
    "○ 機能が網羅的（残業・有給・シフトすべて）",
    "○ TVCMで知名度が高い",
], 0.95, 2.6, 5.3, 1.2, 11, DARK, spacing=1.4)

txt(s, "弱み（派遣会社目線）", 0.85, 3.95, 5.5, 0.3, 13, RED, bold=True)
multiline(s, [
    "× 月200〜400円と幅広く、機能追加で高くなる",
    "× 「機能が多すぎて何を選べばいいか分からない」",
    "× 派遣業特有のニーズに合わない",
    "× 設定画面が複雑、PCに張り付く時間必要",
    "× スタッフ向け管理コスト（パスワード等）",
], 0.95, 4.25, 5.3, 2.5, 11, DARK, spacing=1.4)

# Right
rbox(s, 6.8, 1.15, 6.0, 5.85, LIGHT_GREEN)
txt(s, "ラクラク勤怠 ★", 7.05, 1.4, 5.5, 0.5, 22, DARK_GREEN, bold=True)
txt(s, "派遣会社の現場で必要なものだけに絞る", 7.05, 1.85, 5.5, 0.3, 11, DARK_GREEN)
box(s, 7.05, 2.15, 0.6, 0.04, GOLD)

txt(s, "強み", 7.05, 2.3, 5.5, 0.3, 13, DARK_GREEN, bold=True)
multiline(s, [
    "○ 月150〜200円（ジョブカン上位機能比 半額）",
    "○ 機能を3つに絞る：打刻 / コンディション / 出力",
    "○ 派遣現場の「離職問題」を直接解決",
    "○ 操作が単純、迷いようがない",
    "○ スタッフ管理ゼロ（パスワード不要）",
], 7.15, 2.6, 5.3, 2.6, 11, DARK, spacing=1.4)

txt(s, "メッセージング", 7.05, 5.4, 5.5, 0.3, 13, NAVY, bold=True)
multiline(s, [
    "ジョブカン = 「全部入りで給与まで」",
    "ラクラク勤怠 = 「派遣の打刻と離職対策に集中」",
], 7.15, 5.7, 5.3, 1.0, 11, DARK, spacing=1.4)

footer(s)


# =====================================================
# Slide 11 : 競合に勝てない領域（正直ベース）
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "10", "正直に言う：競合が勝つ領域", "ラクラク勤怠が向かないケース")

txt(s, "全部やろうとして失敗するのが一番危険。", 0.6, 1.15, 12, 0.4, 14, NAVY, bold=True)
txt(s, "「ここはやらない」を明確にすることで、本当に強い領域を作っています。", 0.6, 1.55, 12, 0.4, 12, GRAY)

# 3 cards - こんな会社は他社推奨
title_y = 2.2
txt(s, "次のような会社にはラクラク勤怠は向きません：", 0.6, title_y, 12, 0.4, 14, RED, bold=True)

cases = [
    ("社員数500名以上の大企業",
     "経理・人事の複雑な業務がある大企業は KING OF TIME や ジョブカンの方が機能要件を満たします。\nラクラク勤怠は中小派遣会社（社員5〜100名規模）に特化しています。"),
    ("給与計算まで一気通貫で必要",
     "ラクラク勤怠は勤怠データの記録・出力までです。\n給与計算が必要な場合は CSV を出力して既存ソフトに連携してください。\n（freee・マネーフォワード等とのCSV互換あり）"),
    ("複雑なシフト・残業ルールがある",
     "夜勤手当・深夜残業の自動計算・複数勤務形態の組み合わせは未対応。\n単純な打刻管理＋スタッフのメンタルケアに振り切っています。"),
]

for i, (title, body) in enumerate(cases):
    y = 2.7 + i * 1.45
    rbox(s, 0.6, y, 12.13, 1.3, LIGHT_GRAY)
    rbox(s, 0.8, y + 0.2, 0.9, 0.9, RED)
    txt(s, str(i+1), 0.8, y + 0.32, 0.9, 0.7, 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, 1.95, y + 0.18, 10.5, 0.4, 14, NAVY, bold=True)
    multiline(s, body.split("\n"), 1.95, y + 0.55, 10.5, 0.8, 11, DARK, spacing=1.4)

# Bottom honesty message
rbox(s, 0.6, 7.0, 12.13, 0.4, NAVY)
txt(s, "「派遣スタッフの離職を減らす」一点に集中。それ以外は素直に他社を勧めます。", 0.6, 7.05, 12, 0.3, 11, GOLD, bold=True, align=PP_ALIGN.CENTER)


# =====================================================
# Slide 12 : スイートスポット
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "11", "ベストフィット顧客像", "ラクラク勤怠が圧倒的に効くのはこんな会社")

txt(s, "次の条件に1つでも当てはまれば、必ず効果があります。", 0.6, 1.15, 12, 0.4, 14, NAVY, bold=True)

# 6 cards (3x2)
fits = [
    ("👥", "中小派遣会社", "社員5〜30名／登録スタッフ50〜300名"),
    ("📋", "紙タイムカード現役", "集計に毎月数時間〜数日かかっている"),
    ("📈", "離職率が高い", "登録スタッフの30%以上が1年以内離職"),
    ("🏭", "現場系を扱う", "製造・物流・食品・軽作業など現場業種"),
    ("📱", "スタッフがシニア多め", "LINE以外のアプリが苦手な層"),
    ("💰", "コストを抑えたい", "1人月300円超は払えない"),
]

for i, (em, title, body) in enumerate(fits):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 4.15
    y = 1.85 + row * 2.6
    rbox(s, x, y, 3.95, 2.4, LIGHT_GREEN)
    rbox(s, x + 0.3, y + 0.25, 1.0, 1.0, WHITE)
    txt(s, em, x + 0.3, y + 0.42, 1.0, 0.8, 32, NAVY, align=PP_ALIGN.CENTER)
    txt(s, title, x + 1.45, y + 0.35, 2.4, 0.4, 16, NAVY, bold=True)
    box(s, x + 1.45, y + 0.78, 0.5, 0.04, LINE_GREEN)
    multiline(s, [body], x + 1.45, y + 0.95, 2.4, 1.2, 11, DARK, spacing=1.4)

footer(s)


# =====================================================
# Slide 13 : 想定切り返しトーク
# =====================================================
s = blank_slide()
bg(s, WHITE)
header(s, "12", "競合への切り返しトーク集", "営業現場でよく出る質問への答え")

talks = [
    ("「KING OF TIMEで十分」と言われたら",
     "→ コンディション報告は KING OF TIME にはありません。\n  打刻だけなら確かに KING OF TIME で OK ですが、\n  スタッフの離職予兆を毎日キャッチできるのは弊社だけです。"),
    ("「ジョブカンで給与まで一気通貫したい」",
     "→ 給与連携が最重要なら、ジョブカンが向いています。\n  ラクラク勤怠は CSV 出力で給与ソフトに連携できますが、\n  全部一気通貫が必須ならジョブカンを推奨します。"),
    ("「紙タイムカードで困っていない」",
     "→ 集計時間に何時間使っていますか？\n  月10時間なら年120時間。\n  スタッフ100名なら、月¥20,000で集計時間ほぼゼロになります。"),
    ("「実績が浅くて不安」",
     "→ 30日間無料お試しで判断してください。\n  契約も月単位、違約金なし。\n  合わなければ翌月解約できます。"),
    ("「他社と比べて検討したい」",
     "→ ぜひ比較してください。\n  特に「コンディション報告」と「LINE完結」の2点を、\n  他社にも聞いてみてください。"),
]

y = 1.15
for i, (q, a) in enumerate(talks):
    rbox(s, 0.6, y, 12.13, 1.1, LIGHT_GRAY)
    rbox(s, 0.7, y + 0.1, 0.8, 0.9, ORANGE)
    txt(s, "Q", 0.7, y + 0.22, 0.8, 0.6, 20, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, q, 1.7, y + 0.12, 10.8, 0.4, 14, NAVY, bold=True)
    multiline(s, a.split("\n"), 1.7, y + 0.5, 10.8, 0.6, 10, DARK, spacing=1.3)
    y += 1.18

footer(s)


# =====================================================
# Slide 14 : Closing
# =====================================================
s = blank_slide()
bg(s, NAVY)
box(s, 0, 0, 13.33, 0.15, GOLD)
box(s, 0, 7.35, 13.33, 0.15, GOLD)

txt(s, "ラクラク勤怠は、", 0.6, 1.3, 12, 0.7, 28, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "「派遣スタッフの離職」だけに集中したSaaS。", 0.6, 1.9, 12, 0.8, 26, GOLD, bold=True, align=PP_ALIGN.CENTER)

box(s, 6.0, 2.85, 1.3, 0.04, GOLD)

multiline(s, [
    "競合より機能は少ない。でも、",
    "派遣会社の現場で本当に必要なものに振り切っています。",
], 0.6, 3.1, 12, 1.0, 16, WHITE, align=PP_ALIGN.CENTER, spacing=1.5)

# 3 strengths
y = 4.2
for i, (label, val) in enumerate([
    ("唯一の機能", "コンディション報告"),
    ("圧倒的価格", "150〜200円/人"),
    ("即日導入", "30分でセットアップ"),
]):
    x = 0.6 + i * 4.15
    rbox(s, x, y, 3.95, 1.5, WHITE)
    txt(s, label, x, y + 0.2, 3.95, 0.4, 12, LINE_GREEN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, val, x, y + 0.65, 3.95, 0.7, 18, NAVY, bold=True, align=PP_ALIGN.CENTER)

# Contact
rbox(s, 3.5, 6.05, 6.33, 1.0, GOLD)
txt(s, "30日間無料トライアル受付中", 3.5, 6.15, 6.33, 0.4, 14, NAVY, bold=True, align=PP_ALIGN.CENTER)
txt(s, "小原 健太  /  080-9895-7770  /  biccuri5086@gmail.com", 3.5, 6.55, 6.33, 0.4, 12, NAVY, align=PP_ALIGN.CENTER)


# Save
prs.save("c:/Users/PC_User/Desktop/AI動画/rakuraku-kintai/らくらく勤怠/sales/15_競合比較・差別化資料.pptx")
print("OK: created 15_競合比較・差別化資料.pptx")
