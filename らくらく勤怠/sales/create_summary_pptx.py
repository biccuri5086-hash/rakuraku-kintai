from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

LINE_GREEN  = RGBColor(0x06, 0xC7, 0x55)
DARK_GREEN  = RGBColor(0x04, 0x9A, 0x40)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1A, 0x1A, 0x2E)
GRAY        = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
RED         = RGBColor(0xDC, 0x26, 0x26)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)
BLUE        = RGBColor(0x1D, 0x4E, 0xD8)
LIGHT_GREEN = RGBColor(0xDC, 0xFB, 0xE5)

blank_layout = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, fill=None, line=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, fill=LINE_GREEN)
    add_text(slide, title, 0.4, 0.15, 10, 0.6, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.75, 10, 0.45, size=14, color=RGBColor(0xCC,0xFF,0xDD))

def slide_num(slide, n, total):
    add_text(slide, f"{n} / {total}", 12.3, 7.1, 0.9, 0.35, size=11, color=GRAY, align=PP_ALIGN.RIGHT)

TOTAL = 10

# ========== SLIDE 1: 表紙 ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=LINE_GREEN)
add_rect(s, 0, 4.8, 13.33, 2.7, fill=DARK_GREEN)

add_text(s, "ラクラク勤怠", 1.5, 1.0, 10, 1.4, size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "機能・セキュリティ 総合説明資料", 1.5, 2.5, 10, 0.7, size=24, color=RGBColor(0xCC,0xFF,0xDD), align=PP_ALIGN.CENTER)
add_text(s, "〜 非エンジニアでもわかる、仕組みと安全性の全体像 〜", 1.5, 3.2, 10, 0.5, size=15, color=RGBColor(0xAA,0xFF,0xCC), align=PP_ALIGN.CENTER)
add_text(s, "2026年5月　ラクラク勤怠 開発チーム", 1.5, 5.3, 10, 0.5, size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "CONFIDENTIAL", 1.5, 6.1, 10, 0.45, size=12, color=RGBColor(0xAA,0xFF,0xCC), align=PP_ALIGN.CENTER)
slide_num(s, 1, TOTAL)

# ========== SLIDE 2: 目次 ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
add_header_bar(s, "目次", "この資料でわかること")
slide_num(s, 2, TOTAL)

items = [
    ("1", "サービス全体像", "ラクラク勤怠とは何か"),
    ("2", "スタッフ向け機能", "打刻・コンディション報告の仕組み"),
    ("3", "管理者向け機能", "ダッシュボードとアラートの仕組み"),
    ("4", "データの流れ", "情報がどこをどう通るか"),
    ("5", "セキュリティ設計", "4層の防御でデータを守る"),
    ("6", "過去の脆弱性と修正", "発見された問題と対処"),
    ("7", "法的対応", "個人情報保護・利用規約の整備"),
    ("8", "今後の課題", "残っているセキュリティ強化項目"),
]

for i, (num, title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.5 + row * 1.35
    add_rect(s, x, y, 6.1, 1.1, fill=WHITE, line=LINE_GREEN, line_width=15000)
    add_rect(s, x, y, 0.55, 1.1, fill=LINE_GREEN)
    add_text(s, num, x, y+0.2, 0.55, 0.7, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x+0.65, y+0.08, 5.2, 0.45, size=15, bold=True, color=DARK)
    add_text(s, desc, x+0.65, y+0.55, 5.2, 0.45, size=11, color=GRAY)

# ========== SLIDE 3: サービス全体像 ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
add_header_bar(s, "1. サービス全体像", "ラクラク勤怠とは何か")
slide_num(s, 3, TOTAL)

add_text(s, "「アプリのダウンロード不要・LINEだけで使える」 派遣スタッフ向け勤怠管理サービス", 0.4, 1.4, 12.5, 0.5, size=14, color=DARK)

boxes = [
    (LINE_GREEN, WHITE, "👆", "1タップ打刻", "出勤・退勤ボタンを\n押すだけで時刻記録"),
    (ORANGE,     WHITE, "😊", "コンディション報告", "5段階の絵文字で\n今日の調子を5秒で報告"),
    (BLUE,       WHITE, "📊", "管理者ダッシュボード", "全スタッフの状況を\nリアルタイムで一覧表示"),
    (RED,        WHITE, "⚠️", "離職防止アラート", "調子が悪いスタッフを\n自動でハイライト表示"),
]

for i, (bg, fg, icon, title, desc) in enumerate(boxes):
    x = 0.35 + i * 3.15
    add_rect(s, x, 2.1, 2.95, 3.5, fill=bg)
    add_text(s, icon, x, 2.2, 2.95, 0.7, size=30, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x, 3.0, 2.95, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x, 3.55, 2.95, 0.9, size=11, color=RGBColor(0xEE,0xFF,0xEE), align=PP_ALIGN.CENTER)

add_rect(s, 0.35, 5.8, 12.6, 1.4, fill=LIGHT_GRAY)
add_text(s, "技術スタック（開発者向け参考情報）", 0.55, 5.85, 12, 0.35, size=11, bold=True, color=GRAY)
add_text(s,
    "フロントエンド: Next.js 16 (App Router) + TypeScript + Tailwind CSS  ／  "
    "データベース: Supabase (PostgreSQL)  ／  "
    "認証基盤: LINE LIFF v2  ／  デプロイ: Vercel (CDN + Serverless Functions)",
    0.55, 6.2, 12.3, 0.85, size=11, color=GRAY)

# ========== SLIDE 4: スタッフ向け機能 ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
add_header_bar(s, "2. スタッフ向け機能", "打刻・コンディション報告の仕組み")
slide_num(s, 4, TOTAL)

steps = [
    ("STEP 1", "URL / QRコードを開く", "管理者からLINEで送られたURLを開くだけ\nアプリのダウンロードは一切不要"),
    ("STEP 2", "電話番号を登録（初回のみ）", "携帯電話番号を入力して登録\n次回以降この操作は不要"),
    ("STEP 3", "出勤ボタンをタップ", "大きなボタンを1回タップするだけ\n打刻時刻が自動で記録される"),
    ("STEP 4", "コンディションを報告", "😄良い ／ 😊まあまあ ／ 😐普通 ／ 😔疲れ ／ 😢しんどい\nの5段階から選ぶ（5秒で完了・任意）"),
    ("STEP 5", "退勤ボタンをタップ", "帰るときに退勤ボタンを押す\n勤務時間が自動計算される"),
]

for i, (step, title, desc) in enumerate(steps):
    x = 0.35 + (i % 3) * 4.2
    y = 1.5 + (i // 3) * 2.8
    add_rect(s, x, y, 3.9, 2.45, fill=LIGHT_GRAY)
    add_rect(s, x, y, 3.9, 0.45, fill=LINE_GREEN)
    add_text(s, step, x, y+0.05, 3.9, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x+0.15, y+0.55, 3.6, 0.45, size=13, bold=True, color=DARK)
    add_text(s, desc, x+0.15, y+1.05, 3.6, 1.25, size=11, color=GRAY)

add_rect(s, 8.75, 1.5, 4.2, 2.45, fill=RGBColor(0xE8,0xF5,0xE9))
add_rect(s, 8.75, 1.5, 4.2, 0.45, fill=DARK_GREEN)
add_text(s, "GPS機能（補足）", 8.75, 1.55, 4.2, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s,
    "出勤タップ時、バックグラウンドで\n位置情報を自動取得・記録\n\n"
    "✔ 許可しなくても打刻は正常動作\n✔ 不正打刻（場所確認）に活用\n✔ データはSupabaseに安全保存",
    8.9, 2.0, 3.9, 1.8, size=11, color=DARK)

# ========== SLIDE 5: 管理者向け機能 ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
add_header_bar(s, "3. 管理者向け機能", "ダッシュボードとアラートの仕組み")
slide_num(s, 5, TOTAL)

add_rect(s, 0.35, 1.45, 5.8, 5.65, fill=LIGHT_GRAY)
add_rect(s, 0.35, 1.45, 5.8, 0.42, fill=BLUE)
add_text(s, "ダッシュボード機能", 0.35, 1.47, 5.8, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

dash_items = [
    ("🔢", "登録スタッフ数", "打刻したスタッフの総数"),
    ("✅", "出勤済み人数", "当日出勤済みのスタッフ数"),
    ("⚠️", "要フォロー人数", "コンディション「疲れ・しんどい」の数"),
    ("🕐", "出退勤時刻一覧", "各スタッフの出勤・退勤・勤務時間"),
    ("📅", "日付切り替え", "過去の日付のデータを遡って確認"),
    ("😊", "コンディション表示", "絵文字＋ラベルで状態を一目で把握"),
]
for i, (icon, title, desc) in enumerate(dash_items):
    y = 2.0 + i * 0.82
    add_text(s, icon, 0.45, y, 0.5, 0.6, size=16, color=DARK)
    add_text(s, title, 0.95, y+0.05, 2.2, 0.35, size=12, bold=True, color=DARK)
    add_text(s, desc, 0.95, y+0.38, 4.9, 0.35, size=10, color=GRAY)

add_rect(s, 6.5, 1.45, 6.5, 5.65, fill=RGBColor(0xFF,0xF3,0xF3))
add_rect(s, 6.5, 1.45, 6.5, 0.42, fill=RED)
add_text(s, "⚠️ 要フォローアラートの仕組み", 6.5, 1.47, 6.5, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "コンディションスコアが「2以下」のスタッフを自動検出", 6.65, 2.0, 6.2, 0.5, size=12, bold=True, color=DARK)

score_data = [
    ("😄 絶好調", "5", LINE_GREEN),
    ("😊 良い感じ", "4", RGBColor(0x22, 0xC5, 0x5E)),
    ("😐 普通", "3", ORANGE),
    ("😔 少し疲れ", "2", RED),
    ("😢 しんどい", "1", RED),
]
for i, (label, score, color) in enumerate(score_data):
    y = 2.65 + i * 0.62
    highlight = i >= 3
    if highlight:
        add_rect(s, 6.55, y-0.04, 6.35, 0.56, fill=RGBColor(0xFF,0xE5,0xE5))
    add_text(s, label, 6.65, y, 3.0, 0.5, size=12, color=DARK, bold=highlight)
    add_text(s, f"スコア: {score}", 9.8, y, 1.2, 0.5, size=12, color=color, bold=True)
    if highlight:
        add_text(s, "← 行が赤くなり管理者に通知", 11.1, y, 1.8, 0.5, size=10, color=RED)

add_rect(s, 6.55, 5.8, 6.35, 1.1, fill=RGBColor(0xFF,0xED,0xED))
add_text(s, "管理者へのアクセス制限\nパスワード認証（サーバーサイド検証）済み", 6.7, 5.85, 6.1, 0.95, size=11, color=DARK)

# ========== SLIDE 6: データの流れ ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
add_header_bar(s, "4. データの流れ", "情報がどこをどう通るか（非エンジニア向け）")
slide_num(s, 6, TOTAL)

nodes = [
    (1.0,  3.3, "📱\nスタッフの\nスマートフォン", LINE_GREEN),
    (4.3,  3.3, "🌐\nVercel\n（アプリサーバー）", BLUE),
    (7.6,  1.5, "🔑\nLINE\n（本人確認）", RGBColor(0x00,0xB9,0x00)),
    (7.6,  5.1, "🗄️\nSupabase\n（データベース）", RGBColor(0x1F,0x83,0x5A)),
    (10.9, 3.3, "💻\n管理者の\nパソコン", RGBColor(0x44,0x44,0xFF)),
]
for (x, y, label, color) in nodes:
    add_rect(s, x, y, 2.0, 1.8, fill=color)
    add_text(s, label, x, y+0.1, 2.0, 1.6, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

arrows = [
    (3.05, 4.15, "打刻・コンディション\nデータを送信"),
    (6.35, 2.35, "LINE IDで\n本人確認"),
    (6.35, 6.0,  "出退勤データ\nを保存・取得"),
    (9.65, 4.15, "データを\n表示"),
]
for (x, y, label) in arrows:
    add_text(s, "→", x, y-0.15, 1.2, 0.4, size=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, label, x, y+0.2, 1.25, 0.5, size=9, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(s, 0.35, 6.0, 12.6, 1.25, fill=RGBColor(0xE8,0xF4,0xFD))
add_text(s, "💡 ポイント：通信はすべてHTTPS（暗号化）で保護されています。パスワードや個人情報が途中で盗み見られることはありません。\nデータはすべてSupabaseのサーバー（AWS 東京リージョン）に保存。スタッフの情報はスマートフォン本体には一切残りません。",
    0.5, 6.05, 12.3, 1.1, size=11, color=BLUE)

# ========== SLIDE 7: セキュリティ設計 ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
add_header_bar(s, "5. セキュリティ設計", "4層の防御でデータを守る")
slide_num(s, 7, TOTAL)

layers = [
    ("第1層", "通信の暗号化（HTTPS）",
     "スマートフォン ↔ サーバー間の通信をすべて暗号化。\n"
     "打刻データ・パスワード・位置情報などが第三者に見える心配はない。",
     BLUE, "🔒"),
    ("第2層", "LINE本人確認",
     "打刻にはLINEアカウントが必須。LINE IDはLINE社が厳重に管理しており、\n"
     "他人がなりすまして打刻することは極めて困難。",
     LINE_GREEN, "👤"),
    ("第3層", "管理者パスワード（サーバー検証）",
     "管理画面のパスワードはサーバー内部でのみ照合。\n"
     "ブラウザ側（スマホ・PC）にパスワードの情報は一切残らない設計。",
     ORANGE, "🔑"),
    ("第4層", "データアクセス制限（RLS）",
     "Supabaseのデータベース側でルール設定済み。\n"
     "自分のデータは自分だけが見られ、他スタッフのデータは参照不可。",
     RED, "🛡️"),
]

for i, (layer, title, desc, color, icon) in enumerate(layers):
    y = 1.5 + i * 1.45
    add_rect(s, 0.35, y, 12.6, 1.25, fill=LIGHT_GRAY)
    add_rect(s, 0.35, y, 1.5, 1.25, fill=color)
    add_text(s, icon,  0.35, y+0.05, 1.5, 0.55, size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, layer, 0.35, y+0.65, 1.5, 0.5,  size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, 2.0, y+0.1,  10.7, 0.45, size=14, bold=True, color=DARK)
    add_text(s, desc,  2.0, y+0.62, 10.7, 0.55, size=11, color=GRAY)

# ========== SLIDE 8: 過去の脆弱性と修正 ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
add_header_bar(s, "6. 発見された脆弱性と修正", "開発中に発見・対処したセキュリティ問題")
slide_num(s, 8, TOTAL)

add_text(s, "開発中に発見された問題を迅速に修正しました。透明性の観点からすべて開示します。", 0.4, 1.4, 12.5, 0.45, size=13, color=DARK)

issues = [
    (
        "【重大】管理者パスワードの漏洩リスク",
        "問題の内容",
        "「NEXT_PUBLIC_ADMIN_PASSWORD」という名前で環境変数を設定していた。\n"
        "Next.jsでは「NEXT_PUBLIC_」のついた変数はブラウザ側に自動公開される仕様のため、\n"
        "誰でも管理者パスワードを見られる状態だった。",
        "対処内容",
        "API Route（サーバーサイド機能）を新設。パスワードの照合はサーバー内部のみで行うよう変更。\n"
        "ブラウザにはパスワード情報を一切送らない設計に修正済み。",
        RED, "✅ 修正完了"
    ),
    (
        "【中】データベースのアクセス制限（RLS）",
        "問題の内容",
        "Supabaseのデータベースは初期状態では全データに誰でもアクセス可能。\n"
        "このままでは全スタッフの打刻データが外部から参照・改ざんできる状態だった。",
        "対処内容",
        "Row Level Security（行単位のアクセス制限）を設定。\n"
        "自分のLINE IDに紐づくデータのみ読み書き可能なルールを適用済み。",
        ORANGE, "✅ 修正完了"
    ),
]

for i, (title, label1, desc1, label2, desc2, color, status) in enumerate(issues):
    y = 2.0 + i * 2.45
    add_rect(s, 0.35, y, 12.6, 2.2, fill=LIGHT_GRAY)
    add_rect(s, 0.35, y, 12.6, 0.42, fill=color)
    add_text(s, title,  0.5,  y+0.05, 10.5, 0.35, size=13, bold=True, color=WHITE)
    add_text(s, status, 10.9, y+0.05, 1.9,  0.35, size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(s, f"■ {label1}",  0.5, y+0.52, 2.0, 0.35, size=11, bold=True, color=color)
    add_text(s, desc1,           0.5, y+0.85, 12.2, 0.6,  size=10, color=GRAY)
    add_text(s, f"■ {label2}",  0.5, y+1.4,  2.0, 0.35, size=11, bold=True, color=LINE_GREEN)
    add_text(s, desc2,           0.5, y+1.72, 12.2, 0.55, size=10, color=GRAY)

# ========== SLIDE 9: 法的対応 ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
add_header_bar(s, "7. 法的対応・整備状況", "サービス運営に必要な法的ページの整備")
slide_num(s, 9, TOTAL)

legal_items = [
    ("プライバシーポリシー", "/privacy",
     "✅ 整備済み",
     "・収集する個人情報の種類（LINE情報・電話番号・GPS・打刻データ）を明記\n"
     "・データの保管先（AWS東京リージョン）を明記\n"
     "・GPS取得の同意UIをアプリ内に設置済み\n"
     "・利用目的・第三者提供の条件・問い合わせ先を記載",
     LINE_GREEN),
    ("利用規約", "/terms",
     "✅ 整備済み",
     "・サービスの提供条件・禁止事項を8条で規定\n"
     "・免責事項・サービス変更・解約条件を明記\n"
     "・準拠法（日本法）・管轄裁判所を指定",
     LINE_GREEN),
    ("特定商取引法に基づく表記", "/legal",
     "⏳ 会社情報の記入待ち",
     "・販売業者名・所在地・連絡先のプレースホルダーを設置済み\n"
     "・オーナー様による会社情報の記入が必要\n"
     "・記入後、すぐに公開可能な状態",
     ORANGE),
]

for i, (title, url, status, desc, color) in enumerate(legal_items):
    y = 1.5 + i * 1.8
    add_rect(s, 0.35, y, 12.6, 1.6, fill=LIGHT_GRAY)
    add_rect(s, 0.35, y, 0.25, 1.6, fill=color)
    add_text(s, title,  0.72, y+0.1,  5.0, 0.45, size=14, bold=True, color=DARK)
    add_text(s, url,    0.72, y+0.52, 3.0, 0.35, size=11, color=GRAY)
    add_text(s, status, 9.8,  y+0.1,  2.9, 0.45, size=13, bold=True, color=color, align=PP_ALIGN.RIGHT)
    add_text(s, desc,   0.72, y+0.88, 12.0, 0.65, size=10, color=GRAY)

add_rect(s, 0.35, 7.0, 12.6, 0.35, fill=RGBColor(0xE8,0xF5,0xE9))
add_text(s, "💡 LPページ（/lp）・スタッフ向けマニュアル・管理者向けマニュアルも整備済み", 0.5, 7.02, 12.3, 0.3, size=11, color=DARK_GREEN)

# ========== SLIDE 10: 今後の課題 ==========
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
add_header_bar(s, "8. 今後の課題・強化項目", "残っているセキュリティ・機能強化")
slide_num(s, 10, TOTAL)

todo_items = [
    ("P1 最優先", [
        ("特定商取引法の会社情報記入", "オーナー様", "法的義務のため公開前に必須"),
        ("Vercel環境変数 ADMIN_PASSWORD の設定", "オーナー様", "管理画面ログインに必要"),
        ("Supabaseでadminロール設定", "オーナー様", "管理者権限の正式付与"),
    ], RED),
    ("P2 近日中", [
        ("Supabase Auth（正式認証）への移行", "開発チーム", "管理者権限をDB側で厳密に制御"),
        ("管理者LINE通知機能", "開発チーム", "要フォロースタッフ発生時に即時通知"),
        ("打刻修正機能（管理画面）", "開発チーム", "現在はSupabase直接修正が必要"),
    ], ORANGE),
    ("P3 将来的に", [
        ("独自ドメイン取得", "オーナー様", "ブランディング強化"),
        ("Stripe決済連携", "開発チーム", "有料プラン課金の自動化"),
        ("CSVエクスポート機能", "開発チーム", "給与計算への連携"),
    ], BLUE),
]

for i, (priority, items, color) in enumerate(todo_items):
    x = 0.35 + i * 4.3
    add_rect(s, x, 1.45, 4.05, 5.7, fill=LIGHT_GRAY)
    add_rect(s, x, 1.45, 4.05, 0.42, fill=color)
    add_text(s, priority, x, 1.47, 4.05, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, (task, owner, note) in enumerate(items):
        y = 2.0 + j * 1.65
        add_rect(s, x+0.15, y, 3.75, 1.45, fill=WHITE)
        add_text(s, task,  x+0.25, y+0.08, 3.55, 0.45, size=11, bold=True, color=DARK)
        add_text(s, f"担当: {owner}", x+0.25, y+0.52, 3.55, 0.3, size=10, color=color, bold=True)
        add_text(s, note,  x+0.25, y+0.82, 3.55, 0.55, size=10, color=GRAY)

add_rect(s, 0.35, 7.05, 12.6, 0.35, fill=RGBColor(0xE8,0xF5,0xE9))
add_text(s, "P1の3項目が完了次第、サービスとして正式公開可能な状態になります", 0.5, 7.07, 12.3, 0.3, size=12, bold=True, color=DARK_GREEN)

out = r"C:\Users\PC_User\Desktop\AI動画\rakuraku-kintai\らくらく勤怠\sales\機能_セキュリティ説明資料.pptx"
prs.save(out)
print(f"saved: {out}")
