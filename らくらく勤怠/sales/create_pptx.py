from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

LINE_GREEN = RGBColor(0x06, 0xC7, 0x55)
DARK_GREEN = RGBColor(0x04, 0x9A, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x7B, 0x8D)
LIGHT_GREEN = RGBColor(0xE8, 0xF8, 0xEE)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
RED = RGBColor(0xE5, 0x3E, 0x3E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txb

def multiline(slide, lines, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, spacing=1.2):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(size * (spacing - 1) * 0.5)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold

# =============================================
# スライド1: タイトル
# =============================================
s1 = blank_slide(prs)
bg(s1, LINE_GREEN)

# 左側ホワイトパネル
box(s1, 0, 0, 7.5, 7.5, WHITE)

# 緑アクセントバー
box(s1, 7.5, 0, 0.15, 7.5, DARK_GREEN)

# タイトル
txt(s1, "ラクラク勤怠", 0.6, 1.5, 6.5, 1.5, 54, LINE_GREEN, bold=True)
txt(s1, "派遣スタッフの離職を、LINEで防ぐ。", 0.6, 3.0, 6.5, 0.8, 22, DARK, bold=False)

# サブテキスト
multiline(s1, [
    "✅  LINEで1タップ打刻（アプリDL不要）",
    "✅  毎日5秒のコンディション報告",
    "✅  管理者への離職サインアラート",
], 0.6, 4.0, 6.5, 2.0, 16, GRAY)

# 右側
txt(s1, "営業ご提案資料", 8.0, 1.2, 4.8, 0.6, 13, WHITE, bold=False, align=PP_ALIGN.CENTER)
txt(s1, "2026", 8.0, 6.8, 4.8, 0.5, 12, RGBColor(0xCC, 0xFF, 0xDD), align=PP_ALIGN.CENTER)

# =============================================
# スライド2: 課題提示
# =============================================
s2 = blank_slide(prs)
bg(s2, RGBColor(0xF8, 0xF9, 0xFA))

box(s2, 0, 0, 13.33, 1.2, LINE_GREEN)
txt(s2, "こんなお悩みはありませんか？", 0.5, 0.2, 12, 0.8, 28, WHITE, bold=True)

problems = [
    ("😰", "スタッフが突然来なくなった", "兆候があっても気づけない構造になっている"),
    ("📋", "打刻修正の依頼が毎週来て手間", "紙・LINE・電話での連絡が混在している"),
    ("📱", "専用アプリを入れてくれない", "私物スマホに知らないアプリは入れたくない"),
]

for i, (emoji, title, desc) in enumerate(problems):
    bx = 0.4 + i * 4.3
    box(s2, bx, 1.5, 4.0, 4.5, WHITE)
    box(s2, bx, 1.5, 4.0, 0.08, LINE_GREEN)
    txt(s2, emoji, bx + 1.5, 1.8, 1.5, 1.0, 40, DARK, align=PP_ALIGN.CENTER)
    txt(s2, title, bx + 0.2, 2.9, 3.6, 0.8, 16, DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(s2, desc, bx + 0.2, 3.8, 3.6, 1.5, 12, GRAY, align=PP_ALIGN.CENTER)

txt(s2, "これらすべてを、ラクラク勤怠が解決します", 0.5, 6.3, 12.3, 0.7,
    18, LINE_GREEN, bold=True, align=PP_ALIGN.CENTER)

# =============================================
# スライド3: ソリューション
# =============================================
s3 = blank_slide(prs)
bg(s3, WHITE)

box(s3, 0, 0, 13.33, 1.2, DARK)
txt(s3, "ラクラク勤怠 3つの機能", 0.5, 0.2, 12, 0.8, 28, WHITE, bold=True)

features = [
    ("① LINEで1タップ打刻", LINE_GREEN,
     "LINEを開いてボタンを押すだけ\nアプリDL不要・当日から使える\nGPSを自動で裏側取得"),
    ("② コンディション報告", RGBColor(0xFF, 0x8C, 0x00),
     "5段階絵文字で「今日の調子」を報告\n😄絶好調 😊良い 😐普通\n😔疲れ  😢しんどい"),
    ("③ 管理者アラート", RED,
     "調子が悪いスタッフを即座に可視化\n離職サインを早期にキャッチ\n声かけのタイミングを逃さない"),
]

for i, (title, color, desc) in enumerate(features):
    bx = 0.4 + i * 4.3
    box(s3, bx, 1.5, 4.0, 5.3, LIGHT_GREEN)
    box(s3, bx, 1.5, 4.0, 0.6, color)
    txt(s3, title, bx + 0.15, 1.55, 3.7, 0.5, 15, WHITE, bold=True)
    for j, line in enumerate(desc.split("\n")):
        txt(s3, line, bx + 0.25, 2.4 + j * 0.7, 3.5, 0.6, 14, DARK)

# =============================================
# スライド4: 競合比較
# =============================================
s4 = blank_slide(prs)
bg(s4, RGBColor(0xF8, 0xF9, 0xFA))

box(s4, 0, 0, 13.33, 1.2, LINE_GREEN)
txt(s4, "他社ツールとの比較", 0.5, 0.2, 12, 0.8, 28, WHITE, bold=True)

headers = ["比較項目", "ラクラク勤怠", "KING OF TIME", "ジョブカン", "Touch On Time"]
col_w = [3.2, 2.2, 2.0, 2.0, 2.0]
col_x = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w + 0.05)

rows = [
    ["LINEで完結", "✅", "❌", "❌", "❌"],
    ["アプリDL不要", "✅", "❌", "❌", "❌"],
    ["コンディション報告", "✅", "❌", "❌", "❌"],
    ["離職サインアラート", "✅", "❌", "❌", "❌"],
    ["月額料金（目安）", "150円/人〜", "330円/人", "200〜400円/人", "220円/人〜"],
]

# ヘッダー行
for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    hcolor = LINE_GREEN if i == 1 else DARK
    box(s4, x, 1.3, w, 0.55, hcolor)
    txt(s4, h, x + 0.1, 1.35, w - 0.2, 0.45, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    row_y = 1.9 + ri * 0.72
    bg_color = LIGHT_GREEN if ri % 2 == 0 else WHITE
    for ci, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        if ci == 1:
            box(s4, x, row_y, w, 0.68, RGBColor(0xD4, 0xF5, 0xE2))
        else:
            box(s4, x, row_y, w, 0.68, bg_color)
        color = LINE_GREEN if cell == "✅" else (RED if cell == "❌" else DARK)
        bold = ci == 1
        txt(s4, cell, x + 0.1, row_y + 0.1, w - 0.2, 0.5, 13, color,
            bold=bold, align=PP_ALIGN.CENTER)

# =============================================
# スライド5: 料金プラン
# =============================================
s5 = blank_slide(prs)
bg(s5, WHITE)

box(s5, 0, 0, 13.33, 1.2, DARK)
txt(s5, "料金プラン", 0.5, 0.2, 12, 0.8, 28, WHITE, bold=True)

plans = [
    ("スタータープラン", "150円 × スタッフ数 / 月",
     ["1タップ打刻（出勤・退勤）", "コンディション報告", "管理者ダッシュボード", "GPSによる位置記録"],
     LINE_GREEN, False),
    ("スタンダードプラン", "200円 × スタッフ数 / 月",
     ["スターターの全機能", "コンディションアラート通知", "月次レポートCSV出力", "チャットサポート"],
     DARK_GREEN, True),
    ("エンタープライズ", "要お見積り（100名以上）",
     ["全機能", "AI離職リスクスコア（開発予定）", "既存システム連携", "専任サポート担当"],
     DARK, False),
]

for i, (name, price, items, color, rec) in enumerate(plans):
    bx = 0.4 + i * 4.3
    box(s5, bx, 1.4, 4.0, 5.5, RGBColor(0xF5, 0xF5, 0xF5))
    box(s5, bx, 1.4, 4.0, 1.0, color)
    if rec:
        txt(s5, "★ 推奨", bx + 2.5, 1.42, 1.3, 0.3, 10, WHITE, bold=True)
    txt(s5, name, bx + 0.15, 1.5, 3.7, 0.5, 15, WHITE, bold=True)
    txt(s5, price, bx + 0.15, 2.05, 3.7, 0.4, 12, WHITE)
    for j, item in enumerate(items):
        txt(s5, f"• {item}", bx + 0.2, 2.6 + j * 0.56, 3.6, 0.5, 12, DARK)

box(s5, 0.4, 7.0, 12.5, 0.35, LIGHT_GREEN)
txt(s5, "🎁  30日間・全機能・無料トライアル実施中（クレジットカード不要・違約金なし）",
    0.6, 7.0, 12.1, 0.35, 13, DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)

# =============================================
# スライド6: 使い方（スタッフ側）
# =============================================
s6 = blank_slide(prs)
bg(s6, RGBColor(0xF0, 0xFD, 0xF4))

box(s6, 0, 0, 13.33, 1.2, LINE_GREEN)
txt(s6, "使い方：スタッフ側（たった3ステップ）", 0.5, 0.2, 12, 0.8, 28, WHITE, bold=True)

steps = [
    ("STEP 1", "LINEでURLを開く", "QRコードまたはURLをタップするだけ\nアプリダウンロード不要"),
    ("STEP 2", "「出勤する」を1タップ", "大きなボタンを押すだけで打刻完了\nGPSは裏側で自動取得"),
    ("STEP 3", "コンディションを選ぶ", "5秒で今日の調子を絵文字で報告\n無理に書く必要はありません"),
]

for i, (step, title, desc) in enumerate(steps):
    bx = 0.5 + i * 4.2
    box(s6, bx, 1.5, 3.8, 4.8, WHITE)
    box(s6, bx, 1.5, 3.8, 0.08, LINE_GREEN)
    txt(s6, step, bx + 0.15, 1.65, 3.5, 0.45, 13, LINE_GREEN, bold=True)
    txt(s6, title, bx + 0.15, 2.15, 3.5, 0.6, 18, DARK, bold=True)
    for j, line in enumerate(desc.split("\n")):
        txt(s6, line, bx + 0.15, 3.0 + j * 0.55, 3.5, 0.5, 13, GRAY)

    # 矢印
    if i < 2:
        txt(s6, "→", bx + 3.85, 3.2, 0.5, 0.5, 24, LINE_GREEN, bold=True)

txt(s6, "操作は以上！　1日の所要時間：合計10秒以下",
    0.5, 6.5, 12.3, 0.6, 18, LINE_GREEN, bold=True, align=PP_ALIGN.CENTER)

# =============================================
# スライド7: 管理者画面
# =============================================
s7 = blank_slide(prs)
bg(s7, WHITE)

box(s7, 0, 0, 13.33, 1.2, DARK)
txt(s7, "管理者ダッシュボード：一目でわかる現場状況", 0.5, 0.2, 12, 0.8, 26, WHITE, bold=True)

# サマリーカード3つ
cards = [
    ("登録スタッフ数", "24", "名", LINE_GREEN),
    ("出勤済み", "21", "名", RGBColor(0x22, 0x8B, 0xE6)),
    ("要フォロー", "2", "名", RED),
]
for i, (label, num, unit, color) in enumerate(cards):
    bx = 0.4 + i * 2.7
    box(s7, bx, 1.4, 2.4, 1.8, LIGHT_GREEN)
    box(s7, bx, 1.4, 2.4, 0.08, color)
    txt(s7, label, bx + 0.1, 1.55, 2.2, 0.4, 12, GRAY, align=PP_ALIGN.CENTER)
    txt(s7, num, bx + 0.1, 2.0, 1.8, 0.8, 36, color, bold=True, align=PP_ALIGN.CENTER)
    txt(s7, unit, bx + 1.9, 2.45, 0.4, 0.4, 13, GRAY)

# スタッフ一覧テーブル（サンプル）
headers2 = ["名前", "出勤", "退勤", "勤務時間", "コンディション"]
col_w2 = [2.8, 1.6, 1.6, 2.0, 2.8]
col_x2 = [0.3]
for w in col_w2[:-1]:
    col_x2.append(col_x2[-1] + w + 0.05)

box(s7, 0.3, 3.4, 12.5, 0.5, DARK)
for h, x, w in zip(headers2, col_x2, col_w2):
    txt(s7, h, x + 0.1, 3.45, w, 0.4, 12, WHITE, bold=True)

sample = [
    ["田中 花子", "09:01", "18:03", "9h 2m", "😄 絶好調"],
    ["鈴木 一郎", "08:55", "--:--", "勤務中", "😊 良い"],
    ["山田 太郎 ⚠", "10:12", "--:--", "勤務中", "😢 しんどい"],
]
row_colors = [WHITE, RGBColor(0xFA, 0xFA, 0xFA), RGBColor(0xFF, 0xF0, 0xF0)]
for ri, (row, rc) in enumerate(zip(sample, row_colors)):
    ry = 3.95 + ri * 0.68
    box(s7, 0.3, ry, 12.5, 0.65, rc)
    for ci, (cell, x, w) in enumerate(zip(row, col_x2, col_w2)):
        fc = RED if "⚠" in cell else (RED if "しんどい" in cell else DARK)
        txt(s7, cell, x + 0.1, ry + 0.1, w - 0.2, 0.45, 12, fc)

# =============================================
# スライド8: CTA（クロージング）
# =============================================
s8 = blank_slide(prs)
bg(s8, LINE_GREEN)

txt(s8, "まずは30日間、無料でお試しください", 0.5, 1.0, 12.3, 1.0,
    30, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s8, "クレジットカード不要　違約金なし　スタッフ数無制限", 0.5, 2.0, 12.3, 0.6,
    16, RGBColor(0xCC, 0xFF, 0xDD), align=PP_ALIGN.CENTER)

box(s8, 2.5, 2.9, 8.3, 2.2, WHITE)
txt(s8, "デモURL（スマホのLINEで開いてください）", 2.7, 3.0, 7.9, 0.5,
    13, GRAY, align=PP_ALIGN.CENTER)
txt(s8, "https://liff.line.me/2010014245-i7LMCgYl", 2.7, 3.55, 7.9, 0.6,
    15, LINE_GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(s8, "← LINEアプリで開くと実際の画面をご体験いただけます", 2.7, 4.1, 7.9, 0.5,
    11, GRAY, align=PP_ALIGN.CENTER)

multiline(s8, [
    "📧 お問い合わせ先：（メールアドレスを追記してください）",
    "📞 電話番号：（電話番号を追記してください）",
], 0.5, 5.4, 12.3, 1.0, 14, WHITE, align=PP_ALIGN.CENTER)

txt(s8, "ラクラク勤怠　— 派遣スタッフが辞めない現場をLINEで作る —",
    0.5, 6.9, 12.3, 0.45, 12, RGBColor(0xCC, 0xFF, 0xDD), align=PP_ALIGN.CENTER)

# 保存
out = r"c:\Users\PC_User\Desktop\AI動画\rakuraku-kintai\らくらく勤怠\sales\ラクラク勤怠_営業資料.pptx"
prs.save(out)
print(f"保存完了: {out}")
