# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ======== Colors ========
GREEN  = RGBColor(0x06,0xC7,0x55)
DGREEN = RGBColor(0x04,0x9A,0x40)
LGREEN = RGBColor(0xDC,0xFB,0xE5)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x1A,0x1A,0x2E)
GRAY   = RGBColor(0x55,0x55,0x55)
LGRAY  = RGBColor(0xF4,0xF4,0xF4)
RED    = RGBColor(0xDC,0x26,0x26)
LRED   = RGBColor(0xFE,0xE2,0xE2)
ORANGE = RGBColor(0xF5,0x9E,0x0B)
LORG   = RGBColor(0xFF,0xF4,0xE0)
BLUE   = RGBColor(0x1D,0x4E,0xD8)
LBLUE  = RGBColor(0xEF,0xF6,0xFF)
NAVY   = RGBColor(0x0F,0x17,0x2A)
PURPLE = RGBColor(0x6D,0x28,0xD9)
GOLD   = RGBColor(0xD9,0x7A,0x06)
SAVE_DIR = r"C:\Users\PC_User\Desktop\AI動画\rakuraku-kintai\らくらく勤怠\sales"

def new_prs():
    p = Presentation()
    p.slide_width = Inches(13.33)
    p.slide_height = Inches(7.5)
    return p

def rect(s,l,t,w,h,fill=None,line=None,lw=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        if lw: sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def txt(s, text, l, t, w, h, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def cover(prs, title, subtitle, badge=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
    rect(s, 0, 0, 13.33, 0.08, fill=GREEN)
    rect(s, 0, 7.42, 13.33, 0.08, fill=GREEN)
    for i in range(6): rect(s, 10.5+i*0.45, 0, 0.35, 7.5, fill=RGBColor(0x1F,0x2A,0x3E))
    if badge:
        rect(s, 1.0, 1.0, 3.5, 0.42, fill=GREEN)
        txt(s, badge, 1.0, 1.05, 3.5, 0.35, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 1.0, 2.0, 11, 1.6, size=54, bold=True, color=WHITE)
    txt(s, subtitle, 1.0, 3.8, 11, 0.7, size=20, color=GREEN)
    rect(s, 1.0, 4.85, 4.5, 0.06, fill=GREEN)
    txt(s, "ラクラク勤怠   小原 健太", 1.0, 5.0, 8, 0.4, size=13, color=RGBColor(0xAA,0xCC,0xFF))
    txt(s, "biccuri5086@gmail.com   /   080-9895-7770", 1.0, 5.4, 8, 0.35, size=11, color=RGBColor(0x88,0xAA,0xCC))

def hdr(s, title, sub=None, num=None, total=None):
    rect(s, 0, 0, 13.33, 1.25, fill=NAVY)
    rect(s, 0, 1.25, 0.06, 6.25, fill=GREEN)
    txt(s, title, 0.45, 0.15, 11, 0.6, size=24, bold=True, color=WHITE)
    if sub:
        txt(s, sub, 0.45, 0.78, 11, 0.4, size=12, color=RGBColor(0xAA,0xCC,0xFF))
    if num and total:
        txt(s, f"{num} / {total}", 11.8, 7.1, 1.4, 0.35, size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def card_row(s, items, y=1.6, h=2.4):
    n = len(items)
    width = (13.33 - 0.7 - 0.2 * (n - 1)) / n
    for i, (color, title, desc) in enumerate(items):
        x = 0.35 + i * (width + 0.2)
        rect(s, x, y, width, h, fill=LGRAY)
        rect(s, x, y, width, 0.08, fill=color)
        txt(s, title, x + 0.2, y + 0.2, width - 0.4, 0.55, size=14, bold=True, color=DARK)
        txt(s, desc, x + 0.2, y + 0.85, width - 0.4, h - 1.0, size=11, color=GRAY)

def table_slide(s, headers, rows, x=0.35, y=1.5, w=12.6, header_h=0.5, row_h=0.5):
    cols = len(headers)
    col_w = w / cols
    rect(s, x, y, w, header_h, fill=NAVY)
    for i, h_text in enumerate(headers):
        txt(s, h_text, x + i * col_w + 0.1, y + 0.08, col_w - 0.2, header_h - 0.1, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        ry = y + header_h + r * row_h
        bg = WHITE if r % 2 == 0 else LGRAY
        rect(s, x, ry, w, row_h, fill=bg)
        for i, cell in enumerate(row):
            txt(s, str(cell), x + i * col_w + 0.1, ry + 0.1, col_w - 0.2, row_h - 0.1, size=11, color=DARK, align=PP_ALIGN.CENTER)

def step_row(s, steps_data, y=1.7, h=3.5):
    n = len(steps_data)
    width = (13.33 - 0.7 - 0.3 * (n - 1)) / n
    for i, (color, num, title, desc) in enumerate(steps_data):
        x = 0.35 + i * (width + 0.3)
        rect(s, x, y, width, h, fill=LGRAY)
        rect(s, x, y, width, 0.85, fill=color)
        txt(s, num, x, y + 0.05, width, 0.4, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, x, y + 0.45, width, 0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, x + 0.2, y + 1.05, width - 0.4, h - 1.2, size=11, color=DARK)
        if i < n - 1:
            ax = 0.35 + i * (width + 0.3) + width + 0.02
            txt(s, "->", ax, y + 1.6, 0.25, 0.4, size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

def cta(prs, message, sub_msg="", num=None, total=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
    rect(s, 0, 0, 13.33, 0.08, fill=GREEN)
    rect(s, 0, 7.42, 13.33, 0.08, fill=GREEN)
    for i in range(6): rect(s, 10.5+i*0.45, 0, 0.35, 7.5, fill=RGBColor(0x1F,0x2A,0x3E))
    txt(s, "NEXT STEP", 1.0, 1.5, 12, 0.5, size=14, bold=True, color=GREEN, italic=True)
    txt(s, message, 1.0, 2.2, 12, 2.0, size=32, bold=True, color=WHITE)
    if sub_msg:
        txt(s, sub_msg, 1.0, 4.5, 12, 1.5, size=16, color=RGBColor(0xCC,0xDD,0xFF))
    rect(s, 1.0, 6.3, 11.3, 0.85, fill=GREEN)
    txt(s, "ラクラク勤怠   小原 健太", 1.2, 6.4, 11, 0.4, size=14, bold=True, color=WHITE)
    txt(s, "biccuri5086@gmail.com   /   080-9895-7770", 1.2, 6.85, 11, 0.3, size=11, color=LGREEN)
    if num and total: txt(s, f"{num} / {total}", 11.8, 7.1, 1.4, 0.35, size=10, color=RGBColor(0x66,0x77,0x88), align=PP_ALIGN.RIGHT)

# =====================================================
# 00 営業実行ロードマップ
# =====================================================
def deck_00():
    TOTAL = 6
    prs = new_prs()
    cover(prs, "営業実行ロードマップ", "4週間で月4〜7社獲得", "INTERNAL")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "全体像", "コスト対効果の高い順で同時並行", 2, TOTAL)
    headers = ["順", "施策", "開始日数", "月額コスト", "1顧客獲得コスト", "評価"]
    rows = [
        ["1", "Google広告",        "1日",    "30,000円",  "15,000円", "★★★★★"],
        ["2", "社労士パートナー",   "即日",   "0円",       "3,000円",  "★★★★★"],
        ["3", "リスト作成外注",    "3-5日",  "5,000円(単発)", "-",     "★★★★"],
        ["4", "営業代行(成果報酬)", "1-2週",  "15,000円〜", "10,000円", "★★★★"],
    ]
    table_slide(s, headers, rows, y=2.0, row_h=0.6)
    txt(s, "→ 全部仕掛けると4週間後には月4〜7社獲得ペース構築完了", 0.4, 6.4, 12.5, 0.5, size=14, bold=True, color=DGREEN, align=PP_ALIGN.CENTER)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "Week 1：自動集客の仕掛け", "種を撒く", 3, TOTAL)
    step_row(s, [
        (GREEN, "Day 1-2", "Google広告開始",     "2時間で設定完了\n月3万円で集客開始"),
        (BLUE,  "Day 3-4", "社労士10名打診",     "短文メールを送付\n返信2-3件期待"),
        (ORANGE,"Day 5",   "リスト100社発注",    "クラウドワークス\n¥5,000で1週間納品"),
        (PURPLE,"Day 6-7", "デモ準備",          "問合せ・返信対応\n資料の最終確認"),
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "Week 2〜4：人力営業の仕組み化", "刈り取り", 4, TOTAL)
    step_row(s, [
        (NAVY,  "Week 2", "営業代行発注", "リスト納品\nカクトクに月3アポ依頼"),
        (BLUE,  "Week 3", "デモ実施",     "Google広告流入\n社労士紹介の商談"),
        (GREEN, "Week 4", "成約",         "トライアル開始\n手厚いサポート"),
    ], h=3.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "月の予算と作業時間", "ROIシミュレーション", 5, TOTAL)
    card_row(s, [
        (GREEN, "月コスト ¥50,000",     "・Google広告 ¥30k\n・営業代行 ¥15k〜\n・リスト ¥5k(初月)"),
        (BLUE,  "月作業時間 15〜20h",   "・広告チューニング 1h\n・社労士連絡 2h\n・代行レビュー 1h\n・デモ実施 5-10h\n・サポート 5h"),
        (ORANGE,"期待成果 月5社獲得",   "・Google経由 2社\n・社労士経由 2社\n・営業代行経由 1社\n→ MRR ¥100,000+"),
    ], y=1.5, h=4.5)

    cta(prs, "今日から STEP 1：Google広告", "1日で設定完了、明日から問合せが来始める", 6, TOTAL)
    prs.save(SAVE_DIR + r"\00_営業実行ロードマップ.pptx")
    print("00 saved")

# =====================================================
# 01 サービス紹介資料
# =====================================================
def deck_01():
    TOTAL = 6
    prs = new_prs()
    cover(prs, "ラクラク勤怠", "LINE完結 1タップ打刻 × コンディション管理 SaaS", "SERVICE GUIDE")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "サービス概要", "派遣会社向けの新世代勤怠ツール", 2, TOTAL)
    txt(s, "「LINEを開いてボタンを1回押すだけ」", 0.4, 1.5, 12.5, 0.7, size=22, bold=True, color=DGREEN, align=PP_ALIGN.CENTER)
    txt(s, "アプリ不要・当日から使える派遣会社特化の勤怠管理SaaS", 0.4, 2.2, 12.5, 0.5, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    card_row(s, [
        (GREEN, "1タップ打刻",         "LINEで出退勤\nアプリDL不要"),
        (ORANGE,"コンディション報告",   "5段階の絵文字で\n5秒の体調報告"),
        (RED,   "離職アラート",         "スコア低下で\n管理者に通知"),
        (BLUE,  "管理ダッシュボード",   "全員の状況を\nリアルタイム可視化"),
    ], y=3.2, h=3.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "解決する4つの課題", "派遣業界が抱える構造問題", 3, TOTAL)
    card_row(s, [
        (RED,   "離職率の高さ",       "正社員の2-3倍\n採用・教育コスト圧迫"),
        (ORANGE,"紙タイムカード",     "月末集計に数時間\n電帳法対応未"),
        (BLUE,  "既存ツール導入の壁", "アプリDL拒否\nIT不慣れで挫折"),
        (PURPLE,"スタッフとの接点不足","LINE個人連絡のみ\n組織的ケア不可"),
    ], y=1.7, h=4.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "他社にない3つの強み", "競合ゼロの差別化要素", 4, TOTAL)
    card_row(s, [
        (GREEN, "LINE完結",           "アプリDL不要\n当日から使える\n摩擦ゼロ"),
        (ORANGE,"コンディション報告",  "業界唯一の機能\n離職予兆を可視化\n早期介入が可能"),
        (BLUE,  "中小派遣特化",       "大手SaaSより安価\n導入工数最小\n50-300名規模に最適"),
    ], y=1.7, h=4.8)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "料金プラン", "シンプル・低コストエントリー", 5, TOTAL)
    headers = ["項目", "スタータープラン", "スタンダードプラン", "エンタープライズ"]
    rows = [
        ["月額",       "150円/人",     "200円/人",     "応相談"],
        ["対象規模",   "〜50名",       "50〜300名",    "300名以上"],
        ["初期費用",   "0円",          "0円",          "0円"],
        ["無料期間",   "30日",         "30日",         "30日"],
        ["契約縛り",   "なし",         "なし",         "なし"],
    ]
    table_slide(s, headers, rows, y=2.0, row_h=0.55)

    cta(prs, "今すぐ30日間無料トライアル", "クレジットカード不要・違約金なし", 6, TOTAL)
    prs.save(SAVE_DIR + r"\01_サービス紹介資料.pptx")
    print("01 saved")

# =====================================================
# 02 料金プラン
# =====================================================
def deck_02():
    TOTAL = 5
    prs = new_prs()
    cover(prs, "料金プラン", "ラクラク勤怠", "PRICING")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "3つのプラン", "規模に応じた最適価格", 2, TOTAL)
    plans = [
        (GREEN, "スターター", "150円", "/人/月", "〜50名規模", ["1タップ打刻","コンディション報告","管理ダッシュボード","GPS打刻"]),
        (NAVY,  "スタンダード", "200円", "/人/月", "50〜300名規模", ["スターター全機能","打刻修正(管理画面)","CSVエクスポート","LINE通知アラート"]),
        (PURPLE,"エンタープライズ", "応相談", "", "300名以上", ["スタンダード全機能","専用サポート","給与計算API連携","多言語対応"]),
    ]
    for i, (c, name, price, unit, target, feats) in enumerate(plans):
        x = 0.35 + i * 4.3
        rect(s, x, 1.5, 4.1, 5.7, fill=LGRAY); rect(s, x, 1.5, 4.1, 0.5, fill=c)
        txt(s, name, x, 1.53, 4.1, 0.44, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i == 1:
            rect(s, x+1.5, 1.53, 1.1, 0.3, fill=GOLD)
            txt(s, "★ 推奨", x+1.5, 1.56, 1.1, 0.27, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, price, x, 2.15, 4.1, 0.8, size=32, bold=True, color=c, align=PP_ALIGN.CENTER)
        txt(s, unit, x, 2.9, 4.1, 0.35, size=12, color=GRAY, align=PP_ALIGN.CENTER)
        txt(s, target, x, 3.3, 4.1, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)
        for j, f in enumerate(feats):
            txt(s, "✓  " + f, x+0.3, 3.85+j*0.5, 3.7, 0.45, size=11, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "料金シミュレーション", "規模別の月額目安", 3, TOTAL)
    headers = ["登録スタッフ数", "適用プラン", "月額(税抜)", "年額(税抜)", "1人あたり日額換算"]
    rows = [
        ["30名",  "スターター",    "4,500円",     "54,000円",     "5円/人/日"],
        ["50名",  "スタンダード",  "10,000円",    "120,000円",    "6.7円/人/日"],
        ["100名", "スタンダード",  "20,000円",    "240,000円",    "6.7円/人/日"],
        ["200名", "スタンダード",  "40,000円",    "480,000円",    "6.7円/人/日"],
        ["500名", "エンタープライズ","応相談",    "応相談",       "応相談"],
    ]
    table_slide(s, headers, rows, y=2.0, row_h=0.55)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "支払い方法", "中小派遣会社の経理慣行に合わせた柔軟対応", 4, TOTAL)
    card_row(s, [
        (GREEN, "月払い 銀行振込", "毎月末締め\n翌月末払い\n請求書発行"),
        (BLUE,  "年払い 銀行振込", "1年分一括\n5%割引適用\n年度予算組みやすい"),
        (ORANGE,"クレジットカード", "VISA/Master/JCB\n自動継続課金\n少人数規模で便利"),
    ], y=2.0, h=3.0)
    rect(s, 0.35, 5.5, 12.6, 1.5, fill=LGREEN)
    txt(s, "30日間無料トライアル", 0.5, 5.6, 12.3, 0.5, size=16, bold=True, color=DGREEN)
    txt(s, "・期間中の費用は一切発生しません\n・クレジットカード登録不要\n・トライアル終了後の自動課金もありません（明示的にお申込みいただいてからスタート）", 0.5, 6.1, 12.3, 0.9, size=12, color=DARK)

    cta(prs, "まずは30日間 無料でお試しください", "電話・メールでお気軽にお問い合わせください", 5, TOTAL)
    prs.save(SAVE_DIR + r"\02_料金プラン.pptx")
    print("02 saved")

# =====================================================
# 03 営業トーク台本
# =====================================================
def deck_03():
    TOTAL = 7
    prs = new_prs()
    cover(prs, "営業トーク台本", "派遣会社経営者向けデモ進行", "INTERNAL")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "15分デモの基本構成", "1回の商談で目指すゴール", 2, TOTAL)
    step_row(s, [
        (GREEN, "0-2分", "アイスブレイク", "天気・コロナ後の業況雑談\n相手の業種ヒアリング"),
        (BLUE,  "2-5分", "課題ヒアリング", "現状の勤怠管理方法\n離職率・スタッフ数"),
        (ORANGE,"5-12分","デモ・提案",     "LIFFで実際の打刻\n管理画面・コンディション機能"),
        (RED,   "12-15分","クロージング",  "料金提示\n無料トライアル提案\n次回アクション確定"),
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "アイスブレイク・課題ヒアリング", "信頼関係構築と痛点発見", 3, TOTAL)
    rect(s, 0.35, 1.5, 12.6, 5.5, fill=LGRAY)
    txt(s, "オープニング・トーク", 0.5, 1.55, 12.5, 0.45, size=14, bold=True, color=NAVY)
    txt(s,
        "「本日はお時間ありがとうございます！\n"
        "ラクラク勤怠の小原と申します。\n"
        "派遣業界の勤怠管理について、現状の課題から\n"
        "順番にヒアリングさせていただきたいのですが、\n"
        "現在は紙のタイムカードでしょうか、それとも何か\n"
        "システムをお使いでしょうか？」",
        0.5, 2.05, 12.3, 1.8, size=13, color=DARK)
    txt(s, "ヒアリング必須3項目", 0.5, 4.0, 12.5, 0.45, size=14, bold=True, color=NAVY)
    txt(s,
        "1. 登録スタッフ数（料金見積もりに必要）\n"
        "2. 現在の勤怠管理方法（紙／他社／なし）\n"
        "3. 一番困っているのは何か（離職／集計／打刻ミス）",
        0.5, 4.5, 12.3, 2.3, size=13, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "デモパート（7分）", "実機を見せて納得感を出す", 4, TOTAL)
    card_row(s, [
        (GREEN, "1. LINE打刻", "自分のスマホで\n出勤ボタンを押す\nGPSも見せる"),
        (ORANGE,"2. コンディション", "5段階絵文字を選択\n一言コメント入力\n送信"),
        (BLUE,  "3. 管理画面",  "PCで管理画面を表示\nスタッフ一覧\n要フォロー赤色表示"),
    ], y=1.7, h=3.0)
    rect(s, 0.35, 5.0, 12.6, 2.0, fill=LBLUE)
    txt(s, "デモ中の決め台詞", 0.5, 5.1, 12.5, 0.4, size=13, bold=True, color=BLUE)
    txt(s,
        "「これ、スタッフ側はLINE開くだけで使えるんです。\n"
        "コンディション報告は1日数十秒。スコア低下で離職予兆が見えるので、\n"
        "声かけのタイミングを逃さない。これが他社にない最大の強みです。」",
        0.5, 5.5, 12.3, 1.4, size=12, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "クロージング・料金提示", "数字で具体的に示す", 5, TOTAL)
    txt(s, "ヒアリングした人数で即座に試算", 0.4, 1.45, 12.5, 0.5, size=14, bold=True, color=DARK)
    headers = ["スタッフ数", "プラン", "月額", "年額", "1人あたり日額"]
    rows = [
        ["50名",  "スタンダード", "10,000円", "120,000円", "6.7円/人/日"],
        ["100名", "スタンダード", "20,000円", "240,000円", "6.7円/人/日"],
        ["200名", "スタンダード", "40,000円", "480,000円", "6.7円/人/日"],
    ]
    table_slide(s, headers, rows, y=2.0, row_h=0.55)
    rect(s, 0.35, 5.5, 12.6, 1.5, fill=LGREEN)
    txt(s, "決め台詞", 0.5, 5.6, 12.5, 0.4, size=13, bold=True, color=DGREEN)
    txt(s,
        "「1人あたり1日たった6.7円で、紙の集計工数削減＋離職予兆の早期発見が手に入ります。\n"
        "30日間無料でカード登録も不要なので、リスクなく試せます。来週から始めませんか？」",
        0.5, 6.0, 12.3, 1.0, size=12, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "断り文句への切り返し", "よくある3パターン", 6, TOTAL)
    rows = [
        ["「他社使ってる」", "「現在お使いのサービスは？特徴比較資料お送りしましょうか」", "切り口にして関係維持"],
        ["「料金が高い」",   "「1人/月200円、1日6.7円。離職1人防止で30万円節約」", "コスト→投資視点に転換"],
        ["「持ち帰り検討」", "「ご検討中の不安点は？無料トライアルだけでも始めませんか」", "次のアクションを必ず取る"],
        ["「LINEは不安」",   "「2要素認証・暗号化・監査ログ完備。資料お送りします」", "セキュリティ説明資料を活用"],
        ["「決裁者と相談」", "「決裁者向け資料お送りします。私からも一緒に説明可能です」", "決裁者直接の場を作る"],
    ]
    table_slide(s, ["断り文句", "切り返しトーク", "ポイント"], rows, y=1.7, row_h=0.7)

    cta(prs, "デモ後 即日にフォロー", "メール送信＋トライアル設定支援の打診", 7, TOTAL)
    prs.save(SAVE_DIR + r"\03_営業トーク台本.pptx")
    print("03 saved")

# =====================================================
# 04 FAQ
# =====================================================
def deck_04():
    TOTAL = 6
    prs = new_prs()
    cover(prs, "よくあるご質問", "ラクラク勤怠", "FAQ")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "Q1〜Q4：機能・使い方", "", 2, TOTAL)
    qa_list = [
        ("Q1. スタッフはアプリのインストール必要？", "A. 不要です。LINEを開いてURLをタップするだけで利用開始できます。"),
        ("Q2. 何人まで登録できますか？", "A. 人数制限はありません。エンタープライズプランで1,000名超でも対応可能です。"),
        ("Q3. GPS打刻は必須ですか？", "A. 任意です。スタッフがGPS許可しなくても打刻は通常通り動作します。"),
        ("Q4. オフラインでも打刻できますか？", "A. インターネット接続が必要です。圏外時は接続後にまとめて記録されます。"),
    ]
    for i, (q, a) in enumerate(qa_list):
        y = 1.55 + i * 1.35
        rect(s, 0.35, y, 12.6, 1.2, fill=LGRAY); rect(s, 0.35, y, 0.08, 1.2, fill=GREEN)
        txt(s, q, 0.55, y+0.15, 12.3, 0.45, size=13, bold=True, color=DARK)
        txt(s, a, 0.55, y+0.6, 12.3, 0.55, size=11, color=GRAY)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "Q5〜Q8：料金・契約", "", 3, TOTAL)
    qa_list = [
        ("Q5. 最低契約期間はありますか？", "A. ありません。月単位での解約が可能です。"),
        ("Q6. 途中で人数が増えた場合の料金は？", "A. 翌月から自動的に人数分を反映して請求します。日割り計算なし。"),
        ("Q7. 無料トライアル中の制限は？", "A. 全機能を制限なく30日間使えます。クレジットカード登録も不要。"),
        ("Q8. 解約時のデータはどうなる？", "A. 解約後30日間はCSVで全データダウンロード可能。その後完全削除。"),
    ]
    for i, (q, a) in enumerate(qa_list):
        y = 1.55 + i * 1.35
        rect(s, 0.35, y, 12.6, 1.2, fill=LGRAY); rect(s, 0.35, y, 0.08, 1.2, fill=BLUE)
        txt(s, q, 0.55, y+0.15, 12.3, 0.45, size=13, bold=True, color=DARK)
        txt(s, a, 0.55, y+0.6, 12.3, 0.55, size=11, color=GRAY)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "Q9〜Q12：セキュリティ", "", 4, TOTAL)
    qa_list = [
        ("Q9. データは安全ですか？", "A. AWS東京リージョンに暗号化保存。SSL通信。管理画面は2要素認証対応。"),
        ("Q10. 個人情報の取扱いは？", "A. 個人情報保護方針に基づき、勤怠管理目的以外には使用しません。第三者提供なし。"),
        ("Q11. 不正打刻の対策は？", "A. GPS記録（任意）、LINE本人確認、打刻時刻のサーバー記録で改ざん不可。"),
        ("Q12. 管理者の操作ログは見られますか？", "A. はい。監査ログ機能ですべての管理者操作が記録されています。"),
    ]
    for i, (q, a) in enumerate(qa_list):
        y = 1.55 + i * 1.35
        rect(s, 0.35, y, 12.6, 1.2, fill=LGRAY); rect(s, 0.35, y, 0.08, 1.2, fill=PURPLE)
        txt(s, q, 0.55, y+0.15, 12.3, 0.45, size=13, bold=True, color=DARK)
        txt(s, a, 0.55, y+0.6, 12.3, 0.55, size=11, color=GRAY)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "Q13〜Q16：導入・サポート", "", 5, TOTAL)
    qa_list = [
        ("Q13. 導入にどれくらいかかりますか？", "A. 当日中に運用開始可能。アカウント作成5分、スタッフへのURL配布即時。"),
        ("Q14. 既存の勤怠データを移行できますか？", "A. CSVインポートで対応可能。当社で移行作業のサポートも無料で実施。"),
        ("Q15. サポート体制は？", "A. メール・電話（平日10-18時）で対応。導入時は無料で個別サポートあり。"),
        ("Q16. 給与計算ソフトと連携できますか？", "A. CSV出力対応。マネーフォワード等への直接API連携は順次対応予定。"),
    ]
    for i, (q, a) in enumerate(qa_list):
        y = 1.55 + i * 1.35
        rect(s, 0.35, y, 12.6, 1.2, fill=LGRAY); rect(s, 0.35, y, 0.08, 1.2, fill=ORANGE)
        txt(s, q, 0.55, y+0.15, 12.3, 0.45, size=13, bold=True, color=DARK)
        txt(s, a, 0.55, y+0.6, 12.3, 0.55, size=11, color=GRAY)

    cta(prs, "他のご質問もお気軽に", "メール・電話でいつでもお答えいたします", 6, TOTAL)
    prs.save(SAVE_DIR + r"\04_FAQ.pptx")
    print("04 saved")

# =====================================================
# 05 営業1枚資料
# =====================================================
def deck_05():
    TOTAL = 2
    prs = new_prs()

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    rect(s, 0, 0, 13.33, 1.0, fill=GREEN)
    txt(s, "ラクラク勤怠", 0.4, 0.15, 6, 0.5, size=24, bold=True, color=WHITE)
    txt(s, "派遣会社向け 1タップ打刻 × 離職予兆 SaaS", 0.4, 0.6, 8, 0.4, size=13, color=WHITE)
    txt(s, "30日 無料 / 150円〜", 9.5, 0.25, 3.5, 0.5, size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    rect(s, 0.35, 1.3, 12.6, 1.8, fill=LRED)
    txt(s, "派遣会社のリアルなお悩み", 0.5, 1.4, 12.5, 0.4, size=14, bold=True, color=RED)
    txt(s,
        "❌ スタッフが突然来なくなる（離職率は正社員の2-3倍）\n"
        "❌ 打刻アプリを入れてくれないスタッフが多い\n"
        "❌ 紙タイムカードの集計に月数時間取られる\n"
        "❌ 電子帳簿保存法対応がまだ・スタッフの体調把握ができない",
        0.5, 1.85, 12.3, 1.15, size=11, color=DARK)

    rect(s, 0.35, 3.25, 12.6, 1.8, fill=LGREEN)
    txt(s, "ラクラク勤怠で全部解決", 0.5, 3.35, 12.5, 0.4, size=14, bold=True, color=DGREEN)
    txt(s,
        "✅ LINEだけで完結 → アプリDL不要、IT不慣れスタッフでも即日OK\n"
        "✅ コンディション報告（5段階絵文字）で離職予兆を早期発見（業界唯一）\n"
        "✅ 管理ダッシュボードでリアルタイム可視化、要フォローを赤色ハイライト\n"
        "✅ 2要素認証・監査ログ・RLS本番化済 → セキュリティ大手並",
        0.5, 3.8, 12.3, 1.15, size=11, color=DARK)

    headers = ["項目", "ラクラク勤怠", "KING OF TIME", "ジョブカン"]
    rows = [
        ["月額", "150〜200円/人", "330円/人", "200〜400円/人"],
        ["LINE完結", "◎", "✗", "✗"],
        ["コンディション報告", "◎", "✗", "✗"],
        ["離職アラート", "◎", "✗", "✗"],
        ["派遣特化", "◎", "△", "✗"],
    ]
    table_slide(s, headers, rows, y=5.2, w=12.6, header_h=0.42, row_h=0.36)

    rect(s, 0.35, 7.0, 12.6, 0.4, fill=NAVY)
    txt(s, "お問い合わせ：小原 健太  /  080-9895-7770  /  biccuri5086@gmail.com", 0.5, 7.04, 12.3, 0.32, size=12, bold=True, color=WHITE)

    cta(prs, "30日間無料トライアル", "クレジットカード不要・違約金なし・解約自由", 2, TOTAL)
    prs.save(SAVE_DIR + r"\05_営業1枚資料.pptx")
    print("05 saved")

# Run all
deck_00()
deck_01()
deck_02()
deck_03()
deck_04()
deck_05()
print("Batch 1 complete")