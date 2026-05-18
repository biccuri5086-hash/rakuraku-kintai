# -*- coding: utf-8 -*-
import sys
sys.path.insert(0, r"c:\Users\PC_User\Desktop\AI動画\rakuraku-kintai\らくらく勤怠\sales")
from create_decks_batch1 import (
    new_prs, rect, txt, cover, hdr, card_row, table_slide, step_row, cta,
    GREEN, DGREEN, LGREEN, WHITE, DARK, GRAY, LGRAY, RED, LRED, ORANGE, LORG,
    BLUE, LBLUE, NAVY, PURPLE, GOLD, SAVE_DIR,
)
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# =====================================================
# 11 Google広告セットアップ手順書
# =====================================================
def deck_11():
    TOTAL = 7
    prs = new_prs()
    cover(prs, "Google広告セットアップ手順書", "1〜2時間で月3万円の自動集客開始", "INTERNAL")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "必要なもの", "用意するのは3つだけ", 2, TOTAL)
    card_row(s, [
        (GREEN, "Googleアカウント", "普段使いのGmailで\nそのまま登録可能"),
        (BLUE,  "クレジットカード", "広告費の引落用\nVISA/Master/JCB"),
        (ORANGE,"予算 30,000円〜",  "1ヶ月分\n月単位で停止可能"),
    ], y=1.7, h=4.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "STEP 1〜2：アカウント作成 + キャンペーン", "60分の作業", 3, TOTAL)
    step_row(s, [
        (GREEN, "STEP 1", "アカウント作成", "ads.google.com\n→「開始する」\n→目標：ウェブサイト訪問\n→ラクラク勤怠を登録"),
        (BLUE,  "STEP 2", "キャンペーン作成", "・タイプ：検索広告\n・入札：クリック最大化\n・1日予算 ¥1,000\n・地域 日本全国"),
    ], y=1.7, h=4.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "STEP 3：キーワード15個（コピペで使う）", "派遣業界をピンポイントで狙う", 4, TOTAL)
    rect(s, 0.35, 1.6, 6.0, 5.4, fill=LGRAY)
    txt(s, "派遣業界キーワード", 0.5, 1.7, 5.7, 0.4, size=13, bold=True, color=NAVY)
    txt(s,
        "派遣 勤怠管理\n"
        "派遣会社 勤怠\n"
        "派遣社員 打刻\n"
        "派遣 タイムカード\n"
        "派遣 LINE 打刻\n"
        "派遣 勤怠アプリ\n"
        "派遣 離職防止\n"
        "派遣 出勤管理\n"
        "派遣 シフト管理",
        0.5, 2.2, 5.7, 4.5, size=12, color=DARK)
    rect(s, 6.95, 1.6, 6.0, 5.4, fill=LGRAY)
    txt(s, "一般キーワード", 7.1, 1.7, 5.7, 0.4, size=13, bold=True, color=NAVY)
    txt(s,
        "勤怠管理 SaaS\n"
        "勤怠管理 LINE\n"
        "勤怠管理 アプリ不要\n"
        "タイムカード 電子化\n"
        "紙タイムカード 廃止\n"
        "電子帳簿保存法 勤怠",
        7.1, 2.2, 5.7, 4.5, size=12, color=DARK)

    s = prs.slides.add_slide(prs.slide_invocation if False else prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "STEP 4：広告文（半角30文字）", "そのままコピペでOK", 5, TOTAL)
    rect(s, 0.35, 1.5, 12.6, 4.5, fill=LGRAY)
    txt(s, "見出し（3つ以上）", 0.5, 1.6, 12.5, 0.4, size=13, bold=True, color=NAVY)
    txt(s,
        "・派遣スタッフ向け勤怠管理SaaS\n"
        "・LINEで1タップ打刻 アプリ不要\n"
        "・月150円〜 30日間無料トライアル\n"
        "・コンディション報告で離職防止\n"
        "・中小派遣会社特化 当日導入可能\n"
        "・初期費用0円 GPS打刻にも対応",
        0.5, 2.1, 12.3, 2.0, size=12, color=DARK)
    txt(s, "説明文（半角90文字、2つ以上）", 0.5, 4.2, 12.5, 0.4, size=13, bold=True, color=NAVY)
    txt(s,
        "LINEを開いてボタン1つ、アプリDL不要。紙タイムカード廃止と離職予兆の早期発見を同時に実現。\n"
        "紙タイムカードでお困りの中小派遣会社向け。スタッフの体調を5段階で把握、業界唯一のSaaS。",
        0.5, 4.65, 12.3, 1.3, size=12, color=DARK)
    rect(s, 0.35, 6.15, 12.6, 0.9, fill=LBLUE)
    txt(s, "最終ページURL", 0.5, 6.2, 4, 0.35, size=12, bold=True, color=BLUE)
    txt(s, "https://rakuraku-kintai-frb6.vercel.app/lp", 0.5, 6.55, 12.3, 0.4, size=13, bold=True, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "STEP 5：効果チェック（1週間後）", "数値で改善判断", 6, TOTAL)
    rows = [
        ["表示回数",       "1,000以上/週",    "少ない → キーワード追加"],
        ["クリック率(CTR)", "3%以上",          "低い → 広告文を見直し"],
        ["クリック単価",    "¥100以下",       "高い → 入札戦略を変更"],
        ["コンバージョン",  "1件以上/週",      "少ない → LP改善"],
    ]
    table_slide(s, ["指標", "目標値", "対処"], rows, y=1.7, row_h=0.7)
    rect(s, 0.35, 5.7, 12.6, 1.3, fill=LRED)
    txt(s, "⚠ 最初の1週間は触らない", 0.5, 5.8, 12.5, 0.4, size=14, bold=True, color=RED)
    txt(s, "設定直後はAIの学習期間。最低7日間は数値を見ず放置。焦って調整するとAIの学習がリセットされて非効率に。",
        0.5, 6.25, 12.3, 0.7, size=11, color=DARK)

    cta(prs, "今日2時間で月3万円の自動集客を構築", "明日から問い合わせメールが届き始めます", 7, TOTAL)
    prs.save(SAVE_DIR + r"\11_Google広告セットアップ手順書.pptx")
    print("11 saved")

# =====================================================
# 12 社労士短文打診メール
# =====================================================
def deck_12():
    TOTAL = 6
    prs = new_prs()
    cover(prs, "社労士事務所への打診メール", "1日5-10件送付 → 月5社の紹介ルート", "INTERNAL")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "3段階アプローチ戦略", "返信率を3倍に上げる流れ", 2, TOTAL)
    step_row(s, [
        (GREEN, "1通目", "超短文メール", "本文150字以下\n10秒で読める\n返信率を最大化"),
        (BLUE,  "2通目", "詳細PPT送付",   "提案書(PPT)添付\n15分Zoom面談打診\n提携の合意取得"),
        (ORANGE,"3通目", "正式提案・契約", "業務委託契約締結\n紹介開始\n継続収入スタート"),
    ], y=1.7, h=4.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "1通目：超短文打診メール", "コピペで使える定型文", 3, TOTAL)
    rect(s, 0.35, 1.5, 12.6, 1.0, fill=LBLUE)
    txt(s, "件名：派遣会社向けSaaSの紹介パートナーをお願いできませんか？", 0.5, 1.75, 12.3, 0.6, size=14, bold=True, color=BLUE)
    rect(s, 0.35, 2.7, 12.6, 4.3, fill=LGRAY)
    txt(s,
        "○○社会保険労務士事務所\n"
        "○○ ○○ 様\n\n"
        "突然のご連絡失礼いたします。ラクラク勤怠の小原健太と申します。\n\n"
        "派遣会社向けに「LINEで1タップ打刻＋コンディション報告」が\n"
        "できるSaaSを開発・運営しております。\n\n"
        "派遣業界の社労士の先生方と紹介パートナー契約を進めており、\n"
        "○○先生にも是非ご検討いただきたくご連絡しました。\n\n"
        "【提携の概要】\n"
        "・初回成約ボーナス ¥10,000\n"
        "・継続紹介手数料 月額の15〜20%（12ヶ月支払）\n"
        "・先生の手間は最小限（紹介の一言だけ、商談は弊社が対応）\n\n"
        "詳細をご希望でしたら、提案書PDFをお送りいたします。",
        0.5, 2.85, 12.3, 4.1, size=11, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "社労士リストの作り方", "全て無料・30分で完成", 4, TOTAL)
    rows = [
        ["⭐⭐⭐⭐⭐", "日本社会保険労務士連合会", "shakaihokenroumushi.jp で都道府県別検索"],
        ["⭐⭐⭐⭐",   "Google検索",                "「○○市 社会保険労務士」「派遣業 社労士」"],
        ["⭐⭐⭐",     "商工会議所の会員名簿",       "地域密着型を発見"],
    ]
    table_slide(s, ["評価", "情報源", "使い方"], rows, y=1.7, row_h=0.75)
    rect(s, 0.35, 4.5, 12.6, 2.5, fill=LGREEN)
    txt(s, "ターゲット社労士の特徴", 0.5, 4.6, 12.5, 0.4, size=14, bold=True, color=DGREEN)
    txt(s,
        "・派遣業を顧問先に持っている事務所（即紹介可能）\n"
        "・製造・物流系企業を顧問先に持つ（派遣会社との取引が多い）\n"
        "・HPに労務相談を打ち出している（営業に積極的）\n"
        "・40〜50代の所長（新サービス導入に前向き）\n"
        "・地方の中堅事務所(社員5〜20名)（大手より話を聞いてくれる）",
        0.5, 5.05, 12.3, 1.95, size=12, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "KPIと1週間の実行プラン", "週1〜2件の提携合意ペース", 5, TOTAL)
    headers = ["指標", "目標値"]
    rows = [
        ["1日送付数",       "5〜10件"],
        ["返信率",          "10〜20%"],
        ["面談実施率",      "50%"],
        ["提携契約率",      "30〜50%"],
        ["週の提携合意",    "1〜2件"],
    ]
    table_slide(s, headers, rows, x=0.35, y=1.7, w=6.0, row_h=0.55)
    rect(s, 6.7, 1.7, 6.3, 5.3, fill=LBLUE)
    txt(s, "1週間の実行例", 6.85, 1.8, 6.0, 0.4, size=13, bold=True, color=BLUE)
    txt(s,
        "月 リスト50件作成（1時間）\n"
        "火朝 10件メール送付（30分）\n"
        "水朝 10件メール送付（30分）\n"
        "木朝 10件メール送付（30分）\n"
        "金朝 10件 + 返信対応（1時間）\n"
        "翌月 面談調整・実施（1時間×n）",
        6.85, 2.3, 6.1, 4.5, size=12, color=DARK)

    cta(prs, "今週10件 メール送付してみる", "テンプレはコピペするだけ", 6, TOTAL)
    prs.save(SAVE_DIR + r"\12_社労士短文打診メール.pptx")
    print("12 saved")

# =====================================================
# 13 クラウドワークス発注テンプレート
# =====================================================
def deck_13():
    TOTAL = 6
    prs = new_prs()
    cover(prs, "クラウドワークス発注テンプレ", "派遣会社100社リストを ¥5,000 で外注", "INTERNAL")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "ゴールと発注先", "1週間以内に100社リスト納品", 2, TOTAL)
    rect(s, 0.35, 1.5, 12.6, 1.3, fill=LGREEN)
    txt(s, "「派遣会社100社リスト」を ¥5,000 程度で発注", 0.5, 1.65, 12.5, 0.8, size=20, bold=True, color=DGREEN, align=PP_ALIGN.CENTER)
    card_row(s, [
        (GREEN, "クラウドワークス", "crowdworks.jp\n大手・ワーカー最多\n推奨"),
        (BLUE,  "ランサーズ",        "lancers.jp\n大手・対抗馬"),
        (ORANGE,"ココナラ",          "coconala.com\n個人サービス向け"),
    ], y=3.2, h=3.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "発注内容（コピペ用）", "そのまま貼り付ければOK", 3, TOTAL)
    rect(s, 0.35, 1.5, 12.6, 5.5, fill=LGRAY)
    txt(s, "タイトル：派遣会社100社のリストアップ作業（厚労省サイト・公開情報のみ）", 0.5, 1.6, 12.3, 0.5, size=12, bold=True, color=NAVY)
    txt(s,
        "【依頼内容】\n"
        "日本国内の人材派遣会社のリストをExcelファイルで作成いただきたいです。\n\n"
        "【納品形式】\n"
        "100行のExcelファイル。カラムは：番号 / 会社名 / 都道府県 / 市区町村 / 住所 /\n"
        "電話番号 / HP URL / メールアドレス / 派遣業種 / 推定スタッフ数\n\n"
        "【データソース】公開情報のみ\n"
        "・厚生労働省 人材サービス総合サイト\n・公式HP\n・Google検索\n\n"
        "【絞り込み条件】\n"
        "・社員5-30名規模、登録スタッフ50-300名規模\n・製造/物流/食品/軽作業系\n"
        "・大手（パーソル等）は除外\n\n"
        "【納期】1週間以内　　【予算】¥5,000",
        0.5, 2.15, 12.3, 4.7, size=11, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "応募者選びのポイント", "10件以上来た中から1名選定", 4, TOTAL)
    rows = [
        ["⭐ 過去の評価",        "★4.5以上、20件以上の評価あり"],
        ["⭐ 似た案件の実績",     "「リスト作成」のプロフィール記載"],
        ["⭐ 自己紹介の質",       "テンプレではなく案件に即した提案"],
        ["⭐ サンプル提示意欲",   "「先にサンプル送ります」"],
        ["⭐ レスポンスの速さ",   "24時間以内の応募"],
    ]
    table_slide(s, ["評価項目", "判断基準"], rows, y=1.7, row_h=0.7)
    rect(s, 0.35, 5.7, 12.6, 1.3, fill=LRED)
    txt(s, "⚠ 避けるべき応募者", 0.5, 5.8, 12.5, 0.4, size=13, bold=True, color=RED)
    txt(s, "・評価0または★3未満　・「いくらでもやります」価格訴求のみ　・日本語が不自然　・テンプレ大量送付",
        0.5, 6.25, 12.3, 0.7, size=11, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "適正価格と納品後の活用", "「お金で時間を買う」発想", 5, TOTAL)
    headers = ["件数", "標準予算", "備考"]
    rows = [
        ["50社",   "¥3,000",  "安価でも質の良いワーカーあり"],
        ["100社",  "¥5,000",  "コスパが一番良い"],
        ["200社",  "¥10,000", "慣れたワーカーに"],
        ["500社",  "¥20,000〜", "経験豊富なワーカー必須"],
    ]
    table_slide(s, headers, rows, x=0.35, y=1.7, w=6.0, row_h=0.55)
    rect(s, 6.7, 1.7, 6.3, 5.0, fill=LBLUE)
    txt(s, "リスト納品後の活用法", 6.85, 1.8, 6.0, 0.4, size=13, bold=True, color=BLUE)
    txt(s,
        "Path 1: 自分でメール送付\n"
        "  →「08_営業メールテンプレ」を使用\n\n"
        "Path 2: 営業代行に渡してテレアポ\n"
        "  →「14_営業代行発注ガイド」参照\n\n"
        "Path 3: FAX一斉送信\n"
        "  → ネクスウェイ等で1通¥10〜30\n  → 派遣業界ではまだFAX有効",
        6.85, 2.3, 6.1, 4.2, size=12, color=DARK)

    cta(prs, "今週中にクラウドワークスで発注", "¥5,000で1週間後にリスト100社が手に入る", 6, TOTAL)
    prs.save(SAVE_DIR + r"\13_クラウドワークス発注テンプレート.pptx")
    print("13 saved")

# =====================================================
# 14 営業代行発注ガイド
# =====================================================
def deck_14():
    TOTAL = 7
    prs = new_prs()
    cover(prs, "営業代行発注ガイド", "成果報酬で月3-10件のデモアポ取得", "INTERNAL")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "営業代行の選び方", "重要な3軸", 2, TOTAL)
    headers = ["軸", "推奨", "避けるべき"]
    rows = [
        ["料金体系",   "成果報酬型（1アポ¥5k〜15k）",  "着手金・月額固定型"],
        ["専門領域",   "SaaS・B2B経験あり",            "訪問販売中心の代行会社"],
        ["規模",       "個人〜小規模代行",              "大手（最低契約金額が高い）"],
    ]
    table_slide(s, headers, rows, y=1.7, row_h=0.75)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "おすすめサービス3選", "推奨度順", 3, TOTAL)
    card_row(s, [
        (GREEN, "1.カクトク ⭐⭐⭐⭐⭐", "フリーランス営業マッチング\n1アポ¥5,000〜の成果報酬\n着手金なし\nkakutoku.jp"),
        (BLUE,  "2.ウィルオブ・ワーク",  "SaaS向けインサイドセールス\n月額固定¥20-30万\n体制整ってきたら"),
        (ORANGE,"3.クラウドワークス",    "個人副業の営業マン\n1アポ¥3,000〜から\n質はピンキリで選定重要"),
    ], y=1.7, h=4.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "カクトクへの発注内容", "コピペで使える定型文", 4, TOTAL)
    rect(s, 0.35, 1.5, 12.6, 5.5, fill=LGRAY)
    txt(s, "案件タイトル：派遣会社向けSaaSのテレアポ営業（成果報酬・月3アポから）", 0.5, 1.6, 12.3, 0.5, size=12, bold=True, color=NAVY)
    txt(s,
        "【サービス】ラクラク勤怠 - 派遣会社向けLINE完結勤怠SaaS\n\n"
        "【ターゲット】中小派遣会社(社員5-30名、登録スタッフ50-300名)\n"
        "             製造・物流・食品・軽作業を扱う会社\n"
        "             決裁者：社長 or 営業部長\n\n"
        "【依頼内容】\n"
        "リスト100社からデモアポイントを取得\n"
        "目標：月3〜10件のアポ取得\n\n"
        "【支払い】完全成果報酬：1アポイント¥10,000\n"
        "       （取得アポが商談実施まで進んだ場合のみ支払い）\n\n"
        "【リスト】こちらで用意（クラウドワークスで作成済み）",
        0.5, 2.15, 12.3, 4.7, size=11, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "ROIシミュレーション", "投資対効果は十分", 5, TOTAL)
    card_row(s, [
        (GREEN, "ケース1：月¥30k投資",
         "・月アポ取得：3件\n・成約率30% = 1社\n・月売上：¥20,000\n・1社のLTV：¥240,000\n→ 1ヶ月でROI +¥210,000"),
        (BLUE,  "ケース2：月¥100k投資",
         "・月アポ取得：10件\n・成約率30% = 3社\n・月売上：¥60,000\n・3社のLTV：¥720,000\n→ 半年でROI +¥260,000"),
    ], y=1.7, h=5.0)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "落とし穴 7つ", "知っておくべき注意点", 6, TOTAL)
    rows = [
        ["1", "着手金型の代行は避ける",       "成果報酬型一択"],
        ["2", "安すぎる単価は質が下がる",     "1アポ¥10,000前後"],
        ["3", "リストは自分で用意",           "代行任せると質が落ちる"],
        ["4", "スクリプトを丸投げしない",     "自分の想いを込める"],
        ["5", "月1回モニタリング架電",        "現場感を失わない"],
        ["6", "アポの「質」をフィードバック",  "質を上げてもらう"],
        ["7", "契約解除条件を明記",           "3ヶ月で見切る"],
    ]
    table_slide(s, ["#", "落とし穴", "対策"], rows, y=1.7, row_h=0.55)

    cta(prs, "今月中にカクトクに案件登録", "1週間で営業代行体制が立ち上がる", 7, TOTAL)
    prs.save(SAVE_DIR + r"\14_営業代行発注ガイド.pptx")
    print("14 saved")

deck_11()
deck_12()
deck_13()
deck_14()
print("Batch 3 complete")