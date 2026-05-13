# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

GREEN  = RGBColor(0x06,0xC7,0x55)
DGREEN = RGBColor(0x04,0x9A,0x40)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARK   = RGBColor(0x1A,0x1A,0x2E)
GRAY   = RGBColor(0x55,0x55,0x55)
LGRAY  = RGBColor(0xF4,0xF4,0xF4)
RED    = RGBColor(0xDC,0x26,0x26)
ORANGE = RGBColor(0xF5,0x9E,0x0B)
BLUE   = RGBColor(0x1D,0x4E,0xD8)
LBLUE  = RGBColor(0xEF,0xF6,0xFF)
LGREEN = RGBColor(0xDC,0xFB,0xE5)
GOLD   = RGBColor(0xD9,0x7A,0x06)
NAVY   = RGBColor(0x0F,0x17,0x2A)
PURPLE = RGBColor(0x6D,0x28,0xD9)
blank  = prs.slide_layouts[6]
TOTAL  = 16

def rect(slide,l,t,w,h,fill=None,line=None,lw=None):
    s=slide.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb=fill
    s.line.fill.background()
    if line:
        s.line.color.rgb=line
        if lw: s.line.width=lw
    else: s.line.fill.background()
    return s

def txt(slide,text,l,t,w,h,size=14,bold=False,color=DARK,align=PP_ALIGN.LEFT,wrap=True,italic=False):
    tb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.word_wrap=wrap
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color
    return tb

def hdr(slide,title,sub=None):
    rect(slide,0,0,13.33,1.25,fill=NAVY)
    rect(slide,0,1.25,0.06,6.25,fill=GREEN)
    txt(slide,title,0.45,0.12,11,0.65,size=26,bold=True,color=WHITE)
    if sub: txt(slide,sub,0.45,0.78,11,0.4,size=13,color=RGBColor(0xAA,0xCC,0xFF))

def pnum(slide,n):
    txt(slide,f"{n} / {TOTAL}",11.8,7.1,1.4,0.35,size=11,color=RGBColor(0x99,0x99,0x99),align=PP_ALIGN.RIGHT)

# ---- S1 表紙 ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=NAVY)
rect(s,0,0,13.33,0.06,fill=GREEN)
rect(s,0,7.44,13.33,0.06,fill=GREEN)
for i in range(6): rect(s,10.5+i*0.45,0,0.35,7.5,fill=RGBColor(0x1F,0x2A,0x3E))
txt(s,"INVESTOR PRESENTATION",1.0,0.9,11,0.45,size=13,bold=True,color=GREEN,align=PP_ALIGN.LEFT,italic=True)
txt(s,"ラクラク勤怠",1.0,1.45,11,1.5,size=60,bold=True,color=WHITE,align=PP_ALIGN.LEFT)
txt(s,"LINE完結・1タップ打刻 x コンディション管理 SaaS",1.0,3.0,11,0.6,size=20,color=GREEN,align=PP_ALIGN.LEFT)
txt(s,"派遣会社の「離職防止」と「勤怠DX」を同時に実現する",1.0,3.65,11,0.5,size=16,color=RGBColor(0xCC,0xDD,0xFF),align=PP_ALIGN.LEFT)
rect(s,1.0,4.5,4.5,0.06,fill=GREEN)
txt(s,"2026年5月  ラクラク勤怠 株式会社（準備中）",1.0,4.65,8,0.4,size=12,color=RGBColor(0x88,0xAA,0xCC))
txt(s,"CONFIDENTIAL  本資料の無断転載・複製を禁じます",1.0,5.1,11,0.35,size=11,color=RGBColor(0x66,0x77,0x88),italic=True)
pnum(s,1)

# ---- S2 目次 ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=LGRAY)
hdr(s,"目次","CONTENTS")
pnum(s,2)
items=[
    ("01","課題","派遣会社が抱える構造的な問題"),
    ("02","ソリューション","ラクラク勤怠が解決すること"),
    ("03","市場規模","TAM / SAM / SOM"),
    ("04","競合優位性","競合比較とポジショニング"),
    ("05","ビジネスモデル","価格・収益構造"),
    ("06","Go-to-Market","顧客獲得戦略"),
    ("07","収支計画（1年）","月次損益シミュレーション"),
    ("08","収支計画（3年）","中期成長シナリオ"),
    ("09","KPI・マイルストーン","達成指標と計画"),
    ("10","調達計画","資金調達額・使途"),
    ("11","リスクと対策","主要リスクの認識と軽減策"),
    ("12","チーム・体制","経営陣とアドバイザー"),
]
for i,(num,t,d) in enumerate(items):
    col,row=i%3,i//3
    x,y=0.35+col*4.3,1.45+row*1.45
    rect(s,x,y,4.1,1.3,fill=WHITE,line=NAVY,lw=8000)
    rect(s,x,y,0.6,1.3,fill=NAVY)
    txt(s,num,x,y+0.25,0.6,0.65,size=15,bold=True,color=GREEN,align=PP_ALIGN.CENTER)
    txt(s,t,x+0.7,y+0.1,3.3,0.45,size=14,bold=True,color=DARK)
    txt(s,d,x+0.7,y+0.6,3.3,0.6,size=10,color=GRAY)

# ---- S3 課題 ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"01  課題","派遣業界が抱える「未解決の構造問題」")
pnum(s,3)
txt(s,"日本には約13,000社の派遣会社があり、162万人の派遣労働者が働いている。しかし現場では深刻な問題が放置されている。",0.4,1.35,12.5,0.5,size=13,color=DARK)
problems=[
    (RED,   "離職率が正社員の2〜3倍","突然来なくなる・無断欠勤・早期離職が多発\n体調不良のサインを早期に把握できていない"),
    (ORANGE,"紙タイムカード・手集計が主流","管理者が月末に集計で数時間消費\n電子帳簿保存法対応も急務"),
    (BLUE,  "既存ツールが使えない","アプリDLを嫌がるスタッフが多数\nIT不慣れな中小派遣会社はツール導入できず"),
    (PURPLE,"スタッフとの接点が少ない","派遣元との連絡はLINE個人アカウントのみ\n組織的なケアができずエンゲージメント低下"),
]
for i,(c,ttl,desc) in enumerate(problems):
    x=0.35+(i%2)*6.35; y=2.0+(i//2)*2.45
    rect(s,x,y,6.1,2.2,fill=LGRAY); rect(s,x,y,0.08,2.2,fill=c)
    txt(s,ttl,x+0.2,y+0.15,5.7,0.55,size=14,bold=True,color=DARK)
    txt(s,desc,x+0.2,y+0.75,5.7,1.3,size=11,color=GRAY)
rect(s,0.35,6.85,12.6,0.42,fill=RGBColor(0xFF,0xEE,0xEE))
txt(s,"→ 派遣社員1人の離職・採用・教育コストは平均30万〜50万円。離職を1人防ぐだけで年間数十万円の節約になる。",0.5,6.88,12.3,0.38,size=11,bold=True,color=RED)

# ---- S4 ソリューション ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"02  ソリューション","ラクラク勤怠が解決すること")
pnum(s,4)
txt(s,"「LINEを開いてボタンを1回押すだけ」 アプリ不要・当日から使える派遣特化の勤怠管理SaaS",0.4,1.35,12.5,0.5,size=14,bold=True,color=DARK)
rect(s,0.35,2.0,12.6,4.8,fill=LGRAY)
features=[
    (GREEN, "1タップ打刻",        "LINEを開いてボタン1つ\nアプリDL不要で即日開始"),
    (ORANGE,"コンディション報告",  "5段階の絵文字で\n5秒で今日の調子を報告"),
    (RED,   "離職アラート",        "スコア低下で管理者に\n即時通知・早期介入"),
    (BLUE,  "管理ダッシュボード",  "全スタッフの状況を\nリアルタイムで可視化"),
    (PURPLE,"GPS打刻確認",         "出勤時の位置情報を\n自動記録（バックグラウンド）"),
]
for i,(c,ttl,desc) in enumerate(features):
    x=0.6+i*2.42
    rect(s,x,2.15,2.2,4.45,fill=WHITE); rect(s,x,2.15,2.2,0.08,fill=c)
    txt(s,ttl,x,3.2,2.2,0.5,size=12,bold=True,color=DARK,align=PP_ALIGN.CENTER)
    txt(s,desc,x+0.1,3.8,2.0,1.3,size=10,color=GRAY,align=PP_ALIGN.CENTER)
rect(s,0.35,7.0,12.6,0.38,fill=LGREEN)
txt(s,"競合他社にはない「コンディション報告 x 離職アラート」が最大の差別化。派遣会社の離職防止と勤怠DXを同時解決。",0.5,7.04,12.3,0.32,size=11,bold=True,color=DGREEN)

# ---- S5 市場規模 ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"03  市場規模","TAM / SAM / SOM")
pnum(s,5)
rect(s,1.2,1.5,7.8,5.6,fill=LGRAY)
txt(s,"TAM  日本HRテック市場全体",1.35,1.6,7.5,0.45,size=13,bold=True,color=DARK,align=PP_ALIGN.CENTER)
txt(s,"3,500億円  （CAGR 12〜15%）",1.35,6.5,7.5,0.45,size=12,color=GRAY,align=PP_ALIGN.CENTER)
rect(s,2.0,2.1,6.0,4.0,fill=LGREEN)
txt(s,"SAM  勤怠管理SaaS x 派遣特化",2.15,2.2,5.7,0.45,size=13,bold=True,color=DGREEN,align=PP_ALIGN.CENTER)
txt(s,"450億円（推定）",2.15,5.5,5.7,0.4,size=12,bold=True,color=DGREEN,align=PP_ALIGN.CENTER)
rect(s,3.1,2.75,3.8,2.7,fill=GREEN)
txt(s,"SOM",3.1,2.85,3.8,0.45,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"中小派遣会社（50〜300名）",3.1,3.3,3.8,0.45,size=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"45億円\n（3年内獲得目標）",3.1,3.9,3.8,0.9,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"市場背景",8.6,1.45,4.5,0.4,size=14,bold=True,color=DARK)
drivers=[
    "電子帳簿保存法の厳格化\n紙タイムカード廃止が加速",
    "派遣労働者数：約162万人\n派遣会社登録数：約13,000社",
    "LINE利用率 約86%\nスタッフへの浸透率が高い",
    "HRテック市場 年12〜15%成長\n中小向けは空白地帯",
]
for i,d in enumerate(drivers):
    y=2.0+i*1.22
    rect(s,8.6,y,4.5,1.05,fill=LGRAY)
    txt(s,d,8.75,y+0.1,4.2,0.88,size=11,color=DARK)

# ---- S6 競合優位性 ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"04  競合優位性","「派遣特化 x LINE完結 x コンディション管理」は競合ゼロ")
pnum(s,6)
heads2=["機能","KING OF TIME\n330円/人","ジョブカン\n200円/人","SmartHR\n要見積","ラクラク勤怠\n150〜200円/人"]
rows2=[
    ["LINEで完結（アプリ不要）","x","x","x","◎"],
    ["コンディション報告",      "x","x","x","◎"],
    ["離職アラート機能",        "x","x","x","◎"],
    ["GPS打刻",                "o","△(上位のみ)","x","o"],
    ["初期費用0円",            "o","o","o","◎"],
    ["当日導入可",             "△","△","x","◎"],
    ["派遣特化設計",           "△","x","x","◎"],
    ["中小向け価格",           "△","o","x","◎"],
]
col_x2=[0.35,4.2,6.2,8.15,10.1]; col_w2=[3.8,1.95,1.95,1.95,1.95]
hcols=[NAVY,RGBColor(0x37,0x41,0x51),RGBColor(0x37,0x41,0x51),RGBColor(0x37,0x41,0x51),GREEN]
for j,(cx,cw,cf,h) in enumerate(zip(col_x2,col_w2,hcols,heads2)):
    rect(s,cx,1.35,cw,0.72,fill=cf)
    txt(s,h,cx+0.05,1.38,cw-0.1,0.65,size=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
for i,row in enumerate(rows2):
    y=2.12+i*0.62
    for j,(cx,cw,cell) in enumerate(zip(col_x2,col_w2,row)):
        bg=LGREEN if j==4 else (WHITE if i%2==0 else LGRAY)
        rect(s,cx,y,cw,0.58,fill=bg)
        fc=GREEN if (cell=="◎" and j==4) else (RED if cell=="x" else DARK)
        bld=cell=="◎" and j==4
        txt(s,cell,cx+0.05,y+0.1,cw-0.1,0.4,size=11,bold=bld,color=fc,align=PP_ALIGN.CENTER)
rect(s,0.35,7.08,12.6,0.32,fill=LGREEN)
txt(s,"◎ = 業界唯一。「コンディション報告」「離職アラート」「LINE完結」の3機能を同時に持つ競合は存在しない。",0.5,7.1,12.3,0.28,size=11,bold=True,color=DGREEN)

# ---- S7 ビジネスモデル ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"05  ビジネスモデル","シンプルな月額従量課金・低価格エントリー戦略")
pnum(s,7)
plans=[
    ("スターター","150円","/人/月","〜50名規模",GREEN,["1タップ打刻","コンディション報告","管理ダッシュボード","GPS打刻"]),
    ("スタンダード","200円","/人/月","50〜300名規模",NAVY,["スターター全機能","打刻修正（管理画面）","CSVエクスポート","LINE通知アラート"]),
    ("エンタープライズ","応相談","/人/月","300名以上",PURPLE,["スタンダード全機能","専用サポート","給与計算API連携","多言語対応（英/中/越）"]),
]
for i,(name,price,unit,target,c,feats) in enumerate(plans):
    x=0.35+i*4.3
    rect(s,x,1.4,4.1,5.75,fill=LGRAY); rect(s,x,1.4,4.1,0.5,fill=c)
    txt(s,name,x,1.43,4.1,0.44,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    if i==1:
        rect(s,x+1.55,1.43,1.5,0.3,fill=GOLD)
        txt(s,"★ 推奨",x+1.6,1.45,1.4,0.28,size=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s,price,x,2.05,4.1,0.75,size=32,bold=True,color=c,align=PP_ALIGN.CENTER)
    txt(s,unit,x,2.8,4.1,0.35,size=12,color=GRAY,align=PP_ALIGN.CENTER)
    txt(s,f"対象：{target}",x,3.2,4.1,0.38,size=11,color=GRAY,align=PP_ALIGN.CENTER)
    for j,f in enumerate(feats):
        txt(s,f"  {f}",x+0.3,3.68+j*0.52,3.7,0.45,size=11,color=DARK)
rect(s,0.35,7.08,12.6,0.32,fill=LBLUE)
txt(s,"収益モデルのポイント：顧客1社（平均100名）あたり月2万円。顧客300社でMRR 600万円 = 年商7,200万円",0.5,7.1,12.3,0.28,size=11,bold=True,color=BLUE)

# ---- S8 GTM ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"06  Go-to-Market戦略","3フェーズで中小派遣会社を攻略する")
pnum(s,8)
phases=[
    ("Phase 1  M1〜M3","モニター獲得",GREEN,"無料モニター10社確保\n口コミ・成功事例の作成\nプロダクト改善"),
    ("Phase 2  M4〜M12","有料化と拡大",BLUE,"直販（電話・訪問）開始\n社労士パートナー経由\n45社・MRR 90万円目標"),
    ("Phase 3  Y2〜Y3","スケール",PURPLE,"派遣協会セミナー展示\nYouTube・SNS集客\n300社・MRR 600万円目標"),
]
for i,(ph,ttl,c,desc) in enumerate(phases):
    x=0.35+i*4.3
    rect(s,x,1.45,4.1,3.2,fill=LGRAY); rect(s,x,1.45,4.1,0.85,fill=c)
    txt(s,ph,x,1.48,4.1,0.38,size=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s,ttl,x,1.87,4.1,0.38,size=15,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    for j,line in enumerate(desc.split("\n")):
        txt(s,f"  {line}",x+0.25,2.42+j*0.58,3.75,0.48,size=12,color=DARK)
    if i<2: txt(s,"→",4.5+i*4.3,2.65,0.3,0.5,size=24,bold=True,color=GRAY,align=PP_ALIGN.CENTER)
txt(s,"主要チャネルと優先度",0.35,4.85,5,0.4,size=13,bold=True,color=DARK)
channels=[
    ("中小派遣会社への直販（電話・訪問）","最優先",GREEN),
    ("社労士・中小企業診断士経由の紹介",  "高",    BLUE),
    ("派遣協会（JASSA）セミナー・展示",   "中",    ORANGE),
    ("YouTube・SNS（派遣会社経営者向け）","中長期", GRAY),
]
for i,(ch,pri,c) in enumerate(channels):
    y=5.38+i*0.52
    rect(s,0.35,y,8.8,0.46,fill=LGRAY)
    txt(s,ch,0.5,y+0.07,7.5,0.35,size=11,color=DARK)
    rect(s,9.2,y,1.5,0.46,fill=c)
    txt(s,pri,9.2,y+0.07,1.5,0.35,size=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"初期獲得の鍵：無料90日トライアル → コンディション改善の成功事例を2〜3社作り横展開",0.35,7.12,12.5,0.3,size=11,bold=True,color=DGREEN)

# ---- S9 収支計画 Year1 月次 ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"07  収支計画（Year 1）","月次損益シミュレーション  M7に単月黒字転換")
pnum(s,9)
months=["M1","M2","M3","M4","M5","M6","M7","M8","M9","M10","M11","M12"]
clients=[0,0,0,5,8,12,17,22,27,33,38,45]
staff=  [0,0,0,100,100,110,110,120,120,130,130,140]
rev=    [0,0,0,100,160,264,374,528,648,858,988,1260]
cost=   [300,300,300,300,300,320,320,350,350,380,380,400]
profit= [r-c for r,c in zip(rev,cost)]
sx=0.35
rect(s,sx,1.35,12.6,0.42,fill=NAVY)
txt(s,"項目",sx+0.05,1.38,0.95,0.36,size=10,bold=True,color=WHITE)
for j,h in enumerate(months):
    x=sx+1.05+j*0.965
    txt(s,h,x,1.38,0.9,0.36,size=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
row_labels=["顧客数（社）","スタッフ数（平均）","売上（千円）","コスト（千円）","損益（千円）"]
row_data=[clients,staff,rev,cost,profit]
row_fc=[DARK,DARK,GREEN,ORANGE,None]
for i,(label,data,fc) in enumerate(zip(row_labels,row_data,row_fc)):
    y=1.82+i*0.56
    bg=LGRAY if i%2==0 else WHITE
    rect(s,sx,y,12.6,0.52,fill=bg)
    txt(s,label,sx+0.05,y+0.1,0.95,0.35,size=10,bold=True,color=DARK)
    for j,v in enumerate(data):
        x=sx+1.05+j*0.965
        if label=="損益（千円）":
            c=GREEN if v>=0 else RED
            display=f"+{v}" if v>=0 else f"▲{abs(v)}"
        else:
            c=fc if fc else DARK
            display=str(v) if v>0 else "-"
        txt(s,display,x,y+0.1,0.9,0.35,size=10,bold=(label=="損益（千円）"),color=c,align=PP_ALIGN.CENTER)
base_y=6.85; max_h=1.1
for j,(r,p) in enumerate(zip(rev,profit)):
    bx=0.35+j*0.97
    rh=(r/1400)*max_h if r>0 else 0.02
    rect(s,bx,base_y-rh,0.38,rh,fill=GREEN)
    ch2=(cost[j]/1400)*max_h
    rect(s,bx+0.4,base_y-ch2,0.38,ch2,fill=ORANGE)
rect(s,0.35,base_y,11.65,0.02,fill=DARK)
rect(s,0.35,7.1,12.6,0.32,fill=LGREEN)
txt(s,"Year 1 通期：売上累計 518万円  コスト累計 390万円  営業利益 128万円  （単月黒字転換：Month 7）",0.5,7.12,12.3,0.28,size=11,bold=True,color=DGREEN)

# ---- S10 収支計画 3ヵ年 ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"08  収支計画（3ヵ年）","中期成長シナリオ  Year 3で年商7,200万円・利益率70%")
pnum(s,10)
years=[
    ("Year 1  2026年", 45, 100, 90,  518,  390,  128,  25, GREEN),
    ("Year 2  2027年",150, 130, 390,3120,  900, 2220,  71, BLUE),
    ("Year 3  2028年",350, 150,1050,7140, 2160, 4980,  70, PURPLE),
]
for i,(yr,nc,ns,mrr,rev2,cost2,profit2,margin,c) in enumerate(years):
    x=0.35+i*4.3
    rect(s,x,1.4,4.1,5.6,fill=LGRAY); rect(s,x,1.4,4.1,0.58,fill=c)
    txt(s,yr,x,1.43,4.1,0.52,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    kpis2=[
        ("顧客数",    f"{nc} 社"),
        ("平均スタッフ",f"{ns} 名/社"),
        ("月末MRR",   f"{mrr}万円"),
        ("年間売上",  f"{rev2}万円"),
        ("年間コスト",f"{cost2}万円"),
        ("営業利益",  f"{profit2}万円"),
        ("利益率",    f"{margin}%"),
    ]
    for j,(lbl,val) in enumerate(kpis2):
        y2=2.1+j*0.68
        rect(s,x+0.15,y2,3.8,0.6,fill=WHITE)
        txt(s,lbl,x+0.25,y2+0.1,1.7,0.4,size=10,color=GRAY)
        fc=GREEN if "利益" in lbl else DARK
        bld="利益" in lbl or "MRR" in lbl
        txt(s,val,x+2.05,y2+0.08,1.75,0.44,size=13,bold=bld,color=fc,align=PP_ALIGN.RIGHT)
for i in range(2):
    x=4.5+i*4.3
    txt(s,"→",x,3.85,0.25,0.5,size=20,bold=True,color=GRAY,align=PP_ALIGN.CENTER)
rect(s,0.35,7.05,12.6,0.35,fill=LBLUE)
txt(s,"前提：月額200円/人、解約率5%/月、新規獲得：Y1=45社・Y2=120社・Y3=220社、コストはインフラ・マーケ・人件費",0.5,7.08,12.3,0.3,size=10,color=BLUE)

# ---- S11 KPI ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"09  KPI・マイルストーン","達成すべき指標と時間軸")
pnum(s,11)
kpis3=[
    ("MRR（月次経常収益）","90万円","390万円","1,050万円",GREEN),
    ("顧客数","45社","150社","350社",BLUE),
    ("登録スタッフ数","4,500名","19,500名","52,500名",PURPLE),
    ("解約率（月次）","< 8%","< 6%","< 5%",ORANGE),
    ("CAC（顧客獲得コスト）","15,000円","12,000円","10,000円",NAVY),
    ("LTV（顧客生涯価値）","30万円","40万円","48万円",DGREEN),
]
rect(s,0.35,1.35,12.6,0.48,fill=NAVY)
for j,(lbl,w) in enumerate(zip(["KPI指標","Year 1 目標","Year 2 目標","Year 3 目標"],[3.8,2.7,2.7,2.7])):
    x=0.35+(0 if j==0 else 3.85+(j-1)*2.75)
    txt(s,lbl,x+0.1,1.38,w,0.38,size=12,bold=True,color=WHITE)
for i,(label,v1,v2,v3,c) in enumerate(kpis3):
    y=1.88+i*0.75
    bg=WHITE if i%2==0 else LGRAY
    rect(s,0.35,y,12.6,0.7,fill=bg); rect(s,0.35,y,0.08,0.7,fill=c)
    txt(s,label,0.5,y+0.15,3.25,0.42,size=12,bold=True,color=DARK)
    for j,val in enumerate([v1,v2,v3]):
        x=4.0+j*2.75
        txt(s,val,x,y+0.12,2.6,0.45,size=14,bold=True,color=c,align=PP_ALIGN.CENTER)
rect(s,0.35,6.72,12.6,0.65,fill=LGREEN)
txt(s,"重要マイルストーン",0.5,6.75,3,0.35,size=12,bold=True,color=DGREEN)
milestones=[("M3","モニター10社完了"),("M7","単月黒字転換"),("M12","累計黒字"),("Y2-Q1","MRR 200万突破"),("Y3","年商7,000万円")]
for i,(m,t) in enumerate(milestones):
    x=0.5+i*2.45
    txt(s,f"[{m}]",x,6.88,1.2,0.28,size=9,bold=True,color=DGREEN,align=PP_ALIGN.CENTER)
    txt(s,t,x-0.1,7.12,1.4,0.28,size=9,color=DARK,align=PP_ALIGN.CENTER)

# ---- S12 調達計画 ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"10  調達計画","資金調達額・使途・バリュエーション")
pnum(s,12)
txt(s,"調達希望額",0.4,1.38,3,0.4,size=13,bold=True,color=DARK)
rect(s,0.35,1.85,5.5,1.5,fill=NAVY)
txt(s,"3,000万円",0.35,1.95,5.5,0.9,size=38,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"シードラウンド",0.35,2.85,5.5,0.38,size=14,color=GREEN,align=PP_ALIGN.CENTER)
txt(s,"資金使途内訳",6.2,1.38,6.5,0.4,size=13,bold=True,color=DARK)
uses=[("開発・エンジニア採用",40,BLUE),("マーケティング・営業",30,GREEN),("インフラ・セキュリティ",15,ORANGE),("法務・会計・雑費",10,PURPLE),("予備費",5,GRAY)]
cum=0
for label,pct,c in uses:
    bw=pct/100*6.3
    rect(s,6.2+cum*6.3/100,1.85,bw,0.55,fill=c)
    cum+=pct
for i,(label,pct,c) in enumerate(uses):
    y=2.55+i*0.62
    rect(s,6.2,y,0.38,0.38,fill=c)
    txt(s,f"{label}  {pct}%  （{int(3000*pct/100)}万円）",6.65,y+0.02,6.0,0.38,size=12,color=DARK)
rect(s,0.35,5.55,5.5,1.65,fill=LGRAY)
txt(s,"想定バリュエーション",0.5,5.6,4.5,0.38,size=12,bold=True,color=DARK)
txt(s,"プレマネー：1億5,000万円\nポストマネー：1億8,000万円\n希薄化率：約16.7%",0.5,6.05,4.8,1.1,size=13,color=DARK)
rect(s,6.2,5.55,6.75,1.65,fill=LGRAY)
txt(s,"調達後のRunway",6.35,5.6,6.4,0.38,size=12,bold=True,color=DARK)
txt(s,"月間バーンレート：約250万円\nRunway：約12ヶ月\n（MRR黒字化でRunwayは延長）",6.35,6.05,6.4,1.1,size=13,color=DARK)

# ---- S13 リスク ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"11  リスクと対策","主要リスクの認識と軽減策")
pnum(s,13)
risks=[
    (RED,   "高","競合の参入・模倣",
     "大手がLINE連携を追加した場合\n既存顧客のロックイン強化\nコンディションデータ蓄積による解約防止\n派遣特化の深化（シフト管理等）で差別化維持"),
    (ORANGE,"中","顧客解約率の上昇",
     "コスト削減で解約が増えるリスク\nオンボーディング支援の充実\n月次レポート送付で価値を実感させる\n解約理由の分析と製品改善サイクル"),
    (BLUE,  "中","LINE仕様変更・LIFFポリシー変更",
     "LINE社の方針変更でLIFFが制限されるリスク\nPWAへのフォールバック機能を開発\nLINE以外の認証手段（電話番号SMS）も準備"),
    (PURPLE,"低","個人情報・セキュリティインシデント",
     "打刻データ・GPS情報の漏洩リスク\nSupabase RLSによる行レベルアクセス制御\nサーバーサイドのみでのパスワード検証\n定期的なセキュリティ監査の実施"),
]
for i,(c,level,ttl,desc) in enumerate(risks):
    col,row=i%2,i//2
    x,y=0.35+col*6.35,1.45+row*2.65
    rect(s,x,y,6.1,2.45,fill=LGRAY); rect(s,x,y,6.1,0.45,fill=c)
    txt(s,f"リスク {i+1}：{ttl}",x+0.12,y+0.07,4.5,0.34,size=13,bold=True,color=WHITE)
    rect(s,x+4.7,y+0.07,1.2,0.3,fill=WHITE)
    txt(s,f"影響度：{level}",x+4.72,y+0.08,1.15,0.28,size=10,bold=True,color=c,align=PP_ALIGN.CENTER)
    txt(s,desc,x+0.2,y+0.58,5.7,1.75,size=10,color=DARK)

# ---- S14 チーム ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"12  チーム・体制","実行力のある経営チーム")
pnum(s,14)
members=[
    (NAVY,  "代表取締役\nオーナー",      "派遣業界・HR領域の事業開発経験\n顧客ネットワークと市場理解が強み\n営業・マーケティング統括"),
    (GREEN, "AI CEO  John",             "戦略立案・プロダクト開発・技術統括\n全機能の設計・実装・セキュリティ対応\nデータ分析・事業計画策定"),
    (BLUE,  "技術顧問\n（採用予定）",    "Next.js / Supabase フルスタック\nセキュリティ設計・スケーリング\nシード調達後に採用予定"),
    (PURPLE,"営業・CS\n（採用予定）",    "派遣会社向けの直販営業\nオンボーディング・カスタマーサクセス\nシード調達後に採用予定"),
]
for i,(c,role,desc) in enumerate(members):
    x=0.35+i*3.15
    rect(s,x,1.45,2.95,4.8,fill=LGRAY); rect(s,x,1.45,2.95,0.08,fill=c)
    rect(s,x+0.9,1.6,1.15,1.15,fill=c)
    txt(s,role,x,2.9,2.95,0.65,size=12,bold=True,color=DARK,align=PP_ALIGN.CENTER)
    txt(s,desc,x+0.15,3.65,2.65,2.45,size=10,color=GRAY)
rect(s,0.35,6.5,12.6,0.82,fill=LGREEN)
txt(s,"アドバイザー（予定）",0.5,6.55,3,0.35,size=12,bold=True,color=DGREEN)
advisors=["派遣業界 経営コンサルタント","社会保険労務士","SaaS スタートアップ投資家","LINE ミニアプリ技術専門家"]
for i,a in enumerate(advisors):
    txt(s,f"  {a}",0.5+i*3.1,6.92,2.9,0.32,size=10,color=DARK)

# ---- S15 なぜ今か ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=WHITE)
hdr(s,"Why Now?  なぜ今か","この市場機会は2〜3年以内に消える")
pnum(s,15)
reasons=[
    (GREEN, "電子帳簿保存法の厳格化（2024年）",
     "紙タイムカードが事実上使えなくなった。中小派遣会社は今すぐデジタル化ツールを必要としている。この波に最速で乗れる。"),
    (BLUE,  "LINEの市場支配が続く今が最大チャンス",
     "LINE利用率86%は今後も維持されるが、競合が気づいてLINE LIFF対応を始めれば差別化が薄れる。先行者利益を取れる時間は限られている。"),
    (RED,   "製造・物流DXが急加速中",
     "コロナ後の物流DXブームで現場作業員へのITツール導入ハードルが大きく下がった。リモートや多現場管理のニーズが急増している。"),
    (ORANGE,"AI開発コストがゼロに近い今",
     "AI CEO（Claude）の活用により、通常数千万円かかる開発コストを極限まで圧縮した。この構造的コスト優位は競合が容易に再現できない。"),
]
for i,(c,ttl,desc) in enumerate(reasons):
    col,row=i%2,i//2
    x,y=0.35+col*6.35,1.45+row*2.65
    rect(s,x,y,6.1,2.45,fill=LGRAY); rect(s,x,y,6.1,0.08,fill=c)
    txt(s,ttl,x+0.2,y+0.18,5.7,0.5,size=14,bold=True,color=DARK)
    txt(s,desc,x+0.2,y+0.82,5.75,1.5,size=11,color=GRAY)

# ---- S16 まとめ CTA ----
s=prs.slides.add_slide(blank)
rect(s,0,0,13.33,7.5,fill=NAVY)
rect(s,0,0,13.33,0.06,fill=GREEN); rect(s,0,7.44,13.33,0.06,fill=GREEN)
for i in range(6): rect(s,10.5+i*0.45,0,0.35,7.5,fill=RGBColor(0x1F,0x2A,0x3E))
txt(s,"SUMMARY  なぜラクラク勤怠に投資するか",0.6,0.5,12,0.45,size=14,bold=True,color=GREEN,italic=True)
txt(s,"3つの理由",0.6,1.0,4,0.6,size=24,bold=True,color=WHITE)
summaries=[
    (GREEN, "01","唯一無二のポジション",
     "コンディション報告 x 離職アラート x LINE完結を同時に持つ競合は存在しない"),
    (BLUE,  "02","巨大な未開拓市場",
     "13,000社の中小派遣会社が狙えるSAM 450億円。既存ツールは大企業向けで中小は空白地帯"),
    (ORANGE,"03","実証済みプロダクトと高い実行力",
     "Vercel本番稼働中・LINEでの実打刻確認済み。AI CEOによる圧倒的低コスト開発体制"),
]
for i,(c,num,ttl,desc) in enumerate(summaries):
    y=1.75+i*1.55
    rect(s,0.6,y,9.5,1.38,fill=RGBColor(0x1A,0x2A,0x3E)); rect(s,0.6,y,0.08,1.38,fill=c)
    txt(s,num,0.75,y+0.1,0.8,0.6,size=28,bold=True,color=c)
    txt(s,ttl,1.6,y+0.1,8.3,0.5,size=16,bold=True,color=WHITE)
    txt(s,desc,1.6,y+0.68,8.3,0.65,size=12,color=RGBColor(0xCC,0xDD,0xFF))
rect(s,0.6,6.4,9.5,0.85,fill=GREEN)
txt(s,"ご投資・パートナーシップのご相談はこちらまで",0.8,6.45,7,0.4,size=14,bold=True,color=WHITE)
txt(s,"contact@rakuraku-kintai.jp（準備中）",0.8,6.82,7,0.35,size=12,color=LGREEN)
pnum(s,16)

out = r"C:\Users\PC_User\Desktop\AI動画\rakuraku-kintai\らくらく勤怠\sales\ラクラク勤怠_投資家向けプレゼン資料.pptx"
prs.save(out)
print("saved:", out)