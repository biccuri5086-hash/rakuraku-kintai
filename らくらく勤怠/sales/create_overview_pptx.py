from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

LINE_GREEN = RGBColor(0x06, 0xC7, 0x55)
DARK_GREEN = RGBColor(0x04, 0x9A, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x7B, 0x8D)
LIGHT_GREEN = RGBColor(0xE8, 0xF8, 0xEE)
BLUE = RGBColor(0x22, 0x8B, 0xE6)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xE5, 0x3E, 0x3E)
LIGHT_GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)

W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txb


def multiline(slide, lines, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, spacing=1.2):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(size * (spacing - 1) * 0.5)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def header(slide, title, subtitle=None, num=None, total=10, color=LINE_GREEN):
    box(slide, 0, 0, 13.33, 1.2, color)
    txt(slide, title, 0.5, 0.18, 11.5, 0.6, 26, WHITE, bold=True)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.72, 11.5, 0.4, 13, RGBColor(0xE8, 0xFF, 0xF0))
    if num:
        txt(slide, f"{num} / {total}", 12.5, 0.85, 0.7, 0.3, 11, WHITE, align=PP_ALIGN.RIGHT)


TOTAL = 9

# =============================================
# 1. タイトル
# =============================================
s = blank_slide()
bg(s, LINE_GREEN)
box(s, 0, 0, 7.7, 7.5, WHITE)
box(s, 7.7, 0, 0.15, 7.5, DARK_GREEN)
txt(s, "ラクラク勤怠", 0.6, 1.6, 6.8, 1.2, 46, LINE_GREEN, bold=True)
txt(s, "サービス全体像", 0.6, 2.5, 6.8, 0.9, 30, DARK, bold=True)
txt(s, "〜 仕組みと機能のご紹介 〜", 0.6, 3.3, 6.8, 0.6, 16, GRAY)
multiline(s, [
    "対象：派遣会社向け勤怠管理SaaS",
    "入口：LINE（スタッフ）／Web管理画面（管理者・運営）",
], 0.6, 4.3, 6.8, 1.2, 13, GRAY)
txt(s, "非エンジニア向け説明資料", 8.1, 1.2, 4.7, 0.6, 13, WHITE, align=PP_ALIGN.CENTER)
txt(s, "2026", 8.1, 6.8, 4.7, 0.4, 12, RGBColor(0xCC, 0xFF, 0xDD), align=PP_ALIGN.CENTER)

# =============================================
# 2. このサービスは何か
# =============================================
s = blank_slide()
bg(s, LIGHT_GRAY_BG)
header(s, "1. このサービスは何か", "派遣会社向けの勤怠管理SaaS。3者が関わる", 2, TOTAL)

roles = [
    ("派遣スタッフ", "LINEで打刻・\nコンディション報告", "LINE (LIFF)", LINE_GREEN),
    ("派遣会社の管理者", "出退勤確認・シフト/契約/\n給与/コンプラ管理", "Web管理画面\n/admin", BLUE),
    ("運営（自社）", "契約している\n派遣会社（テナント）の管理", "Web管理画面\n/superadmin", DARK),
]
for i, (title, desc, entry, color) in enumerate(roles):
    bx = 0.4 + i * 4.25
    box(s, bx, 1.6, 3.95, 4.6, WHITE)
    box(s, bx, 1.6, 3.95, 0.65, color)
    txt(s, title, bx + 0.15, 1.68, 3.7, 0.5, 16, WHITE, bold=True)
    for j, line in enumerate(desc.split("\n")):
        txt(s, line, bx + 0.2, 2.5 + j * 0.4, 3.6, 0.4, 13, DARK)
    txt(s, "入口", bx + 0.2, 3.6, 3.6, 0.35, 11, GRAY, bold=True)
    for j, line in enumerate(entry.split("\n")):
        txt(s, line, bx + 0.2, 4.0 + j * 0.4, 3.6, 0.4, 13, color, bold=True)

txt(s, "一つのシステムを、契約している派遣会社ごとにデータを分離して提供している",
    0.5, 6.5, 12.3, 0.6, 15, DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)

# =============================================
# 3. 全体の仕組み（データの流れ）
# =============================================
s = blank_slide()
bg(s, WHITE)
header(s, "2. 全体の仕組み", "情報がどこをどう通るか（データの流れ）", 3, TOTAL, color=DARK)

nodes = [
    ("📱", "スタッフの\nスマホ（LINE）"),
    ("🌐", "Vercel\n（アプリ本体）"),
    ("🔑", "LINE\n（本人確認）"),
    ("🗄️", "Supabase\n（データベース）"),
    ("💻", "管理者の\nパソコン"),
]
nx = 0.3
nw = 2.3
for i, (emoji, label) in enumerate(nodes):
    box(s, nx, 2.2, nw, 2.0, LIGHT_GREEN)
    txt(s, emoji, nx, 2.35, nw, 0.8, 32, DARK, align=PP_ALIGN.CENTER)
    for j, line in enumerate(label.split("\n")):
        txt(s, line, nx + 0.05, 3.2 + j * 0.35, nw - 0.1, 0.35, 12, DARK, bold=True, align=PP_ALIGN.CENTER)
    if i < len(nodes) - 1:
        txt(s, "→", nx + nw, 2.9, 0.5, 0.6, 22, LINE_GREEN, bold=True, align=PP_ALIGN.CENTER)
    nx += nw + 0.5

multiline(s, [
    "・ 通信はすべてHTTPSで暗号化",
    "・ スタッフのスマホ本体にはデータを残さない（サーバー側に集約）",
    "・ 派遣会社ごとにデータが分離される設計（他社のデータは見えない）",
], 0.6, 4.8, 12.0, 1.6, 15, DARK, spacing=1.6)

# =============================================
# 4. データモデルの考え方
# =============================================
s = blank_slide()
bg(s, LIGHT_GRAY_BG)
header(s, "3. データモデルの考え方", "単発の仕事も、長期の仕事も、同じ仕組みで扱う", 4, TOTAL)

chain = [
    ("companies", "派遣元\n（契約している派遣会社）", LINE_GREEN),
    ("clients", "派遣先企業", BLUE),
    ("assignments", "契約・アサイン\n（単発 or 中長期）", ORANGE),
    ("shifts", "シフト\n（予定・確定・完了・欠勤）", RGBColor(0x8E, 0x5B, 0xE0)),
    ("attendance", "実際の打刻データ", DARK_GREEN),
]
cy = 1.7
for i, (name, desc, color) in enumerate(chain):
    box(s, 1.0, cy, 5.2, 0.85, WHITE)
    box(s, 1.0, cy, 0.12, 0.85, color)
    txt(s, name, 1.3, cy + 0.08, 2.2, 0.4, 15, color, bold=True)
    for j, line in enumerate(desc.split("\n")):
        txt(s, line, 3.6, cy + 0.06 + j * 0.32, 2.5, 0.3, 11, GRAY)
    if i < len(chain) - 1:
        txt(s, "↓", 3.4, cy + 0.85, 0.5, 0.3, 16, GRAY, align=PP_ALIGN.CENTER)
    cy += 1.05

txt(s, "「今日1日だけの仕事」も「数ヶ月契約の仕事」も、同じデータ構造の上で\nシフト管理・勤怠集計・給与計算までつながる",
    6.7, 2.3, 6.0, 1.8, 15, DARK, bold=False)
box(s, 6.7, 4.3, 6.0, 1.6, LIGHT_GREEN)
txt(s, "💡 単発＝中長期の特殊ケース（1回きり）として\n同じ仕組みに乗せる設計",
    6.9, 4.5, 5.6, 1.2, 14, DARK_GREEN, bold=True)

# =============================================
# 5. スタッフ機能
# =============================================
s = blank_slide()
bg(s, RGBColor(0xF0, 0xFD, 0xF4))
header(s, "4. スタッフ機能", "LINEで完結。アプリのダウンロード不要", 5, TOTAL)

staff_feats = [
    ("登録", "初回のみ電話番号を入力\n（アプリDL不要）"),
    ("打刻", "「出勤する」「退勤する」を\n1タップ"),
    ("コンディション報告", "5段階の絵文字で\n「今日の調子」を毎日記録（5秒）"),
    ("GPS記録", "打刻時に位置情報を自動取得\n（不正打刻の確認に活用）"),
]
for i, (title, desc) in enumerate(staff_feats):
    bx = 0.4 + i * 3.2
    box(s, bx, 1.6, 2.95, 4.2, WHITE)
    box(s, bx, 1.6, 2.95, 0.08, LINE_GREEN)
    txt(s, title, bx + 0.15, 1.85, 2.7, 0.6, 15, DARK, bold=True)
    for j, line in enumerate(desc.split("\n")):
        txt(s, line, bx + 0.15, 2.6 + j * 0.4, 2.7, 0.4, 12, GRAY)

txt(s, "スマホ操作が苦手なスタッフでも、当日から迷わず使える",
    0.5, 6.3, 12.3, 0.6, 16, LINE_GREEN, bold=True, align=PP_ALIGN.CENTER)

# =============================================
# 6. 管理者機能①：日々の現場管理
# =============================================
s = blank_slide()
bg(s, WHITE)
header(s, "5. 管理者機能①：日々の現場管理", None, 6, TOTAL, color=BLUE)

admin1 = [
    ("ダッシュボード", "全スタッフの出退勤状況と\nコンディションを一覧表示"),
    ("要フォローアラート", "コンディションが悪い\nスタッフを自動で可視化"),
    ("シフト管理", "予定・確定・完了・欠勤を\nステータス管理"),
    ("派遣先・契約管理", "派遣先企業（clients）、\n契約/アサインを登録・管理"),
]
for i, (title, desc) in enumerate(admin1):
    bx = 0.4 + i * 3.2
    box(s, bx, 1.6, 2.95, 4.5, LIGHT_GREEN)
    box(s, bx, 1.6, 2.95, 0.55, BLUE)
    txt(s, title, bx + 0.15, 1.68, 2.7, 0.5, 14, WHITE, bold=True)
    for j, line in enumerate(desc.split("\n")):
        txt(s, line, bx + 0.15, 2.4 + j * 0.4, 2.7, 0.4, 12, DARK)

# =============================================
# 7. 管理者機能②：バックオフィス業務
# =============================================
s = blank_slide()
bg(s, LIGHT_GRAY_BG)
header(s, "6. 管理者機能②：バックオフィス業務", None, 7, TOTAL, color=BLUE)

rows = [
    ("給与集計", "勤怠データから給与を自動集計（設定画面あり）"),
    ("有給管理", "スタッフの有給取得状況を管理"),
    ("派遣法コンプライアンス", "抵触日を基準に「余裕あり／要注意／超過」を自動判定"),
    ("請求（billing）", "派遣先向けの請求管理"),
    ("監査ログ", "ログイン・操作履歴を記録（誰が・いつ・何をしたか）"),
    ("派遣先向けレポート", "client-report として出力"),
]
ry = 1.55
for i, (title, desc) in enumerate(rows):
    rc = WHITE if i % 2 == 0 else RGBColor(0xEE, 0xF2, 0xF6)
    box(s, 0.4, ry, 12.5, 0.78, rc)
    txt(s, title, 0.6, ry + 0.13, 3.4, 0.5, 14, DARK, bold=True)
    txt(s, desc, 4.1, ry + 0.15, 8.6, 0.5, 13, GRAY)
    ry += 0.82

# =============================================
# 8. 運営（自社）機能
# =============================================
s = blank_slide()
bg(s, DARK)
header(s, "7. 運営（自社）機能", None, 8, TOTAL, color=RGBColor(0x33, 0x33, 0x33))

items = [
    ("🏢", "契約している派遣会社（テナント）の追加・管理を /superadmin から実施"),
    ("🔐", "管理者・運営ともに2段階認証（2FA）に対応"),
    ("🔑", "パスワードは12文字以上・英字/数字/記号必須のポリシーで統一管理"),
]
iy = 2.0
for emoji, text in items:
    box(s, 0.6, iy, 12.1, 1.2, RGBColor(0x2A, 0x2A, 0x2A))
    txt(s, emoji, 0.9, iy + 0.3, 0.8, 0.6, 28, WHITE)
    txt(s, text, 1.8, iy + 0.35, 10.6, 0.6, 15, WHITE)
    iy += 1.45

# =============================================
# 9. まとめ
# =============================================
s = blank_slide()
bg(s, LINE_GREEN)
txt(s, "まとめ", 0.5, 0.9, 12.3, 0.8, 30, WHITE, bold=True, align=PP_ALIGN.CENTER)
box(s, 1.5, 2.1, 10.3, 3.6, WHITE)
multiline(s, [
    "「LINEで完結する打刻・コンディション報告」を入口に、",
    "",
    "シフト・契約・給与・派遣法コンプライアンスまでを",
    "1つのデータ基盤でつなげているのが",
    "",
    "ラクラク勤怠の全体像。",
], 2.0, 2.6, 9.3, 2.8, 18, DARK, align=PP_ALIGN.CENTER, spacing=1.5)

txt(s, "技術構成：Next.js（App Router）+ TypeScript ／ Supabase（PostgreSQL）／ LINE LIFF ／ Vercel",
    0.5, 6.1, 12.3, 0.5, 12, RGBColor(0xE8, 0xFF, 0xF0), align=PP_ALIGN.CENTER)
txt(s, "9 / 9", 12.5, 7.05, 0.7, 0.3, 11, RGBColor(0xE8, 0xFF, 0xF0), align=PP_ALIGN.RIGHT)

out = "/home/user/rakuraku-kintai/らくらく勤怠/sales/16_アプリ全体像_仕組みと機能.pptx"
prs.save(out)
print(f"保存完了: {out}")
