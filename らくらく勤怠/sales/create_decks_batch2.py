# -*- coding: utf-8 -*-
import sys
sys.path.insert(0, r"c:\Users\PC_User\Desktop\AI動画\rakuraku-kintai\らくらく勤怠\sales")
from create_decks_batch1 import (
    new_prs, rect, txt, cover, hdr, card_row, table_slide, step_row, cta,
    GREEN, DGREEN, LGREEN, WHITE, DARK, GRAY, LGRAY, RED, LRED, ORANGE, LORG,
    BLUE, LBLUE, NAVY, PURPLE, GOLD, SAVE_DIR,
)
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# =====================================================
# 06 スタッフ向け導入マニュアル
# =====================================================
def deck_06():
    TOTAL = 6
    prs = new_prs()
    cover(prs, "スタッフ向け 導入マニュアル", "LINEで簡単 1タップ打刻", "STAFF GUIDE")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "ご用意いただくもの", "とってもシンプル！", 2, TOTAL)
    rect(s, 0.35, 1.7, 12.6, 5.0, fill=LGREEN)
    txt(s, "📱 スマートフォン1台だけ", 0.5, 2.0, 12.5, 0.7, size=28, bold=True, color=DGREEN, align=PP_ALIGN.CENTER)
    txt(s, "LINEアプリがインストールされていればOK", 0.5, 2.8, 12.5, 0.5, size=14, color=DARK, align=PP_ALIGN.CENTER)
    rect(s, 3.5, 4.0, 6.3, 2.2, fill=WHITE)
    txt(s, "✓ 専用アプリのダウンロード 不要", 3.7, 4.15, 6.0, 0.4, size=13, color=DARK)
    txt(s, "✓ パソコン 不要", 3.7, 4.6, 6.0, 0.4, size=13, color=DARK)
    txt(s, "✓ 特別な設定 不要", 3.7, 5.05, 6.0, 0.4, size=13, color=DARK)
    txt(s, "✓ 担当者からURLを受け取るだけ", 3.7, 5.5, 6.0, 0.4, size=13, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "STEP 1〜3：はじめての登録", "初回のみ・5分で完了", 3, TOTAL)
    step_row(s, [
        (GREEN, "STEP 1", "URLを開く",      "担当者から送られた\nLINEのURLをタップ"),
        (BLUE,  "STEP 2", "電話番号を登録", "携帯電話番号を入力\nお名前も入力"),
        (ORANGE,"STEP 3", "出勤画面が表示", "登録完了！\nすぐに使えます"),
    ], y=1.7, h=3.5)
    rect(s, 0.35, 5.5, 12.6, 1.5, fill=LBLUE)
    txt(s, "💡 ワンポイント", 0.5, 5.6, 12.5, 0.4, size=13, bold=True, color=BLUE)
    txt(s, "電話番号は本人確認のために1回だけ入力します。次回からは画面が出ません。\nお名前は給与計算・契約書類で使われる正式な氏名を入力してください。",
        0.5, 6.0, 12.3, 1.0, size=12, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "STEP 4：毎日の出勤", "ボタン1回で完了", 4, TOTAL)
    rect(s, 0.35, 1.7, 6.0, 5.0, fill=GREEN)
    txt(s, "出勤する", 0.35, 3.5, 6.0, 1.0, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "← この大きなボタンをタップするだけ", 6.6, 3.7, 6.5, 0.5, size=16, color=DARK)
    txt(s, "・打刻時刻が自動で記録されます", 6.6, 4.3, 6.5, 0.4, size=12, color=GRAY)
    txt(s, "・位置情報も自動で記録（任意）", 6.6, 4.7, 6.5, 0.4, size=12, color=GRAY)
    txt(s, "・担当者に出勤が通知されます", 6.6, 5.1, 6.5, 0.4, size=12, color=GRAY)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "STEP 5：コンディション報告（5秒）", "今日の調子を絵文字で", 5, TOTAL)
    txt(s, "出勤後に「今日の調子」を絵文字で選んでください", 0.4, 1.5, 12.5, 0.5, size=14, color=DARK, align=PP_ALIGN.CENTER)
    conditions = [
        (GREEN, "😄", "絶好調！"),
        (GREEN, "😊", "良い感じ"),
        (ORANGE,"😐", "普通"),
        (RED,   "😔", "少し疲れ"),
        (RED,   "😢", "しんどい"),
    ]
    width = 2.4
    for i, (c, emoji, label) in enumerate(conditions):
        x = 0.4 + i * (width + 0.1)
        rect(s, x, 2.3, width, 3.0, fill=LGRAY)
        rect(s, x, 2.3, width, 0.08, fill=c)
        txt(s, emoji, x, 2.5, width, 1.5, size=60, color=DARK, align=PP_ALIGN.CENTER)
        txt(s, label, x, 4.5, width, 0.5, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    rect(s, 0.35, 5.6, 12.6, 1.5, fill=LBLUE)
    txt(s, "コメントは任意です。書きたくなければ空欄でOK。", 0.5, 5.7, 12.5, 0.4, size=13, color=BLUE)
    txt(s, "「しんどい」を選ぶと担当者がすぐに気づいてくれます。無理しないでくださいね。", 0.5, 6.2, 12.5, 0.4, size=12, color=DARK)

    cta(prs, "あとは退勤ボタンを押すだけ", "毎日カンタン！スタッフのあなたを応援しています", 6, TOTAL)
    prs.save(SAVE_DIR + r"\06_スタッフ向け導入マニュアル.pptx")
    print("06 saved")

# =====================================================
# 07 管理者向けマニュアル
# =====================================================
def deck_07():
    TOTAL = 7
    prs = new_prs()
    cover(prs, "管理者向けマニュアル", "ラクラク勤怠 ダッシュボードの使い方", "ADMIN GUIDE")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "管理者ダッシュボードへのアクセス", "PCでもスマホでもOK", 2, TOTAL)
    rect(s, 0.35, 1.7, 12.6, 1.2, fill=NAVY)
    txt(s, "https://rakuraku-kintai-frb6.vercel.app/admin", 0.5, 1.95, 12.3, 0.7, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 0.35, 3.2, 12.6, 3.5, fill=LGRAY)
    txt(s, "アクセス手順", 0.5, 3.3, 12.5, 0.4, size=14, bold=True, color=DARK)
    txt(s,
        "1. 上記URLをブラウザで開く（Chrome / Safari / Edge いずれもOK）\n"
        "2. パスワードを入力（別途ご連絡）\n"
        "3. 2要素認証コード（6桁）を入力\n"
        "4. ダッシュボードが表示される\n\n"
        "※ パスワードは絶対に他人に教えないでください\n"
        "※ 2FAコードは Google Authenticator 等のアプリで表示されます",
        0.5, 3.8, 12.3, 2.8, size=13, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "サマリーカード（画面上部）", "一目で全体状況がわかる", 3, TOTAL)
    card_row(s, [
        (GREEN, "👥 登録スタッフ数", "ラクラク勤怠を\n使っているスタッフの\n総人数"),
        (BLUE,  "🕐 出勤済み",      "本日 出勤打刻した\nスタッフの人数\n(リアルタイム更新)"),
        (RED,   "⚠️ 要フォロー",   "コンディションが\n「疲れ」or「しんどい」\nのスタッフ数"),
    ], y=1.7, h=4.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "スタッフ一覧テーブル", "全員の打刻・コンディションを一覧表示", 4, TOTAL)
    headers = ["名前", "出勤", "退勤", "勤務時間", "コンディション"]
    rows = [
        ["山田 太郎 (LINE:山田/090-1234)", "08:00", "17:00", "9h 0m", "😊 良い感じ"],
        ["佐藤 花子 (LINE:hanako/090-5678)", "09:00", "--:--", "(勤務中)", "😄 絶好調"],
        ["⚠ 鈴木 一郎 (LINE:いちろう/090-9876)", "10:00", "16:00", "6h 0m", "😢 しんどい"],
    ]
    table_slide(s, headers, rows, y=1.7, row_h=0.6)
    rect(s, 0.35, 4.8, 12.6, 2.2, fill=LRED)
    txt(s, "⚠️ 要フォローの見方", 0.5, 4.9, 12.5, 0.4, size=14, bold=True, color=RED)
    txt(s,
        "コンディションが「少し疲れ」または「しんどい」のスタッフは\n"
        "行が赤背景で表示され、名前の前に「⚠」マークが付きます。\n\n"
        "→ 早めに声をかけることで離職防止に繋がります",
        0.5, 5.35, 12.3, 1.6, size=12, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "監査ログ（操作履歴）", "誰が・いつ・何をしたか追跡可能", 5, TOTAL)
    txt(s, "右上の 🛡 アイコンをクリックで監査ログ画面へ", 0.4, 1.5, 12.5, 0.5, size=13, color=DARK)
    headers = ["記録される操作", "用途"]
    rows = [
        ["✅ ログイン成功",         "正規アクセスの記録"],
        ["❌ ログイン失敗",         "不正アクセス試行の検知"],
        ["❌ 2FAコード失敗",        "セキュリティインシデント検知"],
        ["🚫 試行回数上限到達",    "ブルートフォース攻撃検知"],
        ["🚪 ログアウト",          "セッション終了の記録"],
        ["👀 ダッシュボード閲覧",   "誰がデータを見たか追跡"],
    ]
    table_slide(s, headers, rows, y=2.0, row_h=0.5)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "よくある管理者の質問", "", 6, TOTAL)
    qa = [
        ("Q. スタッフの打刻を修正したい",
         "A. 管理画面の打刻修正機能（スタンダードプラン以上）から修正可能。修正履歴は監査ログに残ります。"),
        ("Q. スタッフを削除したい",
         "A. 管理画面 → スタッフ管理 → 該当スタッフを選択 → 削除。データは30日間保持されます。"),
        ("Q. 月のデータをExportしたい",
         "A. 管理画面 → CSV出力ボタン。給与計算ソフトに取り込み可能な形式で出力されます。"),
        ("Q. 緊急時に2FA解除できる？",
         "A. はい。Vercel管理画面から ADMIN_TOTP_SECRET を削除すれば2FA無効化されます。"),
    ]
    for i, (q, a) in enumerate(qa):
        y = 1.55 + i * 1.35
        rect(s, 0.35, y, 12.6, 1.2, fill=LGRAY); rect(s, 0.35, y, 0.08, 1.2, fill=NAVY)
        txt(s, q, 0.55, y+0.15, 12.3, 0.45, size=13, bold=True, color=DARK)
        txt(s, a, 0.55, y+0.6, 12.3, 0.55, size=11, color=GRAY)

    cta(prs, "サポート窓口", "biccuri5086@gmail.com / 080-9895-7770", 7, TOTAL)
    prs.save(SAVE_DIR + r"\07_管理者向けマニュアル.pptx")
    print("07 saved")

# =====================================================
# 08 営業メールテンプレート
# =====================================================
def deck_08():
    TOTAL = 7
    prs = new_prs()
    cover(prs, "営業メールテンプレート集", "5パターン使い分け", "INTERNAL")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "5つのパターン", "状況に応じて選ぶ", 2, TOTAL)
    rows = [
        ["①", "コールドメール", "初回・無関係先",  "本文 250〜400字"],
        ["②", "紹介経由メール", "社労士・知人経由", "本文 200〜300字"],
        ["③", "LP問合せ返信",  "見込み客自発接触", "本文 200〜300字"],
        ["④", "フォローアップ", "1週間返信なし",   "本文 150〜250字"],
        ["⑤", "デモ後フォロー", "商談直後即送信",  "本文 300〜400字"],
    ]
    table_slide(s, ["#", "種類", "使うタイミング", "想定長さ"], rows, y=1.7, row_h=0.65)

    # 各テンプレートのコア部分
    templates = [
        ("①コールドメール", "件名: 派遣スタッフの離職率を下げる無料ツール（LINE完結）",
         "本文骨子: 自己紹介→課題提示→3つの強み→無料デモ提案→連絡先\nゴール: 15分デモアポイント獲得"),
        ("②紹介経由メール", "件名: 【〇〇様よりご紹介】派遣会社向け勤怠管理サービス",
         "本文骨子: 紹介者の名前→紹介の経緯→特徴3点→面談打診\nゴール: 信頼バイアスを使ったアポ取得"),
        ("③LP問合せ返信", "件名: お問い合わせありがとうございます【ラクラク勤怠】",
         "本文骨子: 御礼→デモ可能日時3提案→事前ヒアリング3項目\nゴール: 24時間以内の返信で機会損失防止"),
        ("④フォローアップ", "件名: Re: 派遣スタッフの離職率を下げる無料ツール（再送）",
         "本文骨子: 配慮の一言→簡単に試せる選択肢提示→今後の連絡判断委ねる\nゴール: 押しすぎず最後の機会創出"),
        ("⑤デモ後フォロー", "件名: 本日のデモのお礼【ラクラク勤怠 / 30日無料トライアル案内】",
         "本文骨子: 御礼→デモ振返り→無料トライアル開始手順→無償サポート\nゴール: 即日トライアル開始"),
    ]
    for tpl in templates:
        s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
        hdr(s, tpl[0], "概要・件名・本文骨子")
        rect(s, 0.35, 1.55, 12.6, 1.2, fill=LBLUE)
        txt(s, "件名", 0.5, 1.6, 2.5, 0.4, size=12, bold=True, color=BLUE)
        txt(s, tpl[1], 0.5, 2.0, 12.3, 0.7, size=13, color=DARK)
        rect(s, 0.35, 2.95, 12.6, 3.8, fill=LGRAY)
        txt(s, "本文骨子", 0.5, 3.05, 12.5, 0.4, size=12, bold=True, color=DARK)
        txt(s, tpl[2], 0.5, 3.55, 12.3, 3.0, size=13, color=DARK)
        txt(s, "完全な本文は markdown ファイル「08_営業メールテンプレート.md」参照", 0.4, 7.0, 12.5, 0.3, size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    cta(prs, "テンプレを下書きに保存しておく", "Gmailの下書き機能を使えば、相手企業名置換のみで即送信")
    prs.save(SAVE_DIR + r"\08_営業メールテンプレート.pptx")
    print("08 saved")

# =====================================================
# 09 電話営業スクリプト
# =====================================================
def deck_09():
    TOTAL = 7
    prs = new_prs()
    cover(prs, "電話営業スクリプト", "受付突破からアポ取得まで", "INTERNAL")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "通話の目的", "1回の電話で目指すゴール", 2, TOTAL)
    rect(s, 0.35, 1.7, 12.6, 2.2, fill=LGREEN)
    txt(s, "ゴール：15分のオンラインデモ アポイント獲得", 0.5, 2.0, 12.5, 0.8, size=24, bold=True, color=DGREEN, align=PP_ALIGN.CENTER)
    txt(s, "（電話で売り切ろうとしない。デモまで持っていけば十分）", 0.5, 2.85, 12.5, 0.6, size=14, color=DARK, align=PP_ALIGN.CENTER)
    rect(s, 0.35, 4.2, 6.0, 2.8, fill=LRED)
    txt(s, "❌ 電話でしてはいけないこと", 0.5, 4.3, 5.8, 0.4, size=14, bold=True, color=RED)
    txt(s, "・電話で契約まで決めようとする\n・営業/セールスと名乗る\n・強引なクロージング\n・断られても食い下がる",
        0.5, 4.75, 5.7, 2.0, size=12, color=DARK)
    rect(s, 6.95, 4.2, 6.0, 2.8, fill=LGREEN)
    txt(s, "✅ 電話で必ずすること", 7.1, 4.3, 5.8, 0.4, size=14, bold=True, color=DGREEN)
    txt(s, "・「ご案内/ご紹介」と表現\n・30秒だけと時間明示\n・LINEで・・を強調\n・断られても次の連絡許可取り",
        7.1, 4.75, 5.7, 2.0, size=12, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "受付突破トーク", "30秒で担当者へ取り次いでもらう", 3, TOTAL)
    rect(s, 0.35, 1.7, 12.6, 5.3, fill=LGRAY)
    txt(s, "受付の方への第一声", 0.5, 1.8, 12.5, 0.4, size=14, bold=True, color=NAVY)
    txt(s,
        "「お忙しいところ恐れ入ります。ラクラク勤怠の小原と申します。\n"
        "派遣スタッフ向けの勤怠管理サービスの件で、人事ご担当の方を\n"
        "お願いできますでしょうか？」",
        0.5, 2.3, 12.3, 1.5, size=13, color=DARK)
    txt(s, "「どのようなご用件で？」と聞かれたら", 0.5, 3.95, 12.5, 0.4, size=14, bold=True, color=NAVY)
    txt(s,
        "「派遣スタッフのLINEで打刻できる新しいサービスのご案内です。\n"
        "30秒で概要をお伝えしたいので、ご担当者様をお願いできますでしょうか？」",
        0.5, 4.45, 12.3, 1.2, size=13, color=DARK)
    txt(s, "ポイント", 0.5, 5.85, 12.5, 0.4, size=13, bold=True, color=ORANGE)
    txt(s, "・「営業」と言わない → 「ご案内」「ご紹介」\n・「30秒だけ」と時間を明示 → 心理的負担を下げる",
        0.5, 6.3, 12.3, 0.8, size=11, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "担当者への第一声（15秒）", "40秒で要点を伝えてアポ提案", 4, TOTAL)
    rect(s, 0.35, 1.7, 12.6, 5.3, fill=LBLUE)
    txt(s, "オープニング・トーク（暗記推奨）", 0.5, 1.8, 12.5, 0.4, size=14, bold=True, color=BLUE)
    txt(s,
        "「お忙しいところ恐れ入ります、ラクラク勤怠の小原と申します。\n"
        "30秒だけお時間頂戴できますでしょうか？\n\n"
        "派遣スタッフがLINEを開いてボタンを1回押すだけで打刻できる\n"
        "SaaSのご案内をしておりまして、\n\n"
        "紙タイムカードの集計負担や、スタッフの突然離職に\n"
        "お困りの派遣会社様に多くご導入いただいております。\n\n"
        "価格は150〜200円/人/月で、初期費用ゼロ、\n"
        "30日間全機能無料でお試しいただけます。\n\n"
        "15分のオンラインデモのお時間を頂戴することは可能でしょうか？」",
        0.5, 2.3, 12.3, 4.6, size=12, color=DARK)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "断り文句への切り返し", "よく出る5パターン", 5, TOTAL)
    rows = [
        ["忙しい/興味ない", "現状ヒアリングに転換 → 提案へ繋ぐ"],
        ["資料送って",       "送る → 必ず次回連絡の約束を取る"],
        ["LINEは不安",       "セキュリティ実績を提示 → 別資料案内"],
        ["他社と契約してる", "現サービス比較資料を送付提案"],
        ["決裁者に確認",     "決裁者向け資料を送付 → 同席提案"],
    ]
    table_slide(s, ["断り文句", "切り返しトーク"], rows, y=1.7, row_h=0.7)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "KPI と コール戦略", "数字で押さえる", 6, TOTAL)
    headers = ["指標", "目標値/日", "備考"]
    rows = [
        ["架電数",         "30件",     "効率重視で時間内に多く"],
        ["担当者接続率",   "30% (9件)", "受付突破できた数"],
        ["アポ取得率",     "15% (4-5件)", "担当者と話せた中で"],
        ["デモ参加率",     "60% (3件)",  "アポ取得した中で参加"],
        ["成約率",         "30% (1件)",  "デモ参加した中で成約"],
    ]
    table_slide(s, headers, rows, y=1.7, row_h=0.55)
    rect(s, 0.35, 5.5, 12.6, 1.5, fill=LGREEN)
    txt(s, "→ 30件架電で1社獲得が現実的な目安", 0.5, 5.6, 12.5, 0.5, size=16, bold=True, color=DGREEN, align=PP_ALIGN.CENTER)
    txt(s, "ベストタイム：火・水・木の10:00-11:00 / 14:00-15:00", 0.5, 6.2, 12.5, 0.5, size=13, color=DARK, align=PP_ALIGN.CENTER)

    cta(prs, "今日10件架電してみる", "完璧を求めず、まず数をこなす", 7, TOTAL)
    prs.save(SAVE_DIR + r"\09_電話営業スクリプト.pptx")
    print("09 saved")

# =====================================================
# 10 ターゲット企業リストアップガイド
# =====================================================
def deck_10():
    TOTAL = 6
    prs = new_prs()
    cover(prs, "ターゲット企業リストアップガイド", "派遣会社100社を無料でリスト化", "INTERNAL")

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "理想のターゲット像", "1社あたりLTV ¥48万", 2, TOTAL)
    headers = ["項目", "詳細"]
    rows = [
        ["業種",         "製造・物流・食品加工・軽作業系の派遣"],
        ["規模",         "社員5〜30名、登録スタッフ50〜300名"],
        ["エリア",       "愛知・岐阜・三重(製造)／関東・関西(物流)"],
        ["決裁者",       "社長または営業部長（即決傾向あり）"],
        ["IT成熟度",     "低〜中（紙タイムカード or Excel管理）"],
    ]
    table_slide(s, headers, rows, y=1.7, row_h=0.7)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "リストアップ方法 5選", "全て無料・公開情報", 3, TOTAL)
    rows = [
        ["⭐⭐⭐⭐⭐", "厚労省 人材サービス総合サイト", "派遣業許可事業者の全公開DB"],
        ["⭐⭐⭐⭐",   "Google検索 + マップ",             "「○○市 派遣会社」で発見"],
        ["⭐⭐⭐⭐",   "LinkedIn",                         "決裁者に直接アプローチ"],
        ["⭐⭐⭐",     "商工会議所 会員名簿",              "地域密着・中小中心"],
        ["⭐⭐⭐",     "派遣業界団体 (JASSA等)",          "業界誌・セミナー活用"],
    ]
    table_slide(s, ["評価", "情報源", "特徴"], rows, y=1.7, row_h=0.7)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "厚労省サイトの活用法（最速）", "全国13,000社の派遣業者リストが見られる", 4, TOTAL)
    rect(s, 0.35, 1.5, 12.6, 1.0, fill=NAVY)
    txt(s, "https://jinzai-sougou.mhlw.go.jp/", 0.5, 1.7, 12.3, 0.6, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    step_row(s, [
        (GREEN, "STEP 1", "サイトを開く", "「事業所を検索」をクリック"),
        (BLUE,  "STEP 2", "都道府県を選択", "例: 愛知県"),
        (ORANGE,"STEP 3", "条件を絞る", "派遣事業のみチェック\n業種で絞り込み"),
        (PURPLE,"STEP 4", "情報を抽出", "会社名・住所・電話\nExcelに転記"),
    ], y=2.7, h=3.0)
    rect(s, 0.35, 6.0, 12.6, 1.0, fill=LGREEN)
    txt(s, "1時間で50〜100社のリスト構築可能", 0.5, 6.15, 12.5, 0.7, size=16, bold=True, color=DGREEN, align=PP_ALIGN.CENTER)

    s = prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,13.33,7.5,fill=WHITE)
    hdr(s, "リスト管理スプレッドシート", "推奨カラム構成", 5, TOTAL)
    headers = ["カラム", "用途", "例"]
    rows = [
        ["会社名",            "メイン情報", "○○派遣株式会社"],
        ["所在地",            "地域分析",   "愛知県名古屋市"],
        ["電話/メール",       "アプローチ手段", "052-xxx-xxxx"],
        ["推定スタッフ数",     "プラン提示", "50-100名"],
        ["初回コンタクト日",   "進捗管理",   "2026/06/01"],
        ["状態",              "ステージ",   "アポ取得 / トライアル中 / 成約"],
        ["次回アクション",     "失念防止",   "6/10 デモ実施"],
    ]
    table_slide(s, headers, rows, y=1.7, row_h=0.55)

    cta(prs, "まず厚労省サイトで100社リスト化", "もしくはクラウドワークスに5,000円で外注（次資料参照）", 6, TOTAL)
    prs.save(SAVE_DIR + r"\10_ターゲット企業リストアップガイド.pptx")
    print("10 saved")

deck_06()
deck_07()
deck_08()
deck_09()
deck_10()
print("Batch 2 complete")